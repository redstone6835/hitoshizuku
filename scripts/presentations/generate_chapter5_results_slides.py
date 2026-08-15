#!/usr/bin/env python3
"""生成答辩第五章“结果结论”的三页正文。"""

from __future__ import annotations

import argparse
from copy import deepcopy
from io import BytesIO
import os
from pathlib import Path
import tempfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from generate_engineering_structure_slide import (
    BG,
    BODY,
    BLUE,
    INK,
    LINE,
    MUTED,
    NAVY,
    PALE_BLUE,
    PALE_GRAY,
    PALE_PURPLE,
    PALE_TEAL,
    PURPLE,
    TEAL,
    WHITE,
    add_arrow_tip,
    add_line,
    add_rect,
    add_text,
    rgb,
    set_run_typefaces,
    style_run,
)


EMU_PER_INCH = 914400
MIN_FONT_PT = 14.0
CONTENT_X = 3.00
CONTENT_RIGHT = 12.60
CONTENT_W = CONTENT_RIGHT - CONTENT_X

CHAPTER5_TRANSITION = "第五章 · 结果结论"
CHAPTER6_TRANSITION = "第六章 · 总结展望"
CHAPTER5_TITLES = (
    "系统运行情况",
    "性能问题与优化取舍",
    "工程实践总结",
)


def inches(value: int) -> float:
    return value / EMU_PER_INCH


def clone_slide(prs: Presentation, source):
    destination = prs.slides.add_slide(source.slide_layout)
    for shape in list(destination.shapes):
        element = shape._element
        element.getparent().remove(element)
    destination.background.fill.solid()
    destination.background.fill.fore_color.rgb = rgb(BG)
    for shape in source.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture = destination.shapes.add_picture(
                BytesIO(shape.image.blob),
                shape.left,
                shape.top,
                shape.width,
                shape.height,
            )
            picture.crop_left = shape.crop_left
            picture.crop_right = shape.crop_right
            picture.crop_top = shape.crop_top
            picture.crop_bottom = shape.crop_bottom
            picture.rotation = shape.rotation
            continue
        destination.shapes._spTree.insert_element_before(
            deepcopy(shape._element), "p:extLst"
        )
    return destination


def slide_texts(slide) -> set[str]:
    return {
        shape.text.strip()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    }


def find_slide(prs: Presentation, exact_text: str):
    for slide in prs.slides:
        if exact_text in slide_texts(slide):
            return slide
    raise RuntimeError(f"没有找到幻灯片：{exact_text}")


def remove_slide(prs: Presentation, slide) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        if prs.part.related_part(slide_id.rId) is slide.part:
            prs.part.drop_rel(slide_id.rId)
            prs.slides._sldIdLst.remove(slide_id)
            return
    raise RuntimeError("没有找到待删除幻灯片关系")


def is_template_chrome(shape) -> bool:
    x, y = inches(shape.left), inches(shape.top)
    w, h = inches(shape.width), inches(shape.height)
    if x + w <= 2.20:
        return True
    if 2.55 <= x <= 2.75 and 0.52 <= y <= 0.75 and h <= 0.90:
        return True
    if 2.80 <= x <= 3.15 and 0.45 <= y <= 0.75 and w >= 8.0:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.LINE and 1.50 <= y <= 1.66:
        return True
    return False


def clean_body_template(slide) -> None:
    for shape in list(slide.shapes):
        if not is_template_chrome(shape):
            element = shape._element
            element.getparent().remove(element)


def set_title(slide, title: str) -> None:
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        x, y = inches(shape.left), inches(shape.top)
        if 2.80 <= x <= 3.15 and 0.45 <= y <= 0.75:
            candidates.append(shape)
    if not candidates:
        raise RuntimeError("没有找到正文标题文本框")
    frame = candidates[0].text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = title
    style_run(
        run,
        size=30,
        color=INK,
        bold=True,
        chinese_font="SimHei",
        latin_font="Times New Roman",
    )


def body(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 15.0,
    color=BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    if size < MIN_FONT_PT:
        raise ValueError(f"正文文字不得小于 {MIN_FONT_PT:g} pt：{text}")
    box = add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        chinese_font="SimSun",
        latin_font="Times New Roman",
        align=align,
        valign=valign,
        margin=0.04,
    )
    box.text_frame.word_wrap = True
    return box


def heading(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float = 0.34,
    *,
    size: float = 18.0,
    color=INK,
    align=PP_ALIGN.LEFT,
):
    box = add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=max(size, MIN_FONT_PT),
        color=color,
        bold=True,
        chinese_font="SimHei",
        latin_font="Times New Roman",
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.02,
    )
    box.text_frame.word_wrap = True
    return box


def section_label(slide, text: str, x: float, y: float, w: float, *, accent=BLUE):
    add_rect(slide, x, y + 0.04, 0.07, 0.32, accent)
    heading(slide, text, x + 0.18, y, w - 0.18, 0.40, size=18, color=NAVY)


def arrow(slide, x1: float, y1: float, x2: float, y2: float, *, color=BLUE):
    add_line(slide, x1, y1, x2, y2, color, 1.45)
    if abs(x2 - x1) >= abs(y2 - y1):
        add_arrow_tip(slide, x2, y2, "right" if x2 >= x1 else "left", color, 0.10)
    else:
        add_arrow_tip(slide, x2, y2, "down" if y2 >= y1 else "up", color, 0.10)


def compact_metric(
    slide,
    x: float,
    y: float,
    w: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
):
    add_rect(slide, x, y, w, 0.86, fill)
    add_rect(slide, x, y, 0.07, 0.86, accent)
    heading(slide, title, x + 0.18, y + 0.09, w - 0.28, 0.27, size=16, color=NAVY)
    body(
        slide,
        detail,
        x + 0.18,
        y + 0.38,
        w - 0.28,
        0.41,
        size=14,
        color=BODY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def path_node(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
):
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.07, h, accent)
    heading(slide, title, x + 0.16, y + 0.10, w - 0.24, 0.30, size=16, color=NAVY)
    body(
        slide,
        detail,
        x + 0.16,
        y + 0.43,
        w - 0.24,
        h - 0.51,
        size=14,
        color=BODY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def draw_01(slide) -> None:
    set_title(slide, CHAPTER5_TITLES[0])

    compact_metric(slide, 3.00, 1.82, 2.12, "运行架构", "LoongArch64 · RISC-V64")
    compact_metric(slide, 5.27, 1.82, 2.15, "用户接口", "Linux syscall · POSIX", fill=PALE_TEAL, accent=TEAL)
    compact_metric(slide, 7.57, 1.82, 2.08, "并发环境", "SMP · 进程 · 信号", fill=PALE_PURPLE, accent=PURPLE)
    compact_metric(slide, 9.80, 1.82, 2.80, "基础子系统", "VM · VFS · 文件系统 · 网络", fill=PALE_GRAY, accent=MUTED)

    section_label(slide, "实际运行路径", 3.00, 2.90, 5.98, accent=BLUE)
    path_node(slide, 3.00, 3.40, 1.32, 0.94, "设备发现", "DTB / PCI")
    path_node(slide, 4.55, 3.40, 1.32, 0.94, "PnP绑定", "match / probe", fill=PALE_TEAL, accent=TEAL)
    path_node(slide, 6.10, 3.40, 1.32, 0.94, "能力发布", "DeviceFunction", fill=PALE_PURPLE, accent=PURPLE)
    path_node(slide, 7.65, 3.40, 1.33, 0.94, "能力消费", "类型化调用 · 投影", fill=PALE_GRAY, accent=MUTED)
    arrow(slide, 4.34, 3.87, 4.51, 3.87, color=BLUE)
    arrow(slide, 5.89, 3.87, 6.06, 3.87, color=TEAL)
    arrow(slide, 7.44, 3.87, 7.61, 3.87, color=PURPLE)

    path_node(slide, 3.00, 4.78, 1.32, 0.94, "配置声明", "Modules.toml", fill=PALE_PURPLE, accent=PURPLE)
    path_node(slide, 4.55, 4.78, 1.32, 0.94, "ELM管理", "依赖 · 代际 · 退出", fill=PALE_PURPLE, accent=PURPLE)
    path_node(slide, 6.10, 4.78, 1.32, 0.94, "服务实例", "net.stack / virtio", fill=PALE_TEAL, accent=TEAL)
    path_node(slide, 7.65, 4.78, 1.33, 0.94, "常驻接口", "VFS / socket / fd")
    arrow(slide, 4.34, 5.25, 4.51, 5.25, color=PURPLE)
    arrow(slide, 5.89, 5.25, 6.06, 5.25, color=PURPLE)
    arrow(slide, 7.44, 5.25, 7.61, 5.25, color=TEAL)

    add_rect(slide, 3.00, 6.08, 5.98, 0.60, NAVY)
    body(
        slide,
        "用户程序  →  Linux ABI / 原生态接口  →  内核服务",
        3.20,
        6.18,
        5.58,
        0.38,
        size=15,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    section_label(slide, "net.stack重载行为", 9.34, 2.90, 3.26, accent=TEAL)
    states = (
        ("运行代际", "当前StackHandle", PALE_BLUE, BLUE),
        ("静默阶段", "拒绝新turn", PALE_PURPLE, PURPLE),
        ("旧fd", "ENETDOWN", PALE_GRAY, MUTED),
        ("等待队列", "ERR | HUP", PALE_TEAL, TEAL),
        ("新代际", "新的StackHandle", PALE_BLUE, BLUE),
    )
    y = 3.40
    for index, (title, detail, fill, accent) in enumerate(states):
        add_rect(slide, 9.34, y, 3.26, 0.55, fill)
        add_rect(slide, 9.34, y, 0.07, 0.55, accent)
        heading(slide, title, 9.52, y + 0.10, 0.98, 0.28, size=15, color=NAVY)
        body(slide, detail, 10.55, y + 0.07, 1.91, 0.36, size=14, color=INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if index != len(states) - 1:
            arrow(slide, 10.97, y + 0.57, 10.97, y + 0.71, color=accent)
        y += 0.67


def cost_row(
    slide,
    y: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
):
    add_rect(slide, 3.00, y, 3.66, 0.80, fill)
    add_rect(slide, 3.00, y, 0.07, 0.80, accent)
    heading(slide, title, 3.18, y + 0.09, 1.44, 0.28, size=16, color=NAVY)
    body(slide, detail, 4.57, y + 0.07, 1.95, 0.62, size=14, color=BODY, bold=True, valign=MSO_ANCHOR.MIDDLE)


def table_cell(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 14,
    color=INK,
    bold: bool = False,
    fill=WHITE,
    align=PP_ALIGN.LEFT,
    chinese_font="SimSun",
):
    add_rect(slide, x, y, w, h, fill, line=LINE, line_width=0.7)
    box = add_text(
        slide,
        text,
        x + 0.08,
        y + 0.03,
        w - 0.16,
        h - 0.06,
        size=size,
        color=color,
        bold=bold,
        chinese_font=chinese_font,
        latin_font="Times New Roman",
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.01,
    )
    box.text_frame.word_wrap = True


def draw_02(slide) -> None:
    set_title(slide, CHAPTER5_TITLES[1])
    section_label(slide, "BuildStorm负载中的内核工作", 3.00, 1.82, 3.66, accent=BLUE)
    rows = (
        ("用户内核边界", "syscall分发；trap现场保存与恢复", PALE_BLUE, BLUE),
        ("缺页处理", "匿名页分配；文件页读取与cache fill", PALE_TEAL, TEAL),
        ("地址空间", "ASID分配；satp切换；TLB维护", PALE_PURPLE, PURPLE),
        ("并发运行", "任务调度；跨核同步；shootdown", PALE_BLUE, BLUE),
        ("系统服务", "VFS路径；文件系统；设备I/O", PALE_GRAY, MUTED),
    )
    y = 2.34
    for index, (title, detail, fill, accent) in enumerate(rows):
        cost_row(slide, y, title, detail, fill=fill, accent=accent)
        if index != len(rows) - 1:
            arrow(slide, 4.83, y + 0.82, 4.83, y + 0.94, color=accent)
        y += 0.91

    section_label(slide, "优化方案处理记录", 6.96, 1.82, 5.64, accent=PURPLE)
    columns = (6.96, 8.73, 11.26, 12.60)
    headers = ("修改对象", "运行表现", "处理")
    for index, text in enumerate(headers):
        table_cell(
            slide,
            text,
            columns[index],
            2.34,
            columns[index + 1] - columns[index],
            0.55,
            size=15,
            color=WHITE,
            bold=True,
            fill=NAVY,
            align=PP_ALIGN.CENTER,
            chinese_font="SimHei",
        )
    decisions = (
        ("ASID与地址空间切换", "减少无效全局TLB处理", "保留", PALE_BLUE, BLUE),
        ("resident PTE快速路径", "长窗口稳定退化", "回退", PALE_PURPLE, PURPLE),
        ("自定义memset", "局部指令缩短，整体负载变慢", "回退", PALE_PURPLE, PURPLE),
        ("memcpy局部优化", "多轮运行结果不稳定", "暂不合入", PALE_GRAY, MUTED),
        ("陷阱扩展状态保存", "高频路径存在重复搬运", "继续处理", PALE_TEAL, TEAL),
    )
    y = 2.89
    for name, observed, decision, fill, accent in decisions:
        table_cell(slide, name, columns[0], y, columns[1] - columns[0], 0.76, size=14, bold=True, fill=fill)
        table_cell(slide, observed, columns[1], y, columns[2] - columns[1], 0.76, size=14, fill=WHITE)
        table_cell(slide, decision, columns[2], y, columns[3] - columns[2], 0.76, size=14, color=accent, bold=True, fill=fill, align=PP_ALIGN.CENTER, chinese_font="SimHei")
        y += 0.76


def draw_03(slide) -> None:
    set_title(slide, CHAPTER5_TITLES[2])

    columns = (3.00, 4.58, 8.10, 12.60)
    headers = ("设计原则", "工程实现", "实际表现")
    for index, text in enumerate(headers):
        table_cell(
            slide,
            text,
            columns[index],
            1.88,
            columns[index + 1] - columns[index],
            0.60,
            size=16,
            color=WHITE,
            bold=True,
            fill=NAVY,
            align=PP_ALIGN.CENTER,
            chinese_font="SimHei",
        )

    rows = (
        ("可拓展性", "抽象接口、实现组件与管理机制分离", "新增设备能力和内核模块进入统一发现、装载与管理路径", PALE_BLUE, BLUE),
        ("安全性", "能力边界、生命周期与确定的失效语义", "组件退出时停止新绑定，旧句柄、等待队列和新实例分别处理", PALE_TEAL, TEAL),
        ("可维护性", "general、hal、arch、drivers职责分层", "通用机制、架构实现与可配置组件在各自目录独立演进", PALE_PURPLE, PURPLE),
        ("可追踪性", "组件身份、调用路径、符号与动态指令映射", "异常与开销能够归属到模块、函数、指令及其执行上下文", PALE_BLUE, BLUE),
        ("可验证性", "固定负载、成对对照与多轮复测", "候选修改按运行结果保留、回退或继续观察", PALE_GRAY, MUTED),
    )
    y = 2.48
    for principle, implementation, behavior, fill, accent in rows:
        table_cell(
            slide,
            principle,
            columns[0],
            y,
            columns[1] - columns[0],
            0.86,
            size=16,
            color=accent,
            bold=True,
            fill=fill,
            align=PP_ALIGN.CENTER,
            chinese_font="SimHei",
        )
        table_cell(slide, implementation, columns[1], y, columns[2] - columns[1], 0.86, size=14, bold=True, fill=fill)
        table_cell(slide, behavior, columns[2], y, columns[3] - columns[2], 0.86, size=14, fill=WHITE)
        y += 0.86


DRAWERS = (draw_01, draw_02, draw_03)


def update_transition(slide) -> None:
    candidates = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and 3.10 <= inches(shape.top) <= 3.45
        and shape.text.strip()
    ]
    if not candidates:
        return
    frame = candidates[0].text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "系统运行、性能取舍与工程实践"
    style_run(
        run,
        size=18,
        color=BODY,
        bold=False,
        chinese_font="SimSun",
        latin_font="Times New Roman",
    )


def enforce_font_floor(slide) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                if run.font.size is None or run.font.size.pt < MIN_FONT_PT:
                    run.font.size = Pt(MIN_FONT_PT)
                if run.font.name is None:
                    set_run_typefaces(run, "SimSun", "Times New Roman")


def insert_chapter(prs: Presentation) -> list:
    transition = find_slide(prs, CHAPTER5_TRANSITION)
    next_transition = find_slide(prs, CHAPTER6_TRANSITION)
    transition_index = list(prs.slides).index(transition)
    next_index = list(prs.slides).index(next_transition)
    if next_index <= transition_index + 1:
        raise RuntimeError("第五章缺少可复用正文模板")

    existing = list(prs.slides)[transition_index + 1 : next_index]
    template = existing[0]
    update_transition(transition)

    added = []
    for drawer in DRAWERS:
        slide = clone_slide(prs, template)
        clean_body_template(slide)
        drawer(slide)
        enforce_font_floor(slide)
        added.append(slide)

    for slide in existing:
        remove_slide(prs, slide)

    transition_index = list(prs.slides).index(transition)
    slide_ids = prs.slides._sldIdLst
    added_ids = list(slide_ids)[-len(added) :]
    for slide_id in added_ids:
        slide_ids.remove(slide_id)
    for offset, slide_id in enumerate(added_ids, 1):
        slide_ids.insert(transition_index + offset, slide_id)
    return added


def atomic_save(prs: Presentation, output: Path) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        prefix=f".{output.stem}-",
        suffix=".pptx",
        dir=output.parent,
        delete=False,
    ) as temporary:
        temporary_path = Path(temporary.name)
    try:
        prs.save(temporary_path)
        os.replace(temporary_path, output)
        output.chmod(0o644)
    except Exception:
        temporary_path.unlink(missing_ok=True)
        raise


def build_topic(full_output: Path, topic_output: Path) -> None:
    prs = Presentation(full_output)
    keep_titles = set(CHAPTER5_TITLES)
    for slide in list(prs.slides):
        if not (slide_texts(slide) & keep_titles):
            remove_slide(prs, slide)
    if len(prs.slides) != len(CHAPTER5_TITLES):
        raise RuntimeError(f"第五章专题页数错误：{len(prs.slides)}")
    atomic_save(prs, topic_output)


def validate(path: Path, *, expected_slides: int | None = None) -> None:
    prs = Presentation(path)
    if expected_slides is not None and len(prs.slides) != expected_slides:
        raise RuntimeError(f"{path} 页数错误：{len(prs.slides)}")
    titles_found = []
    for slide_number, slide in enumerate(prs.slides, 1):
        texts = slide_texts(slide)
        titles_found.extend(title for title in CHAPTER5_TITLES if title in texts)
        if not (texts & set(CHAPTER5_TITLES)):
            continue
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip() or run.font.size is None:
                        continue
                    if run.font.size.pt < MIN_FONT_PT:
                        raise RuntimeError(
                            f"第 {slide_number} 页文字小于 {MIN_FONT_PT:g} pt：{run.text!r}"
                        )
    if sorted(titles_found) != sorted(CHAPTER5_TITLES):
        raise RuntimeError("第五章标题集合不完整或重复")


def main() -> int:
    parser = argparse.ArgumentParser()
    root = Path(__file__).resolve().parents[2]
    parser.add_argument(
        "--base",
        type=Path,
        default=root / "output/presentations/mygo-defense-full.pptx",
    )
    parser.add_argument(
        "--full-output",
        type=Path,
        default=root / "output/presentations/mygo-defense-full.pptx",
    )
    parser.add_argument(
        "--topic-output",
        type=Path,
        default=root / "output/presentations/mygo-defense-chapter5-results-3pages.pptx",
    )
    args = parser.parse_args()

    base = args.base.resolve()
    full_output = args.full_output.resolve()
    topic_output = args.topic_output.resolve()
    prs = Presentation(base)
    insert_chapter(prs)
    atomic_save(prs, full_output)
    validate(full_output)
    build_topic(full_output, topic_output)
    validate(topic_output, expected_slides=len(CHAPTER5_TITLES))
    print(full_output)
    print(topic_output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
