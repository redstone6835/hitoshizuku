#!/usr/bin/env python3
"""重做答辩稿中的 ELM 与网络栈专题，并接入正式全量稿。

精简专题共 6 页：ELM 4 页、网络栈 2 页。所有专题页可见文字均不小于 14 pt，
并删除正文模板右下角的章节数字。内容以 docs/chapters/chapter-12.typ 与
chapter-13.typ 以及对应实现为事实来源。
"""

from __future__ import annotations

import argparse
from copy import deepcopy
from io import BytesIO
import os
from pathlib import Path
import tempfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from generate_engineering_structure_slide import (
    BG,
    BODY,
    BLUE,
    INK,
    LINE,
    MUTED,
    NAVY,
    PALE_BLUE,
    PALE_GRAY,
    PALE_PURPLE,
    PALE_TEAL,
    PURPLE,
    TEAL,
    WHITE,
    add_arrow_tip,
    add_line,
    add_rect,
    add_text,
    rgb,
    set_run_typefaces,
    style_run,
)


EMU_PER_INCH = 914400
MIN_FONT_PT = 14.0
CONTENT_X = 3.00
CONTENT_RIGHT = 12.60
CONTENT_W = CONTENT_RIGHT - CONTENT_X
TOP_Y = 1.80
BOTTOM_Y = 6.78
TOPIC_START_TITLE = "设备抽象能力闭环"
TOPIC_END_TITLE = "第四章 · 调试方法"

OLD_TOPIC_TITLES = {
    "可拓展内核单元（ELM）的概念定位",
    "责任单元与管理自举",
    "身份、代际与关系拓扑",
    "能力契约闭环",
    "调用路径分层",
    "镜像投影与装载证明",
    "生命周期与代际替换",
    "策略、预算与资源所有权",
    "原生执行与故障边界",
    "可复核运行证据",
    "工程接入与完成边界",
    "网络子系统责任分层",
    "FlowShard 数据通路",
    "套接字兼容与控制边界",
    "网络生命周期与性能边界",
    # Titles retained by the immediately preceding full-deck version.
    "Busy 拒绝记录样例",
    "数据通路与受管调用边界",
}

TOPIC_TITLES = OLD_TOPIC_TITLES | {
    "可拓展内核单元（ELM）的概念定位",
    "常驻 Core、elm-mgr 与格式解析器",
    "Cell 身份与实现代际",
    "关系图中的四种责任",
    "能力发布：Contract、Port 与 Provider",
    "能力连接：Binding 与 Lease",
    "四条调用路径与各自成本",
    "从 EKI 文件到候选 EBI",
    "装载事务：提交前外部不可见",
    "生命周期状态与调用门禁",
    "Pause 与 Detach 的执行步骤",
    "Replace：同一 Cell 切换到下一代",
    "Policy、Budget、Lease 与长期资源",
    "原生调用门如何从故障返回",
    "一次 Busy 拒绝如何被复核",
    "开发入口、构建形态与当前接入范围",
    "网络对象分别由谁持有",
    "一次 sendmsg / recvmsg 经过哪些对象",
    "FlowShard 的并行模型与两条执行路径",
    "SocketFacade 如何保持 POSIX 语义稳定",
    "网络模块退出语义与性能开销",
    "装载事务：提交前外部不可见",
    "一次 Busy 拒绝如何被复核",
    "开发入口、构建形态与当前接入范围",
    "一次 sendmsg / recvmsg 经过哪些对象",
    "装载证明：来源与接口兼容",
    "装载提交：镜像准备与唯一公开点",
    "装载提交：镜像准备与一次性提交",
    "装载提交：两级门禁与失败回滚",
    "装载提交：一次性公开",
    "ELM 运行证据的五种视角",
    "ELM 的五类运行证据",
    "Busy 拒绝的逐项复核",
    "ElmModule 与 y / m / n 构建形态",
    "实现接入策略与边界",
    "当前接入范围与未完成边界",
    "完整接收路径：VirtIO DMA 到 recvmsg",
    "完整发送路径：sendmsg 到 VirtIO completion",
    "原生调用门的故障返回机制",
    "网络对象的所有权划分",
    "SocketFacade 的 POSIX 稳定边界",
    "调用边界、核验位置与运行成本",
    "Replace 事务：ElmId 不变，Generation 递增",
    "原生调用门的故障收束边界",
    "Busy 拒绝的证据闭环",
    "网络对象的所有权与退出责任",
    "FlowShard 单写者模型与双执行路径",
    "网络退役协议、用户语义与成本归因",
    "网络能力与责任结构",
    "网络热路径的受管调用边界",
}


def inches(value: int) -> float:
    return value / EMU_PER_INCH


def clone_slide(prs: Presentation, source):
    destination = prs.slides.add_slide(source.slide_layout)
    for shape in list(destination.shapes):
        element = shape._element
        element.getparent().remove(element)

    destination.background.fill.solid()
    destination.background.fill.fore_color.rgb = rgb(BG)
    for shape in source.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture = destination.shapes.add_picture(
                BytesIO(shape.image.blob),
                shape.left,
                shape.top,
                shape.width,
                shape.height,
            )
            picture.crop_left = shape.crop_left
            picture.crop_right = shape.crop_right
            picture.crop_top = shape.crop_top
            picture.crop_bottom = shape.crop_bottom
            picture.rotation = shape.rotation
            continue
        destination.shapes._spTree.insert_element_before(
            deepcopy(shape._element), "p:extLst"
        )
    return destination


def slide_text(slide) -> str:
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def find_slide(prs: Presentation, exact_text: str):
    for slide in prs.slides:
        if any(
            getattr(shape, "has_text_frame", False)
            and shape.text.strip() == exact_text
            for shape in slide.shapes
        ):
            return slide
    raise RuntimeError(f"没有找到幻灯片：{exact_text}")


def remove_slide(prs: Presentation, slide) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        if prs.part.related_part(slide_id.rId) is slide.part:
            prs.part.drop_rel(slide_id.rId)
            prs.slides._sldIdLst.remove(slide_id)
            return
    raise RuntimeError("没有找到待删除幻灯片关系")


def existing_topic_slides(prs: Presentation) -> list:
    slides = list(prs.slides)
    try:
        start = slides.index(find_slide(prs, TOPIC_START_TITLE))
        end = slides.index(find_slide(prs, TOPIC_END_TITLE))
    except (RuntimeError, ValueError):
        start = end = -1
    if 0 <= start < end:
        return slides[start + 1 : end]

    return [
        slide
        for slide in slides
        if {
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        }
        & TOPIC_TITLES
    ]


def remove_existing_topic(prs: Presentation) -> None:
    for slide in existing_topic_slides(prs):
        remove_slide(prs, slide)


def is_page_marker(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    text = shape.text.strip()
    x, y = inches(shape.left), inches(shape.top)
    return x >= 12.0 and y >= 6.9 and text in {"01", "02", "03", "04", "05", "06"}


def remove_page_markers(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if is_page_marker(shape):
                element = shape._element
                element.getparent().remove(element)


def is_template_chrome(shape) -> bool:
    x, y = inches(shape.left), inches(shape.top)
    w, h = inches(shape.width), inches(shape.height)
    if x + w <= 2.20:
        return True
    if 2.55 <= x <= 2.75 and 0.52 <= y <= 0.75 and h <= 0.90:
        return True
    if 2.80 <= x <= 3.15 and 0.45 <= y <= 0.75 and w >= 8.0:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.LINE and 1.50 <= y <= 1.66:
        return True
    return False


def clean_body_template(slide) -> None:
    for shape in list(slide.shapes):
        if is_page_marker(shape):
            element = shape._element
            element.getparent().remove(element)
            continue
        if not is_template_chrome(shape):
            element = shape._element
            element.getparent().remove(element)


def set_title(slide, title: str) -> None:
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        x, y = inches(shape.left), inches(shape.top)
        if 2.80 <= x <= 3.15 and 0.45 <= y <= 0.75:
            candidates.append(shape)
    if not candidates:
        raise RuntimeError("没有找到正文标题文本框")
    box = candidates[0]
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = title
    style_run(
        run,
        size=30,
        color=INK,
        bold=True,
        chinese_font="SimHei",
        latin_font="Times New Roman",
    )


def body(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 15.0,
    color=BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    if size < MIN_FONT_PT:
        raise ValueError(f"正文文字不得小于 {MIN_FONT_PT} pt: {size}")
    box = add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        chinese_font="SimSun",
        latin_font="Times New Roman",
        align=align,
        valign=valign,
        margin=0.05,
    )
    box.text_frame.word_wrap = True
    return box


def heading(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float = 0.34,
    *,
    size: float = 20.0,
    color=INK,
    align=PP_ALIGN.LEFT,
):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=max(size, MIN_FONT_PT),
        color=color,
        bold=True,
        chinese_font="SimHei",
        latin_font="Times New Roman",
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.03,
    )


def lead(slide, text: str) -> None:
    body(
        slide,
        text,
        CONTENT_X,
        1.78,
        CONTENT_W,
        0.62,
        size=16.0,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
    title_size: float = 19.0,
    detail_size: float = 15.0,
    center: bool = False,
):
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.09, h, accent)
    heading(
        slide,
        title,
        x + 0.23,
        y + 0.14,
        w - 0.42,
        0.34,
        size=title_size,
        color=NAVY if fill != NAVY else WHITE,
        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
    )
    body(
        slide,
        detail,
        x + 0.23,
        y + 0.61,
        w - 0.44,
        h - 0.72,
        size=detail_size,
        color=WHITE if fill == NAVY else BODY,
        bold=True,
        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
    )


def chip(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill=BLUE,
    color=WHITE,
    size: float = 16.0,
):
    add_rect(slide, x, y, w, 0.38, fill)
    add_text(
        slide,
        text,
        x,
        y,
        w,
        0.38,
        size=max(size, MIN_FONT_PT),
        color=color,
        bold=True,
        chinese_font="SimHei",
        latin_font="Times New Roman",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.01,
    )


def boundary(
    slide,
    text: str,
    *,
    fill=PALE_GRAY,
    accent=MUTED,
    label: str | None = None,
    label_w: float = 1.42,
) -> None:
    add_rect(slide, CONTENT_X, 6.14, CONTENT_W, 0.64, fill)
    add_rect(slide, CONTENT_X, 6.14, 0.09, 0.64, accent)
    text_x = CONTENT_X + 0.23
    text_w = CONTENT_W - 0.42
    align = PP_ALIGN.CENTER
    if label:
        heading(
            slide,
            label,
            CONTENT_X + 0.23,
            6.25,
            label_w,
            0.32,
            size=17.0,
            color=NAVY,
        )
        text_x = CONTENT_X + 0.34 + label_w
        text_w = CONTENT_W - label_w - 0.56
        align = PP_ALIGN.LEFT
    body(
        slide,
        text,
        text_x,
        6.18,
        text_w,
        0.56,
        size=14.0,
        color=INK,
        bold=True,
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
    )


def flow(slide, x1, y1, x2, y2, *, color=BLUE, direction="right", width=1.6):
    add_line(slide, x1, y1, x2, y2, color, width)
    add_arrow_tip(slide, x2, y2, direction, color, 0.11)


def step(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    detail: str = "",
    *,
    fill=PALE_BLUE,
    accent=BLUE,
):
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.08, h, accent)
    heading(
        slide,
        title,
        x + 0.15,
        y + 0.13,
        w - 0.25,
        0.30,
        size=17.0,
        align=PP_ALIGN.CENTER,
        color=WHITE if fill == NAVY else INK,
    )
    if detail:
        body(
            slide,
            detail,
            x + 0.13,
            y + 0.53,
            w - 0.22,
            h - 0.61,
            size=16.0,
            color=WHITE if fill == NAVY else BODY,
            align=PP_ALIGN.CENTER,
        )


def enforce_font_floor(slide) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                if run.font.size is None or run.font.size.pt < MIN_FONT_PT:
                    run.font.size = Pt(MIN_FONT_PT)
                if run.font.name is None:
                    set_run_typefaces(run, "SimSun", "Times New Roman")


def compact_row(
    slide,
    y: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
    h: float = 0.62,
    title_w: float = 2.10,
) -> None:
    add_rect(slide, CONTENT_X, y, CONTENT_W, h, fill)
    add_rect(slide, CONTENT_X, y, 0.09, h, accent)
    text_color = WHITE if fill == NAVY else INK
    heading(
        slide,
        title,
        CONTENT_X + 0.24,
        y + 0.14,
        title_w,
        0.31,
        size=18,
        color=text_color,
    )
    body(
        slide,
        detail,
        CONTENT_X + 0.28 + title_w,
        y + 0.10,
        CONTENT_W - title_w - 0.52,
        h - 0.16,
        size=15,
        color=text_color,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def numbered_row(
    slide,
    y: float,
    number: str,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
    h: float = 0.78,
    title_w: float = 2.08,
) -> None:
    """全宽流程行；给 14 pt 正文保留两行空间。"""
    add_rect(slide, CONTENT_X, y, CONTENT_W, h, fill)
    add_rect(slide, CONTENT_X, y, 0.09, h, accent)
    add_rect(slide, CONTENT_X + 0.20, y + 0.19, 0.40, 0.40, accent)
    body(
        slide,
        number,
        CONTENT_X + 0.20,
        y + 0.19,
        0.40,
        0.40,
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    heading(
        slide,
        title,
        CONTENT_X + 0.76,
        y + 0.13,
        title_w,
        h - 0.24,
        size=17,
        color=NAVY,
    )
    body(
        slide,
        detail,
        CONTENT_X + 0.88 + title_w,
        y + 0.09,
        CONTENT_W - title_w - 1.12,
        h - 0.16,
        size=15,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def definition_row(
    slide,
    y: float,
    term: str,
    definition: str,
    runtime_rule: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
    h: float = 0.82,
    term_w: float = 1.72,
    rule_w: float = 2.56,
) -> None:
    """术语、定义与运行规则三列并列，避免只展示名词。"""
    add_rect(slide, CONTENT_X, y, CONTENT_W, h, fill)
    add_rect(slide, CONTENT_X, y, 0.09, h, accent)
    heading(
        slide,
        term,
        CONTENT_X + 0.24,
        y + 0.13,
        term_w,
        h - 0.24,
        size=17,
        color=NAVY,
    )
    body(
        slide,
        definition,
        CONTENT_X + 0.34 + term_w,
        y + 0.10,
        CONTENT_W - term_w - rule_w - 0.66,
        h - 0.17,
        size=14,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_rect(
        slide,
        CONTENT_RIGHT - rule_w - 0.08,
        y + 0.11,
        0.03,
        h - 0.22,
        accent,
    )
    body(
        slide,
        runtime_rule,
        CONTENT_RIGHT - rule_w + 0.06,
        y + 0.10,
        rule_w - 0.18,
        h - 0.17,
        size=14,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def process_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    number: str,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
) -> None:
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.09, h, accent)
    add_rect(slide, x + 0.17, y + 0.14, 0.38, 0.38, accent)
    body(
        slide,
        number,
        x + 0.17,
        y + 0.14,
        0.38,
        0.38,
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    heading(slide, title, x + 0.65, y + 0.16, w - 0.82, 0.32, size=17, color=NAVY)
    body(
        slide,
        detail,
        x + 0.20,
        y + 0.61,
        w - 0.38,
        h - 0.72,
        size=15,
        color=BODY,
        bold=True,
        align=PP_ALIGN.LEFT,
    )


def arrow_caption(slide, x1: float, y: float, x2: float, label: str, *, color=BLUE) -> None:
    flow(slide, x1, y, x2, y, color=color)
    if label:
        body(
            slide,
            label,
            (x1 + x2) / 2 - 0.72,
            y - 0.43,
            1.44,
            0.30,
            size=16,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )


def draw_01(slide) -> None:
    set_title(slide, "可拓展内核单元（ELM）的概念定位")
    lead(slide, "传统模块已经具备装载、符号、依赖和签名；ELM 的改变，是把分散约束收敛为一个可计算的责任边界。")

    panel(
        slide,
        3.00,
        2.66,
        2.55,
        2.78,
        "传统模块已有能力",
        "装入镜像\n解析符号\n维护依赖\n签名与引用保护",
        fill=PALE_GRAY,
        accent=MUTED,
        detail_size=17.0,
        center=True,
    )
    panel(
        slide,
        5.88,
        2.66,
        2.82,
        2.78,
        "约束分散位置",
        "装载器\n子系统注册表\n引用计数与回调\n日志与运维工具",
        fill=PALE_PURPLE,
        accent=PURPLE,
        detail_size=17.0,
        center=True,
    )
    panel(
        slide,
        9.03,
        2.66,
        3.57,
        2.78,
        "ELM 统一责任面",
        "身份与代际\n关系与能力契约\n生命周期与资源\n故障与运行证据",
        fill=NAVY,
        accent=BLUE,
        detail_size=17.0,
        center=True,
    )
    flow(slide, 5.55, 4.05, 5.84, 4.05, color=MUTED)
    flow(slide, 8.70, 4.05, 8.99, 4.05, color=PURPLE)
    boundary(slide, "代码能够装入只是起点；能被发现、约束、复核和安全退役，才构成可治理的扩展。")


def draw_02(slide) -> None:
    set_title(slide, "责任单元与管理自举")
    lead(slide, "常驻 Core 持有硬约束；管理器和格式投影器进入同一拓扑，但不能反向改写 Core 不变量。")

    panel(
        slide,
        3.00,
        2.62,
        3.05,
        2.62,
        "ELM Core",
        "状态机与关系图\n端口、租约与预算\n装载证明与事务提交\n故障恢复与审计事实",
        fill=NAVY,
        accent=BLUE,
        detail_size=16.0,
        center=True,
    )
    panel(
        slide,
        6.48,
        2.62,
        2.70,
        2.62,
        "elm-mgr · ID 1",
        "根管理 Cell\n接收外部管理意图\n编排装载、策略、菜单和生命周期",
        fill=PALE_BLUE,
        accent=BLUE,
        center=True,
    )
    panel(
        slide,
        9.61,
        2.62,
        2.99,
        2.62,
        "eki · ID 2",
        "elm-mgr 的内建子单元\n把 EKI 投影为 EBI\n来源显示为 <builtin>",
        fill=PALE_PURPLE,
        accent=PURPLE,
        center=True,
    )
    flow(slide, 6.05, 3.93, 6.44, 3.93, color=BLUE)
    flow(slide, 9.18, 3.93, 9.57, 3.93, color=PURPLE)
    chip(slide, "启动内建端口", 3.00, 5.57, 1.58, fill=TEAL)
    body(
        slide,
        "core.log@1   core.event@1   mgr.menu.item@1   mgr.action.invoke@1",
        4.78,
        5.53,
        7.82,
        0.48,
        size=16.0,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    boundary(slide, "自举表示管理能力也有身份和状态，不表示最小可信 Core 可以被普通模块替换。")


def draw_03(slide) -> None:
    set_title(slide, "身份、代际与关系拓扑")
    lead(slide, "逻辑 Cell 身份保持稳定，具体实现由 Generation 区分；长期引用必须同时匹配二者。")

    heading(slide, "代际时间轴", 3.00, 2.54, 3.40, size=20)
    for x, label, fill, accent in (
        (3.00, "Gen 1", PALE_GRAY, MUTED),
        (4.68, "Gen 2", PALE_PURPLE, PURPLE),
        (6.36, "Gen 3", PALE_BLUE, BLUE),
    ):
        step(slide, x, 3.08, 1.36, 0.78, label, fill=fill, accent=accent)
    flow(slide, 4.36, 3.47, 4.64, 3.47, color=MUTED)
    flow(slide, 6.04, 3.47, 6.32, 3.47, color=PURPLE)
    body(
        slide,
        "ElmId(42) 始终不变\n旧代句柄在提交点后明确失效",
        3.00,
        4.20,
        4.72,
        1.02,
        size=16.0,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    heading(slide, "BindingGraph 的五类关系", 8.15, 2.54, 4.45, size=20)
    panel(slide, 8.15, 3.03, 2.04, 1.28, "父子", "管理归属\n预算委派", fill=PALE_BLUE, accent=BLUE, title_size=18)
    panel(slide, 10.47, 3.03, 2.13, 1.28, "依赖", "激活前置条件", fill=PALE_TEAL, accent=TEAL, title_size=18)
    panel(slide, 8.15, 4.38, 2.04, 1.28, "拓展点", "目标开放位置", fill=PALE_PURPLE, accent=PURPLE, title_size=18)
    panel(slide, 10.47, 4.38, 2.13, 1.28, "拓展项", "契约附着实现", fill=PALE_GRAY, accent=MUTED, title_size=18)
    chip(slide, "能力绑定：消费者 → Port → Provider", 8.15, 5.78, 4.45, fill=NAVY)
    boundary(slide, "关系提交检查端点、契约、重复和有向环；运行 ID 不是地址，也不保证跨启动稳定。")


def draw_04(slide) -> None:
    set_title(slide, "能力契约闭环")
    lead(slide, "一个能力必须从声明走到真实执行，并且在仍有活动引用时不能被提前撤销。")

    stages = [
        ("绑定请求", "消费者 + PortId", PALE_PURPLE, PURPLE, 1.70),
        ("Flow Contract", "name@version", PALE_BLUE, BLUE, 1.72),
        ("PortId", "连接点", PALE_TEAL, TEAL, 1.46),
        ("Binding + Lease", "连接与使用权", PALE_GRAY, MUTED, 1.88),
        ("Provider", "真实执行后端", NAVY, BLUE, 1.62),
    ]
    x = 3.00
    for index, (title, detail, fill, accent, width) in enumerate(stages):
        step(slide, x, 2.88, width, 1.24, title, detail, fill=fill, accent=accent)
        if index < len(stages) - 1:
            flow(slide, x + width, 3.50, x + width + 0.20, 3.50, color=accent)
        x += width + 0.20

    panel(
        slide,
        3.00,
        4.45,
        2.93,
        1.51,
        "Port",
        "方向、模式、访问策略\nowner Generation 与状态",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=18,
    )
    panel(
        slide,
        6.26,
        4.45,
        2.93,
        1.51,
        "Binding",
        "消费者—端口之间\n已提交、可审计的连接",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=18,
    )
    panel(
        slide,
        9.52,
        4.45,
        3.08,
        1.51,
        "Lease",
        "Active → Revoking → Revoked\n引用归零后完成撤销",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=18,
    )
    boundary(slide, "Provider / 受管调用使用 256 B 固定帧；管理 ABI 使用有界缓冲，硬上限 256 KiB。")


def draw_05(slide) -> None:
    set_title(slide, "调用路径分层")
    lead(slide, "ELM 不强迫所有调用经过通用 RPC；发现需求、ABI 稳定性和热路径敏感度共同决定路径。")

    panel(
        slide,
        3.00,
        2.60,
        2.90,
        1.58,
        "Provider",
        "动态发现、授权、审计\n同步 / 异步、取消与背压",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=19,
        center=True,
    )
    panel(
        slide,
        6.15,
        2.60,
        2.90,
        1.58,
        "Managed Import",
        "受管 ELM 接口\n每次调用核验双方代际",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=19,
        center=True,
    )
    panel(
        slide,
        9.30,
        2.60,
        3.30,
        1.58,
        "direct-pinned",
        "固定 export 与 Generation\n装载期精确校验 Rust ABI",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=19,
        center=True,
    )
    panel(
        slide,
        3.00,
        4.38,
        4.40,
        1.42,
        "kernel-symbol",
        "从登记目录解析常驻内核真实实现；校验名称、契约、版本、权限和完整 Rust ABI。",
        fill=NAVY,
        accent=BLUE,
        title_size=19,
    )
    panel(
        slide,
        7.72,
        4.38,
        4.88,
        1.42,
        "Runtime / Management 根 API",
        "普通 Cell 只取得运行时命名空间；Manager 取得管理命名空间，但每次 dispatch 仍重新鉴权。",
        fill=PALE_GRAY,
        accent=MUTED,
        title_size=19,
    )
    boundary(slide, "治理检查前置到装载、绑定和状态切换；稳定热路径不重复经过 elm-mgr。")


def draw_06(slide) -> None:
    set_title(slide, "镜像投影与装载证明")
    lead(slide, "文件格式解释与 Core 不变量分离：任何来源只能产出候选 EBI，最终执行资格仍由 Core 独立证明。")

    stages = [
        ("Upload", "分段写入", PALE_GRAY, MUTED, 1.28),
        ("Seal", "范围 + SHA-256", PALE_PURPLE, PURPLE, 1.42),
        ("EKI Source", "格式投影", PALE_BLUE, BLUE, 1.55),
        ("EBI", "装载协议对象", PALE_TEAL, TEAL, 1.42),
        ("Loader", "校验 + 重定位", PALE_BLUE, BLUE, 1.52),
        ("Commit", "初始化后公开", NAVY, BLUE, 1.45),
    ]
    x = 3.00
    for index, (title, detail, fill, accent, width) in enumerate(stages):
        step(slide, x, 2.72, width, 1.13, title, detail, fill=fill, accent=accent)
        if index < len(stages) - 1:
            flow(slide, x + width, 3.29, x + width + 0.17, 3.29, color=accent)
        x += width + 0.17

    panel(
        slide,
        3.00,
        4.12,
        4.46,
        1.88,
        "装载证明",
        "来源证明：签名或 BuildBound 绑定\n签名来源的 release epoch\n目标架构、EKI / EBI 与 Rust ABI\nimports / exports、W^X 与 I-cache",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=19,
    )
    panel(
        slide,
        7.78,
        4.12,
        4.82,
        1.88,
        "发布不变量",
        "create / initialize 成功前不对外可见；关系、Provider 与 imports 先暂存，唯一提交点之后才进入 Active。",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=19,
    )
    boundary(slide, "EBI 不是文件格式；soyo 仍是未来来源，签名也不能证明模块业务逻辑正确。")


def draw_07(slide) -> None:
    set_title(slide, "生命周期与代际替换")
    lead(slide, "状态机给出合法方向；真正安全性来自预检、锁外钩子、复核、唯一提交点和失败回滚。")

    load_states = [
        ("Discovered", PALE_GRAY, MUTED),
        ("Verified", PALE_PURPLE, PURPLE),
        ("Loaded", PALE_BLUE, BLUE),
        ("Linked", PALE_TEAL, TEAL),
        ("Ready", PALE_BLUE, BLUE),
        ("Active", NAVY, BLUE),
    ]
    x = 3.00
    for index, (name, fill, accent) in enumerate(load_states):
        width = 1.39 if name != "Discovered" else 1.57
        step(slide, x, 2.64, width, 0.65, name, fill=fill, accent=accent)
        if index < len(load_states) - 1:
            flow(slide, x + width, 2.97, x + width + 0.16, 2.97, color=accent)
        x += width + 0.16

    step(slide, 3.00, 3.64, 2.65, 1.13, "Quiescing", "停止新工作并等待排空", fill=PALE_BLUE, accent=BLUE)
    step(slide, 5.98, 3.64, 2.65, 1.13, "Paused", "保留镜像，可恢复", fill=PALE_PURPLE, accent=PURPLE)
    step(slide, 8.96, 3.64, 1.73, 1.13, "Detached", "摘除拓扑", fill=PALE_GRAY, accent=MUTED)
    step(slide, 10.99, 3.64, 1.61, 1.13, "Retired", "完成回收", fill=PALE_TEAL, accent=TEAL)
    flow(slide, 5.65, 4.20, 5.94, 4.20, color=BLUE)
    flow(slide, 8.63, 4.20, 8.92, 4.20, color=PURPLE)
    flow(slide, 10.69, 4.20, 10.95, 4.20, color=MUTED)

    heading(slide, "Replace：Generation N → N+1", 3.00, 5.03, 3.50, size=19)
    for x, title, fill, accent in (
        (3.00, "影子装载", PALE_BLUE, BLUE),
        (5.10, "静默旧代", PALE_PURPLE, PURPLE),
        (7.20, "可选迁移", PALE_TEAL, TEAL),
        (9.30, "唯一提交点", NAVY, BLUE),
        (11.40, "回收旧代", PALE_GRAY, MUTED),
    ):
        step(slide, x, 5.47, 1.20, 0.59, title, fill=fill, accent=accent)
    boundary(slide, "迁移状态上限 64 KiB；默认迁移钩子返回不支持，direct-pinned importer 与长期资源可阻断替换。")


def draw_08(slide) -> None:
    set_title(slide, "策略、预算与资源所有权")
    lead(slide, "权限、容量和长期副作用是三个不同问题，分别由 Policy、Budget 与 Owned Resource 管理。")

    panel(
        slide,
        3.00,
        2.62,
        2.93,
        2.93,
        "Policy",
        "父级给出能力上限\n子级只能缩小权限\n管理能力不自动继承\n策略更新校验代际与 epoch",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=20,
        detail_size=16,
        center=True,
    )
    panel(
        slide,
        6.27,
        2.62,
        2.93,
        2.93,
        "Budget",
        "端口与队列\n订阅与 pending load\n镜像、栈、并发调用\n动态内存与故障记录",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=20,
        detail_size=16,
        center=True,
    )
    panel(
        slide,
        9.54,
        2.62,
        3.06,
        2.93,
        "Owned Resource",
        "任务、定时器与工作项\n回调、IRQ 与异步请求\n设备和自定义对象\n常驻操作表负责安全收束",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=20,
        detail_size=16,
        center=True,
    )
    chip(slide, "退役顺序", 3.00, 5.64, 1.32, fill=NAVY)
    body(
        slide,
        "stop admission → quiesce → cancel → drain → release",
        4.57,
        5.60,
        8.03,
        0.46,
        size=17.0,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    boundary(slide, "CPU 时间当前只记账与统计超额，不实施调度节流；部分设备资源可 Detach，但会阻断可回滚 Pause。")


def draw_09(slide) -> None:
    set_title(slide, "原生执行与故障边界")
    lead(slide, "调用门能够把受支持的 fault、panic 与 timeout 收束到固定出口，并隔离责任 Cell；它不是地址空间沙箱。")

    stages = [
        ("执行预检", "代际 / 策略 / 预算", PALE_BLUE, BLUE, 1.58),
        ("ElmGuard", "64 KiB 栈 + 双 Guard 页", PALE_PURPLE, PURPLE, 1.82),
        ("原生 ELM", "Hook · Provider · Entry", PALE_TEAL, TEAL, 1.72),
        ("异常", "fault / panic / timeout", PALE_GRAY, MUTED, 1.66),
        ("恢复出口", "重写 PC / SP / 返回值", NAVY, BLUE, 1.82),
    ]
    x = 3.00
    for index, (title, detail, fill, accent, width) in enumerate(stages):
        step(slide, x, 2.73, width, 1.32, title, detail, fill=fill, accent=accent)
        if index < len(stages) - 1:
            flow(slide, x + width, 3.35, x + width + 0.20, 3.35, color=accent)
        x += width + 0.20

    panel(
        slide,
        3.00,
        4.38,
        4.58,
        1.49,
        "Fault Dump",
        "Cell · 阶段 · PC · fault address · cause · recovery PC / SP",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=19,
        center=True,
    )
    panel(
        slide,
        7.90,
        4.38,
        4.70,
        1.49,
        "隔离标志与状态迁移",
        "普通原生故障先设置 isolated，阻止新 Provider / Binding / import；生命周期收束失败才进入 Quarantined。",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=19,
    )
    boundary(slide, "共享特权地址空间中的写入无法由调用门撤销；常驻 Kernel Provider 也不属于原生 ELM Guard。")


def draw_10(slide) -> None:
    set_title(slide, "可复核运行证据")
    lead(slide, "可追踪性不是增加日志：记录以 ElmId 与 sequence 关联，Generation、Binding、ticket 按类型补充。")

    add_rect(slide, 4.10, 2.58, 7.40, 0.62, NAVY)
    add_rect(slide, 4.10, 2.58, 0.09, 0.62, BLUE)
    heading(slide, "共同锚点", 4.36, 2.72, 1.34, 0.30, size=18, color=WHITE)
    body(
        slide,
        "ElmId · sequence；其余身份字段按记录类型补充",
        5.75,
        2.69,
        5.48,
        0.33,
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    evidence = [
        (3.00, "Snapshot", "当前事实", PALE_BLUE, BLUE),
        (4.98, "Event", "顺序变化", PALE_TEAL, TEAL),
        (6.96, "Audit", "主体与决策", PALE_PURPLE, PURPLE),
        (8.94, "Trace", "执行阶段", PALE_BLUE, BLUE),
        (10.92, "Journal", "哈希链顺序", PALE_GRAY, MUTED),
    ]
    for x, title, detail, fill, accent in evidence:
        step(slide, x, 3.52, 1.68, 1.02, title, detail, fill=fill, accent=accent)

    for x, w, title, detail, fill, accent in (
        (3.00, 4.58, "管理链", "elmctl → sys_elm_ctl → elm-mgr → Core", PALE_BLUE, BLUE),
        (7.88, 4.72, "只读诊断面", "/sys/kernel/elm · 19 个正式节点", PALE_TEAL, TEAL),
    ):
        add_rect(slide, x, 4.90, w, 0.86, fill)
        add_rect(slide, x, 4.90, 0.09, 0.86, accent)
        heading(slide, title, x + 0.25, 5.15, 1.42, 0.30, size=18, color=NAVY)
        body(
            slide,
            detail,
            x + 1.67,
            5.11,
            w - 1.90,
            0.36,
            size=16,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    boundary(slide, "Journal 默认运行在内存易失模式；当前 replay 主要恢复 trust epoch，不重建 Cell、队列或执行现场。")


def draw_11(slide) -> None:
    set_title(slide, "工程接入与完成边界")
    lead(slide, "ELM 已进入真实驱动和网络执行链；完成度按“运行闭环、受限实现、TODO”区分，而不是用目标态替代现状。")

    panel(
        slide,
        3.00,
        2.62,
        3.02,
        2.30,
        "Rust 开发框架",
        "声明：ElmModule + 模块宏\n启动：create / initialize\n退出：quiesce / finalize\npause、migrate、entry 为可选钩子",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=19,
        center=True,
    )
    panel(
        slide,
        6.29,
        2.62,
        3.02,
        2.30,
        "y / m / n",
        "m：生成受管 EKI\ny：静态集成，不创建动态 Cell\nn：不构建或打包\n实际模式由 .config 决定",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=19,
        center=True,
    )
    panel(
        slide,
        9.58,
        2.62,
        3.02,
        2.30,
        "真实模块链",
        "net.stack · net.loopback\nvirtio.framework · virtio.block\nnet.virtio\nBuildBound 进入统一治理",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=19,
        center=True,
    )
    chip(slide, "当前实现", 3.00, 5.12, 1.42, fill=TEAL)
    body(slide, "Cell / Generation、EKI、调用分流、资源归属与故障收束", 4.60, 5.09, 8.00, 0.43, size=16, bold=True, valign=MSO_ANCHOR.MIDDLE)
    chip(slide, "受限 / TODO", 3.00, 5.61, 1.42, fill=MUTED)
    body(slide, "设备 / VFS Provider 需显式注册；IRQ、DMA、MMIO、block、packet 的通用 Provider 契约尚未接通", 4.60, 5.58, 8.00, 0.43, size=16, bold=True, valign=MSO_ANCHOR.MIDDLE)
    boundary(slide, "下一部分用网络栈说明：治理可以进入真实数据面，但不会把每个网络包送进通用管理分发。")


def draw_12(slide) -> None:
    set_title(slide, "网络子系统责任分层")
    lead(slide, "常规 TCP/IP 能力只作完整性证明；核心设计是把用户 ABI、协议状态和设备队列分给三个稳定责任域。")

    layers = [
        ("POSIX syscall", "socket · bind · connect · sendmsg · recvmsg · poll", PALE_BLUE, BLUE),
        ("VFS / Socket facade", "fd、等待队列、readiness 与用户 ABI", PALE_BLUE, BLUE),
        ("常驻 Host / Broker", "worker、registrar、代际与 pinned call slot", PALE_GRAY, MUTED),
        ("net.stack ELM", "FlowShard、路由、邻居、PMTU 与协议定时器", PALE_PURPLE, PURPLE),
        ("net.virtio / loopback", "queue、buffer pool、DMA / IRQ 与设备生命周期", PALE_TEAL, TEAL),
    ]
    y = 2.48
    for index, (title, detail, fill, accent) in enumerate(layers):
        add_rect(slide, 3.00, y, 9.60, 0.57, fill)
        add_rect(slide, 3.00, y, 0.09, 0.57, accent)
        heading(slide, title, 3.25, y + 0.11, 2.68, 0.30, size=18, color=NAVY)
        body(slide, detail, 5.88, y + 0.09, 6.44, 0.36, size=16, color=INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if index < len(layers) - 1:
            flow(slide, 7.80, y + 0.57, 7.80, y + 0.72, color=accent, direction="down")
        y += 0.74
    boundary(slide, "IPv4 / IPv6、TCP、UDP、ICMP、Raw、路由、邻居、PMTU、分片重组与 VirtIO-net 已具备。")


def draw_13(slide) -> None:
    set_title(slide, "FlowShard 数据通路")
    lead(slide, "网络性能结构来自分片单写者、有界批次和代际固定调用，而不是增加一个全局轮询线程。")

    stages = [
        ("VirtIO 队列", "IRQ / 设备事件", PALE_TEAL, TEAL, 1.76),
        ("Host worker", "批量收集", PALE_GRAY, MUTED, 1.62),
        ("shard-turn", "direct-pinned", PALE_PURPLE, PURPLE, 1.68),
        ("FlowShard", "单写者协议状态", NAVY, BLUE, 1.70),
        ("TxPlan", "批量提交", PALE_BLUE, BLUE, 1.50),
    ]
    x = 3.00
    for index, (title, detail, fill, accent, width) in enumerate(stages):
        step(slide, x, 2.68, width, 1.16, title, detail, fill=fill, accent=accent)
        if index < len(stages) - 1:
            flow(slide, x + width, 3.26, x + width + 0.24, 3.26, color=accent)
        x += width + 0.24

    panel(
        slide,
        3.00,
        4.18,
        3.00,
        1.64,
        "分片原则",
        "shard 数 ≤ 活动 CPU 与 queue pair；单队列设备不制造虚假并行。",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=18,
    )
    panel(
        slide,
        6.30,
        4.18,
        3.00,
        1.64,
        "本地快速路径",
        "Busy、预算 / 代际 / 调用失败，或大块连续流量时回落 owner worker。",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=18,
    )
    panel(
        slide,
        9.60,
        4.18,
        3.00,
        1.64,
        "调用约束",
        "每 CPU 固定 slot；校验代际、结构、提交位与宿主允许地址范围。",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=18,
    )
    boundary(slide, "当前真实热路径不是 Nexus packet.rx / packet.tx；准确表述是池化缓冲、租约和批量边界。")


def draw_14(slide) -> None:
    set_title(slide, "套接字兼容与控制边界")
    lead(slide, "POSIX 兼容保留在常驻 VFS / facade，协议内部使用类型化命令、配置快照和分片状态，两侧可以分别演进。")

    panel(
        slide,
        3.00,
        2.58,
        4.48,
        2.84,
        "数据与等待通路",
        "fd → FileOps → Socket facade → flow command\n\nsocket / bind / listen / accept\nconnect / sendmsg / recvmsg / shutdown\npoll / epoll / close / signal interruption",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=20,
        detail_size=16,
    )
    panel(
        slide,
        8.12,
        2.58,
        4.48,
        2.84,
        "配置与控制通路",
        "SIOC ioctl → ConfigStore → ConfigSnapshot\nnetlink：查询与 dump\n\n地址、路由、邻居与 PMTU\nMTU、管理启停与链路状态",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=20,
        detail_size=16,
    )
    flow(slide, 7.48, 4.00, 8.08, 4.00, color=PURPLE)
    chip(slide, "常驻 facade", 5.82, 5.68, 3.96, fill=NAVY)
    boundary(slide, "SIOC 写路径主要覆盖 IPv4，netlink 只查询 / dump；正常就绪精确唤醒，终止事件按 socket 广播。")


def draw_15(slide) -> None:
    set_title(slide, "网络生命周期与性能边界")
    lead(slide, "网络实例证明 ELM 管理的是实际状态、调用和资源，而不只是模块名称；安全退出优先于伪装成无损迁移。")

    rows = [
        ("net.stack", "建 FlowShard → 注册 shard / local endpoint → quiesce → begin_remove → 销毁代际", PALE_PURPLE, PURPLE),
        ("net.virtio", "校验 framework → 注册 MMIO / PCI driver → 停队列 → 分离设备 → 注销 driver", PALE_TEAL, TEAL),
        ("net.loopback", "建回环队列 → 注册设备 → 静默队列 → 销毁状态 → 注销设备", PALE_BLUE, BLUE),
    ]
    y = 2.55
    for title, detail, fill, accent in rows:
        add_rect(slide, 3.00, y, 9.60, 0.78, fill)
        add_rect(slide, 3.00, y, 0.09, 0.78, accent)
        heading(slide, title, 3.25, y + 0.18, 1.65, 0.34, size=19)
        body(slide, detail, 4.95, y + 0.14, 7.34, 0.45, size=16, color=INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        y += 0.96

    for x, w, title, detail, fill, accent in (
        (3.00, 3.02, "性能来源", "复制 · 推进 · 竞争 · 唤醒", PALE_BLUE, BLUE),
        (6.28, 3.02, "稳定失败", "NetworkDown / hangup", PALE_PURPLE, PURPLE),
        (9.56, 3.04, "后续方向", "等待者选择 · packet Provider", PALE_GRAY, MUTED),
    ):
        add_rect(slide, x, 5.54, w, 0.50, fill)
        add_rect(slide, x, 5.54, 0.09, 0.50, accent)
        heading(slide, title, x + 0.20, 5.65, 1.00, 0.28, size=17, color=NAVY)
        body(
            slide,
            detail,
            x + 1.13,
            5.63,
            w - 1.27,
            0.30,
            size=16,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    boundary(slide, "已验证 Detach + Reload：旧连接稳定失效，新 Cell / handle 生效；原位 Replace 与状态迁移尚未完成。")


def explain_01(slide) -> None:
    set_title(slide, "可拓展内核单元（ELM）的概念定位")
    lead(slide, "ELM 把扩展从一段可装入代码提升为可治理的责任单元：身份、能力、资源、故障和退出都围绕同一条 Cell 记录组织。")

    panel(
        slide,
        3.00,
        2.55,
        4.42,
        2.45,
        "Cell · 一项扩展的运行记录",
        "Cell 是一项逻辑服务的稳定责任边界；实现可以替换，但身份、准入和退出责任保持可追踪。所有公开能力、资源和故障都归属于该边界。",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=20,
    )
    panel(
        slide,
        7.75,
        2.55,
        4.85,
        2.45,
        "LKM / KLD · 约束分散在不同机制",
        "传统模块的装载、注册、回调和资源往往由不同机制分别管理，状态难以统一复核，安全退出也依赖各子系统自行约定。",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=20,
    )
    compact_row(
        slide,
        5.22,
        "Cell 记录的运行事实",
        "身份与代际 · 能力关系 · 活动引用 · 资源归属 · 退役阻断项",
        fill=PALE_TEAL,
        accent=TEAL,
        h=0.70,
        title_w=2.74,
    )
    boundary(
        slide,
        "网络示例：net.stack Cell 统一持有协议状态与资源责任，FlowShard 的运行活动可被同一责任边界追踪。",
        label="net.stack 例",
    )


def explain_02(slide) -> None:
    set_title(slide, "常驻 Core、elm-mgr 与格式解析器")
    lead(slide, "管理自举的关键是分层：常驻 Core 保持不变量，管理器组织意图，格式解析器提供输入投影；每一层都受同一责任拓扑约束。")

    panel(
        slide,
        3.00,
        2.52,
        2.90,
        2.62,
        "Core · 常驻事实与提交",
        "Core 是不可替换的最小可信根，统一持有关系、策略和事务不变量，只负责核验与提交，不承载具体子系统业务。",
        fill=NAVY,
        accent=BLUE,
        title_size=19,
    )
    panel(
        slide,
        6.24,
        2.52,
        2.85,
        2.62,
        "elm-mgr · ElmId 1",
        "elm-mgr 是受管的管理入口，负责把外部意图整理为 Core 可验证的操作；它不能绕过 Core 直接改写权威状态。",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=19,
    )
    panel(
        slide,
        9.43,
        2.52,
        3.17,
        2.62,
        "eki · ElmId 2",
        "eki 是受管的格式投影器，把某种镜像表示转换为统一的候选装载对象；格式解释与是否允许执行由 Core 分离判断。",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=19,
    )
    flow(slide, 5.90, 3.83, 6.20, 3.83, color=BLUE)
    flow(slide, 9.09, 3.83, 9.39, 3.83, color=PURPLE)
    compact_row(
        slide,
        5.42,
        "外部管理请求路径",
        "管理意图 → 统一预检 → 原子提交 → 返回状态与阻断原因",
        fill=PALE_GRAY,
        accent=MUTED,
        h=0.70,
        title_w=2.66,
    )
    boundary(
        slide,
        "管理器与解析器也受 Cell 生命周期管理；只有保存全局不变量的常驻 Core 不可被普通扩展替换。",
        label="不可替换部分",
    )


def explain_03(slide) -> None:
    set_title(slide, "Cell 身份与实现代际")
    lead(slide, "稳定身份与可替换实现分离：Cell 表示逻辑服务，Generation 表示当前实现；只有新实现完整通过验证并提交后才切换代次。")

    panel(
        slide,
        3.00,
        2.52,
        4.05,
        2.74,
        "一条 Cell 记录保存的关键事实",
        "逻辑身份\n实现代际\n当前状态\n管理归属\n权限与额度\n来源证明",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=19,
    )
    heading(slide, "一般动态 Cell 的代次切换示意", 7.42, 2.57, 5.18, 0.36, size=20)
    for x, label, detail, fill, accent in (
        (7.42, "G1", "初次装载", PALE_GRAY, MUTED),
        (9.19, "G2", "首次成功切换", PALE_PURPLE, PURPLE),
        (10.96, "G3", "再次成功切换", PALE_TEAL, TEAL),
    ):
        step(slide, x, 3.12, 1.46, 1.02, label, detail, fill=fill, accent=accent)
    flow(slide, 8.88, 3.63, 9.15, 3.63, color=MUTED)
    flow(slide, 10.65, 3.63, 10.92, 3.63, color=PURPLE)
    body(
        slide,
        "同一逻辑服务的身份保持稳定；新代成功提交后才递增，旧引用随之失效，避免新旧实现同时接收调用。",
        7.42,
        4.48,
        5.18,
        0.80,
        size=16,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    compact_row(
        slide,
        5.52,
        "强类型 ID",
        "不同对象使用不同的强类型身份，不能相互混用，也不把身份等同于实现地址。",
        fill=PALE_TEAL,
        accent=TEAL,
        h=0.62,
        title_w=1.62,
    )
    boundary(
        slide,
        "长期引用必须与当前 Cell 和 Generation 重新匹配；代际变化自动形成清晰的失效边界。",
        label="失效规则",
    )


def explain_04(slide) -> None:
    set_title(slide, "关系图中的四种责任")
    lead(slide, "关系图把不同责任分开表达：管理归属、服务依赖、允许拓展的位置和已提交的调用许可各自有边界，避免用一种关系代替另一种权限。")

    panel(
        slide,
        3.00,
        2.50,
        4.48,
        1.48,
        "父子关系",
        "表达管理归属与额度委派；父子关系本身不授予业务调用权。",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=19,
    )
    panel(
        slide,
        7.80,
        2.50,
        4.80,
        1.48,
        "依赖关系",
        "表达服务依赖与兼容关系；依赖可被追踪并成为安全退出的约束，但不授予管理权限。",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=19,
    )
    panel(
        slide,
        3.00,
        4.18,
        4.48,
        1.48,
        "拓展点与拓展项",
        "目标先公开受控的拓展位置和契约，实现只能按声明挂接，避免任意修改核心行为。",
        fill=PALE_PURPLE,
        accent=PURPLE,
        title_size=19,
    )
    panel(
        slide,
        7.80,
        4.18,
        4.80,
        1.48,
        "能力绑定",
        "Binding 表示已经提交的能力连接；调用许可可审计、可撤销，而不是暴露裸函数入口。",
        fill=PALE_GRAY,
        accent=MUTED,
        title_size=19,
    )
    boundary(
        slide,
        "关系在提交前统一校验，撤销时先解除连接再回收对象；图一致性成为可复核的不变量。",
        label="图一致性",
    )


def explain_05(slide) -> None:
    set_title(slide, "能力发布：Contract、Port 与 Provider")
    lead(slide, "ELM 将服务拆成契约、端口和执行后端：契约稳定语义，端口表达可发现的能力，Provider 承担具体执行，使核心与实现解耦。")

    process_box(slide, 3.00, 2.58, 2.92, 2.05, "1", "流契约（Contract）", "统一命名版本和语义，约束双方对数据与错误的共同理解。", fill=PALE_PURPLE, accent=PURPLE)
    process_box(slide, 6.26, 2.58, 2.92, 2.05, "2", "连接点（Port）", "公开方向、访问范围和发现属性，隔离调用者与实现地址。", fill=PALE_BLUE, accent=BLUE)
    process_box(slide, 9.52, 2.58, 3.08, 2.05, "3", "执行者（Provider）", "承载真实服务，可由常驻能力或受管 ELM 提供；缺少后端时明确返回未实现。", fill=PALE_TEAL, accent=TEAL)
    flow(slide, 5.92, 3.49, 6.22, 3.49, color=PURPLE)
    flow(slide, 9.18, 3.49, 9.48, 3.49, color=BLUE)
    compact_row(slide, 4.92, "管理能力", "管理请求也通过版本化契约和受控 Provider 进入统一治理。", fill=PALE_GRAY, accent=MUTED, h=0.66, title_w=2.74)
    compact_row(slide, 5.57, "调用优势", "契约稳定、实现可替换；调用前检查状态与权限，避免悬空或越权连接。", fill=PALE_TEAL, accent=TEAL, h=0.66, title_w=2.74)
    boundary(
        slide,
        "连接只有在契约、状态、权限和后端同时满足时才生效；撤销可以沿同一责任链追踪。",
        label="调用门槛",
    )


def explain_06(slide) -> None:
    set_title(slide, "能力连接：Binding 与 Lease")
    lead(slide, "Binding 表达已经批准的能力连接，Lease 表达仍在使用的责任；二者让动态调用可撤销、可追踪，并把短期引用与长期资源区分开。")

    numbered_row(slide, 2.43, "1", "提出连接", "声明消费者、目标能力和预期契约，明确双方责任。", fill=PALE_BLUE, accent=BLUE, h=0.65)
    numbered_row(slide, 3.16, "2", "资格预检", "统一核对状态、策略、契约、代际和资源边界，拒绝不一致关系。", fill=PALE_PURPLE, accent=PURPLE, h=0.65)
    numbered_row(slide, 3.89, "3", "提交使用权", "预检成功后公开连接，并记录可撤销的使用责任。", fill=PALE_TEAL, accent=TEAL, h=0.65)
    compact_row(slide, 4.64, "调用期间", "动态或异步调用持有引用；调用结束后释放，长期资源另由所属 Cell 负责清理。", fill=PALE_BLUE, accent=BLUE, h=0.66, title_w=2.30)
    compact_row(slide, 5.39, "安全撤销", "仍有活动引用时返回 Busy 且不改变对象；引用归零后再撤销连接并回收。", fill=PALE_GRAY, accent=MUTED, h=0.66, title_w=2.30)
    boundary(
        slide,
        "连接只传递受约束的契约数据，不暴露内核裸地址；统一边界让撤销、审计和故障归属保持一致。",
        label="设计优势",
    )


def explain_07(slide) -> None:
    set_title(slide, "调用边界、核验位置与运行成本")
    lead(slide, "调用路径按治理需求分层：动态 Provider 便于发现与撤销，受管导入强调代际安全，固定直连服务热路径，常驻符号保持最小可信面。")

    definition_row(slide, 2.42, "Provider 枢纽", "适合动态发现、授权、异步和审计；撤销边界清晰。", "治理完整、成本较高", fill=PALE_BLUE, accent=BLUE, h=0.76)
    definition_row(slide, 3.31, "受管导入", "适合代际敏感接口；实现变化会让旧连接明确失效。", "逐次保持一致", fill=PALE_PURPLE, accent=PURPLE, h=0.76)
    definition_row(slide, 4.20, "固定导入", "适合稳定热路径；装载时证明兼容，运行时仍保留必要门禁。", "性能高、替换受约束", fill=PALE_TEAL, accent=TEAL, h=0.76)
    definition_row(slide, 5.09, "内核符号", "只连接常驻、明确登记的核心能力，缩小可信范围。", "不参与动态代际路由", fill=PALE_GRAY, accent=MUTED, h=0.76)
    boundary(
        slide,
        "网络热路径采用受管直连以避免通用分发开销，同时保留状态、代际和参数边界核验。",
        label="网络创新",
    )


def explain_08(slide) -> None:
    set_title(slide, "从 EKI 文件到候选 EBI")
    lead(slide, "格式解释与装载决策分离：解析器只产生统一的候选描述，Core 再独立判断来源、兼容性和执行资格。")

    panel(slide, 3.00, 2.52, 2.88, 2.62, "镜像表示", "描述代码、依赖和能力的来源；封存后保持不可变，便于复核。", fill=PALE_BLUE, accent=BLUE, title_size=19)
    panel(slide, 6.24, 2.52, 2.92, 2.62, "投影源", "将不同格式转换为统一候选对象，只负责解释，不拥有激活权限。", fill=PALE_PURPLE, accent=PURPLE, title_size=18)
    panel(slide, 9.52, 2.52, 3.08, 2.62, "统一装载对象", "集中表达清单、能力、生命周期和接口证明，供 Core 使用同一套规则审查。", fill=PALE_TEAL, accent=TEAL, title_size=19)
    flow(slide, 5.88, 3.83, 6.20, 3.83, color=BLUE)
    flow(slide, 9.16, 3.83, 9.48, 3.83, color=PURPLE)
    compact_row(slide, 5.40, "可演进性", "新增格式只需增加投影源，Core 无需理解具体文件细节；统一装载对象保持稳定。", fill=PALE_GRAY, accent=MUTED, h=0.70, title_w=2.28)
    boundary(
        slide,
        "投影与执行权限分离，所有候选都经过同一套来源、接口和策略审查。",
        label="设计优势",
    )


def explain_09(slide) -> None:
    set_title(slide, "装载证明：来源与接口兼容")
    lead(slide, "候选实现必须证明来源、完整性与接口兼容，证明通过后才获得进入运行态的资格。")

    numbered_row(slide, 2.46, "1", "内容完整", "输入保持不可变且可追踪，保证后续判断针对同一份实现。", fill=PALE_BLUE, accent=BLUE, h=0.75)
    numbered_row(slide, 3.35, "2", "来源可信", "来源与信任关系满足策略，避免未经授权的实现进入核心。", fill=PALE_PURPLE, accent=PURPLE, h=0.75)
    numbered_row(slide, 4.24, "3", "构建一致", "实现与当前核心契约、依赖和发布状态一致，避免接口漂移。", fill=PALE_TEAL, accent=TEAL, h=0.75)
    numbered_row(slide, 5.13, "4", "目标兼容", "架构、能力范围和执行边界兼容，才允许继续初始化。", fill=PALE_GRAY, accent=MUTED, h=0.75)
    boundary(
        slide,
        "证明解决来源与接口风险，不替代业务正确性验证；任一不一致都在公开前拒绝。",
        label="证明范围",
    )


def explain_09b(slide) -> None:
    set_title(slide, "装载提交：一次性公开")
    lead(slide, "装载采用准备、初始化、公开提交的事务结构；任何阶段失败都不让半成品进入公共拓扑。")

    numbered_row(slide, 2.46, "1", "候选准备", "准备实现及其能力关系，保持外部不可见。", fill=PALE_BLUE, accent=BLUE, h=0.75)
    numbered_row(slide, 3.35, "2", "初始化", "在受控上下文中完成初始化；失败时销毁候选并恢复原状态。", fill=PALE_PURPLE, accent=PURPLE, h=0.75)
    numbered_row(slide, 4.24, "3", "公开门禁", "初始化成功后才开放关系和调用资格，避免半初始化服务被消费。", fill=PALE_TEAL, accent=TEAL, h=0.75)
    numbered_row(slide, 5.13, "4", "原子提交", "一次性发布状态与所有权；失败回滚并留下可追踪的隔离事实。", fill=PALE_GRAY, accent=MUTED, h=0.75)
    boundary(
        slide,
        "唯一公开点把状态、关系和所有权绑定在一起；调用者不会观察到半提交状态。",
        label="可见性边界",
    )


def explain_10(slide) -> None:
    set_title(slide, "生命周期状态与调用门禁")
    lead(slide, "生命周期状态回答扩展是否可用、可暂停或可退出；调用门同时检查身份、权限、关系和资源责任。")

    compact_row(slide, 2.42, "准备阶段", "从发现、验证到链接，逐步形成可运行的候选实现。", fill=PALE_BLUE, accent=BLUE, h=0.82, title_w=1.62)
    compact_row(slide, 3.38, "激活阶段", "初始化成功后公开能力；只有满足调用门禁才进入 Active。", fill=PALE_TEAL, accent=TEAL, h=0.82, title_w=1.62)
    compact_row(slide, 4.34, "退出阶段", "先停止新工作并排空，再选择可恢复的暂停或最终摘除。", fill=PALE_PURPLE, accent=PURPLE, h=0.82, title_w=1.62)
    compact_row(slide, 5.30, "故障处理", "故障先隔离责任单元；无法安全收束时进入受控的终止状态。", fill=PALE_GRAY, accent=MUTED, h=0.82, title_w=1.62)
    boundary(
        slide,
        "Policy、Generation、关系、Lease、活动执行与长期资源同时满足后才提交状态迁移。",
        label="提交条件",
    )


def explain_11(slide) -> None:
    set_title(slide, "Pause 与 Detach 的执行步骤")
    lead(slide, "Pause 与 Detach 是两种不同的退役策略：前者保留恢复可能，后者完成最终摘除；共同原则是先静默、再收束。")

    numbered_row(slide, 2.42, "1", "资格判断", "确认当前状态、依赖和活动责任允许开始退役。", fill=PALE_BLUE, accent=BLUE, h=0.64)
    numbered_row(slide, 3.18, "2", "停止接纳", "关闭新的调用入口，让在途工作自然收束。", fill=PALE_PURPLE, accent=PURPLE, h=0.64)
    numbered_row(slide, 3.94, "3", "排空与回收", "清理活动引用和长期资源，再提交最终状态。", fill=PALE_TEAL, accent=TEAL, h=0.64)
    compact_row(slide, 4.76, "Pause · 可回滚", "保留实现与关系，失败时恢复为可运行状态。", fill=PALE_BLUE, accent=BLUE, h=0.60, title_w=2.14)
    compact_row(slide, 5.46, "Detach · 不可逆", "完成资源收束后摘除公开关系，后续可重新装载新代。", fill=PALE_GRAY, accent=MUTED, h=0.60, title_w=2.14)
    boundary(
        slide,
        "Pause 与 Detach 使用不同 blocker 集合；钩子或资源回滚失败时记录故障阶段并隔离 Cell。",
        label="阻断条件",
    )


def explain_12(slide) -> None:
    set_title(slide, "Replace 事务：ElmId 不变，Generation 递增")
    lead(slide, "Replace 保留逻辑身份，在新代准备完成后切换实现；提交失败不改变旧代，确保更新具有原子性。")

    stages = [
        ("1", "准备新代", "完成来源、接口和初始化证明"),
        ("2", "静默旧代", "停止新工作并等待活动责任结束"),
        ("3", "处理迁移", "仅在明确支持时迁移可转移状态"),
        ("4", "一次提交", "切换当前实现与能力关系"),
        ("5", "回收旧代", "释放旧实现并保留更新事实"),
    ]
    x = 3.00
    widths = [1.72, 1.72, 1.72, 1.72, 1.72]
    for index, (number, title, detail) in enumerate(stages):
        step(slide, x, 2.52, widths[index], 1.70, f"{number}  {title}", detail, fill=PALE_TEAL if number in {"1", "4"} else PALE_PURPLE, accent=TEAL if number in {"1", "4"} else PURPLE)
        if number != "5":
            flow(slide, x + widths[index], 3.37, x + widths[index] + 0.20, 3.37, color=PURPLE)
        x += widths[index] + 0.22
    panel(slide, 3.00, 4.48, 4.56, 1.58, "失败处理", "提交前失败：销毁新代并恢复旧代；收束失败：隔离责任单元并保留失败事实。", fill=PALE_BLUE, accent=BLUE, title_size=18)
    panel(slide, 7.86, 4.48, 4.74, 1.58, "替换阻断项", "仍在使用的连接、不可迁移的长期状态或接口不兼容都会阻止原位切换。", fill=PALE_GRAY, accent=MUTED, title_size=18)
    boundary(
        slide,
        "新旧代次之间保持清晰失效边界；不能安全迁移时宁可拒绝替换，也不破坏用户可见语义。",
        label="设计优势",
    )


def explain_13(slide) -> None:
    set_title(slide, "Policy、Budget、Lease 与长期资源")
    lead(slide, "权限、容量、活动引用和长期资源是不同责任，分别建模后才能判断扩展能否安全运行与退出。")

    panel(slide, 3.00, 2.46, 4.56, 1.56, "策略（Policy）· 能否执行", "限定能力范围，父级给出上限，子级只能进一步收窄。", fill=PALE_PURPLE, accent=PURPLE, title_size=19)
    panel(slide, 7.86, 2.46, 4.74, 1.56, "预算（Budget）· 最多占用多少", "限定并发、队列和内存等资源额度，防止局部扩展无界消耗。", fill=PALE_BLUE, accent=BLUE, title_size=19)
    panel(slide, 3.00, 4.24, 4.56, 1.56, "租约（Lease）· 是否仍被使用", "记录短期活动引用；引用未归零时拒绝撤销，避免悬空调用。", fill=PALE_TEAL, accent=TEAL, title_size=19)
    panel(slide, 7.86, 4.24, 4.74, 1.56, "长期资源（Owned Resource）· 谁回收", "记录需要异步清理的对象，并由常驻责任方执行收束。", fill=PALE_GRAY, accent=MUTED, title_size=19)
    boundary(
        slide,
        "停止接纳 → 静默 → 取消 → 排空 → 逆序释放；资源归属清晰使局部故障不会扩散。",
        label="退役协议",
    )


def explain_14(slide) -> None:
    set_title(slide, "原生调用门的故障收束边界")
    lead(slide, "原生扩展与内核共享特权空间，不能承诺恶意写入隔离；调用门的创新在于把可恢复故障收束到固定责任边界。")

    panel(slide, 3.00, 2.48, 2.92, 2.34, "进入前校验", "核对身份、策略、入口和能力范围，确认调用属于当前责任单元。", fill=PALE_BLUE, accent=BLUE, title_size=19)
    panel(slide, 6.24, 2.48, 2.92, 2.34, "受控执行现场", "在独立的受控上下文中执行扩展，避免故障状态直接污染调用者。", fill=PALE_PURPLE, accent=PURPLE, title_size=19)
    panel(slide, 9.48, 2.48, 3.12, 2.34, "异常记录与固定出口", "记录故障位置、原因和责任归属，统一返回错误并释放执行责任。", fill=PALE_TEAL, accent=TEAL, title_size=19)
    compact_row(slide, 5.02, "故障后的 Cell", "故障单元先隔离并拒绝新调用；无法安全收束时进入受控终止状态。", fill=PALE_BLUE, accent=BLUE, h=0.54, title_w=2.12)
    boundary(
        slide,
        "调用门解决可恢复故障和责任归属，不把共享特权地址空间误称为安全沙箱。",
        label="恢复边界",
        label_w=1.72,
    )


def explain_15(slide) -> None:
    set_title(slide, "ELM 的五类运行证据")
    lead(slide, "运行时提供五类独立查询视图：当前状态、变化顺序、操作主体、控制路径和管理日志。它们各自排序，不能把一张快照误当成完整执行历史。")

    definition_row(slide, 2.38, "快照 Snapshot", "Cell / Port 基础快照给出状态；Binding、拓扑、执行与资源由对应查询视图补齐。", "拼出当时条件", fill=PALE_BLUE, accent=BLUE, h=0.66, rule_w=1.88, term_w=1.92)
    definition_row(slide, 3.14, "事件 Event", "一类状态变化的顺序记录；读取方可以从该事件流的 sequence 继续消费。", "还原该事件流顺序", fill=PALE_TEAL, accent=TEAL, h=0.66, rule_w=1.88, term_w=1.92)
    definition_row(slide, 3.90, "审计 Audit", "记录操作主体、管理动作、结果、拒绝码和具名阻断项。", "核对谁做了什么", fill=PALE_PURPLE, accent=PURPLE, h=0.66, rule_w=1.88, term_w=1.92)
    definition_row(slide, 4.66, "路径 Trace", "生命周期、Provider、拓展、Replace、策略和资源的操作级记录：动作、对象、结果、阻断项。", "定位失败类别", fill=PALE_BLUE, accent=BLUE, h=0.66, rule_w=1.88, term_w=1.92)
    definition_row(slide, 5.42, "日志链 Journal", "记录管理过程的连续证据，支持顺序和完整性复核。", "检查顺序与完整性", fill=PALE_GRAY, accent=MUTED, h=0.66, rule_w=1.88, term_w=1.92)
    boundary(
        slide,
        "各视图以对象身份和事务关联；证据用于复核责任与顺序，不替代持久化状态或业务结果。",
        label="关联与持久性",
    )


def explain_15b(slide) -> None:
    set_title(slide, "Busy 拒绝记录样例")
    lead(slide, "当替换或撤销遇到活动引用时，系统返回 Busy 并保持原状态；多类证据共同解释拒绝原因。")

    numbered_row(slide, 2.46, "1", "请求事实", "记录谁发起了替换、目标是谁以及返回的阻断类别。", fill=PALE_BLUE, accent=BLUE, h=0.75)
    numbered_row(slide, 3.35, "2", "关系事实", "指出哪条能力连接或活动引用仍然存在。", fill=PALE_PURPLE, accent=PURPLE, h=0.75)
    numbered_row(slide, 4.24, "3", "执行事实", "把阻断引用关联到正在进行的调用或异步工作。", fill=PALE_TEAL, accent=TEAL, h=0.75)
    numbered_row(slide, 5.13, "4", "状态事实", "证明拒绝没有半提交，引用归零后才允许再次尝试。", fill=PALE_GRAY, accent=MUTED, h=0.75)
    boundary(
        slide,
        "请求、关系、执行和状态四类证据闭合后，Busy 不再是黑盒错误，而是可定位、可复核的安全决策。",
        label="本例所需记录",
    )


def explain_16(slide) -> None:
    set_title(slide, "实现接入策略与边界")
    lead(slide, "同一业务能力可以选择受管扩展、常驻集成或暂不接入；选择决定它获得的治理能力和替换边界。")

    panel(slide, 3.00, 2.48, 4.48, 2.32, "受管扩展", "获得身份、能力关系、生命周期和故障证据；适合需要独立演进与可撤销的功能。", fill=PALE_BLUE, accent=BLUE, title_size=20)
    panel(slide, 7.80, 2.48, 4.80, 2.32, "常驻集成", "直接成为核心的一部分，路径短、可信边界小，但不具备动态代际替换。", fill=PALE_PURPLE, accent=PURPLE, title_size=20)
    compact_row(slide, 5.10, "选择原则", "把治理收益、热路径成本和替换需求放在同一决策中权衡。", fill=PALE_TEAL, accent=TEAL, h=0.76, title_w=1.62)
    compact_row(slide, 5.99, "未接入能力", "尚未登记的能力不伪装成已完成；先保留清晰边界，再按同一契约接入。", fill=PALE_GRAY, accent=MUTED, h=0.64, title_w=1.62)


def explain_16b(slide) -> None:
    set_title(slide, "当前接入范围与未完成边界")
    lead(slide, "把“代码已有”“生产已接通”和“仍需接入”分开列出；模块声明、运行时注册和真实数据路径不是同一件事。")

    panel(slide, 3.00, 2.46, 4.56, 1.72, "治理能力已具备", "身份、代际、能力关系、生命周期、资源归属、故障隔离和运行证据形成完整框架。", fill=PALE_TEAL, accent=TEAL, title_size=19)
    panel(slide, 7.86, 2.46, 4.74, 1.72, "网络路径已接通", "协议状态与设备状态分别受管，热路径采用受控直连，用户接口保持稳定。", fill=PALE_BLUE, accent=BLUE, title_size=19)
    panel(slide, 3.00, 4.46, 4.56, 1.46, "接入边界", "设备和 VFS 等能力仍需显式登记，不能由模块声明自动推断。", fill=PALE_PURPLE, accent=PURPLE, title_size=19)
    panel(slide, 7.86, 4.46, 4.74, 1.46, "后续方向", "通用数据面 Provider、更多格式投影源和公共发布体系继续沿同一契约扩展。", fill=PALE_GRAY, accent=MUTED, title_size=19)
    boundary(
        slide,
        "创新已在网络路径落地；其他能力按需接入，避免把未完成范围误写成现状。",
        label="完成边界",
    )


def explain_17(slide) -> None:
    set_title(slide, "网络能力与责任结构")
    lead(slide, "网络栈作为 ELM 的真实实现，重点不在罗列协议，而在把用户可见语义与可退役的协议、设备责任分开。")

    panel(slide, 3.00, 2.45, 4.55, 1.66, "稳定用户接口", "POSIX socket 与事件等待继续由常驻文件对象解释；协议实现可以演进，fd 语义不随代际漂移。", fill=PALE_BLUE, accent=BLUE, title_size=19, detail_size=15)
    panel(slide, 7.86, 2.45, 4.74, 1.66, "可退役网络资源", "协议状态、设备队列和回环路径分别由网络 ELM 持有，退出时各自排空并释放。", fill=PALE_TEAL, accent=TEAL, title_size=19, detail_size=15)
    compact_row(slide, 4.38, "解决的问题", "传统全局网络对象难以热替换、难以归因；责任拆分让故障、代际和资源都有明确归属。", fill=PALE_GRAY, accent=MUTED, h=0.68, title_w=1.76)
    compact_row(slide, 5.20, "网络创新", "稳定 SocketFacade、代际 Broker、FlowShard 单写者和受控热路径共同兼顾可演进性与性能。", fill=PALE_PURPLE, accent=PURPLE, h=0.68, title_w=1.76)
    boundary(
        slide,
        "网络栈是 ELM 的生产案例：治理面负责代际和退役，数据面保持低开销固定入口。",
        label="网络定位",
    )


def explain_18(slide) -> None:
    set_title(slide, "网络对象的所有权与退出责任")
    lead(slide, "网络状态不集中在一个全局对象中：常驻对象维持用户可见语义，可退役 Cell 持有协议或设备状态，代际路由器只把调用送往当前有效实现。")

    compact_row(slide, 2.42, "SocketFacade", "常驻于 VFS，维持 fd、缓冲、等待和错误等用户可见语义，不绑定具体协议实现。", fill=PALE_BLUE, accent=BLUE, h=0.76, title_w=2.12)
    compact_row(slide, 3.31, "Host", "常驻协调网络配置、工作推进和排空，但不拥有可替换协议状态。", fill=PALE_GRAY, accent=MUTED, h=0.76, title_w=2.12)
    compact_row(slide, 4.20, "Broker", "把调用路由到当前有效代际，只允许 Active 且就绪的网络实现接收工作。", fill=PALE_PURPLE, accent=PURPLE, h=0.76, title_w=2.12)
    compact_row(slide, 5.09, "网络 ELM", "协议、设备和回环各自持有自己的状态与资源，形成可独立退役的责任单元。", fill=PALE_TEAL, accent=TEAL, h=0.76, title_w=2.12)
    boundary(
        slide,
        "退出时先关闭 Cell 入口并排空在途工作；SocketFacade 继续存在，负责把失效代际转换为稳定的 POSIX 错误与就绪状态。",
        label="退出责任",
    )


def explain_19(slide) -> None:
    set_title(slide, "FlowShard 单写者模型与双执行路径")
    lead(slide, "FlowShard 是协议状态的单写者分区：同一分区同时只有一个执行者修改状态，把并发复杂度收束到清晰的所有权边界。")

    panel(slide, 3.00, 2.46, 4.48, 2.42, "分区所有权", "相关流量进入固定分区，协议状态由该分区统一维护；状态、队列和定时责任不会散落到多个全局锁。", fill=PALE_BLUE, accent=BLUE, title_size=20, detail_size=15)
    panel(slide, 7.80, 2.46, 4.80, 2.42, "双执行路径", "后台工作者负责持续推进，短调用在可行时直接加入；竞争失败的一方交给唯一所有者处理，不并发写入。", fill=PALE_PURPLE, accent=PURPLE, title_size=20, detail_size=15)
    compact_row(slide, 4.98, "创新", "单写者模型同时保留批量推进和低延迟入口，减少锁竞争与状态交错。", fill=PALE_TEAL, accent=TEAL, h=0.50, title_w=1.72)
    compact_row(slide, 5.56, "优势", "并行度由真实资源约束决定，避免用虚假线程数换取不可控竞争。", fill=PALE_BLUE, accent=BLUE, h=0.46, title_w=1.72)
    boundary(
        slide,
        "并行度与硬件队列和在线资源匹配，避免为网络状态制造无效并发。",
        label="并行优势",
    )


def explain_19b(slide) -> None:
    set_title(slide, "数据通路与受管调用边界")
    lead(slide, "ELM 治理不要求每个报文经过 Core：Core 管理装载、关系和代际，收发数据沿固定入口流动，调用前仍核验目标是否有效。")

    panel(slide, 3.00, 2.42, 4.56, 1.48, "管理面：Core 提交关系", "装载、更新、暂停和退役只改变责任关系与代际门禁，不把管理逻辑混入每个报文。", fill=PALE_PURPLE, accent=PURPLE, title_size=19, detail_size=15)
    panel(slide, 7.86, 2.42, 4.74, 1.48, "数据面：固定入口", "收发沿稳定的数据路径进入当前有效网络实现，避免每个报文重复执行通用管理分发。", fill=PALE_TEAL, accent=TEAL, title_size=19, detail_size=15)
    compact_row(slide, 4.12, "收发链路", "接收和发送都经过设备、协议分区、稳定用户代理和文件描述符之间的责任链。", fill=PALE_TEAL, accent=TEAL, h=0.54, title_w=1.72)
    compact_row(slide, 4.77, "逐次门禁", "进入协议状态前核对当前代际和参数范围；失效实现不会继续接收数据。", fill=PALE_BLUE, accent=BLUE, h=0.54, title_w=1.72)
    compact_row(slide, 5.42, "竞争回退", "短调用未取得分区所有权时转为待处理，由唯一所有者继续推进。", fill=PALE_GRAY, accent=MUTED, h=0.54, title_w=1.72)
    boundary(
        slide,
        "治理面与数据面分离，让 Core 保持轻量而不牺牲代际和参数安全。",
        label="网络创新",
    )


def explain_20(slide) -> None:
    set_title(slide, "SocketFacade 的 POSIX 稳定边界")
    lead(slide, "SocketFacade 是常驻 VFS 代理：协议 Cell 可以退出，但 fd、等待关系和可观察错误仍由常驻对象解释；TCP 控制块仍归 net.stack 的 FlowShard。")

    compact_row(slide, 2.42, "POSIX 与 VFS", "用户接口先进入稳定的文件对象，由边界层解释阻塞、错误和等待语义。", fill=PALE_BLUE, accent=BLUE, h=0.76, title_w=2.30)
    compact_row(slide, 3.31, "SocketFacade", "保存用户可见的 fd、端点、缓冲、就绪和关闭状态，不把协议控制状态绑死在 fd 上。", fill=PALE_TEAL, accent=TEAL, h=0.76, title_w=2.30)
    compact_row(slide, 4.20, "Broker / FlowShard", "Broker 选择当前代际，FlowShard 以单写者方式推进协议，再把结果交回稳定代理。", fill=PALE_PURPLE, accent=PURPLE, h=0.76, title_w=2.30)
    compact_row(slide, 5.08, "等待与退出", "正常数据只唤醒必要的等待者；退出、关闭和错误会让所有相关等待者看到稳定的终止语义。", fill=PALE_BLUE, accent=BLUE, h=0.92, title_w=2.38)
    boundary(
        slide,
        "协议 Cell 退出后旧 fd 仍是有效 VFS 对象；网络调用返回 NetworkDown，poll / epoll 观察 ERROR 或 HANGUP。",
        label="退出可见性",
    )


def explain_21(slide) -> None:
    set_title(slide, "网络退役协议、用户语义与成本归因")
    lead(slide, "net.stack 与 net.virtio 分别拥有协议状态和设备资源，因此有不同的退出顺序；共同原则是先关闭入口，再排空在途工作，最后回收本代对象。")

    compact_row(slide, 2.42, "协议 ELM", "关闭新入口、排空在途协议工作，再回收本代状态。", fill=PALE_PURPLE, accent=PURPLE, h=0.88, title_w=1.58)
    compact_row(slide, 3.44, "设备 ELM", "停止设备数据流、排空队列并释放设备责任，避免协议层继续引用旧资源。", fill=PALE_TEAL, accent=TEAL, h=0.88, title_w=1.58)
    compact_row(slide, 4.46, "用户语义", "旧 fd 仍是有效文件对象，但把失效代际转换为稳定的错误和挂断事件；新建对象进入新代。", fill=PALE_BLUE, accent=BLUE, h=0.88, title_w=1.58)
    compact_row(slide, 5.50, "性能归因", "用户复制、协议推进、队列竞争、等待唤醒和设备完成分别计入责任边界。", fill=PALE_GRAY, accent=MUTED, h=0.52, title_w=2.42)
    compact_row(slide, 6.12, "网络优势", "协议与设备可独立退役，POSIX 语义保持稳定；无法安全迁移时明确拒绝而非隐式破坏连接。", fill=PALE_PURPLE, accent=PURPLE, h=0.52, title_w=2.42)


DRAWERS = [
    explain_01,
    explain_04,
    explain_12,
    explain_14,
    explain_17,
    explain_19,
]


def create_topic_slides(prs: Presentation, template) -> list:
    slides = []
    for drawer in DRAWERS:
        slide = clone_slide(prs, template)
        clean_body_template(slide)
        drawer(slide)
        enforce_font_floor(slide)
        slides.append(slide)
    return slides


def build_topic(base: Path, output: Path) -> None:
    source = Presentation(base)
    remove_existing_topic(source)
    template = find_slide(source, "设备抽象能力闭环")
    slides = create_topic_slides(source, template)
    keep = {slide.part for slide in slides}
    for slide in list(source.slides):
        if slide.part not in keep:
            remove_slide(source, slide)

    output.parent.mkdir(parents=True, exist_ok=True)
    source.save(output)


def insert_topic(full: Presentation) -> None:
    anchor = find_slide(full, TOPIC_START_TITLE)
    anchor_index = list(full.slides).index(anchor)
    old_topic = existing_topic_slides(full)

    # Create the replacement pages while the old pages are still present.  This
    # keeps python-pptx from reusing the old slide part names, which can produce
    # duplicate ppt/slides/slide*.xml members in the saved ZIP package.
    added = create_topic_slides(full, anchor)
    added_parts = {slide.part for slide in added}

    for slide in old_topic:
        remove_slide(full, slide)

    slide_ids = full.slides._sldIdLst
    added_ids = [
        slide_id
        for slide_id in list(slide_ids)
        if full.part.related_part(slide_id.rId) in added_parts
    ]
    if len(added_ids) != len(added):
        raise RuntimeError(
            f"新专题页关系数量错误：expected={len(added)} actual={len(added_ids)}"
        )
    for slide_id in added_ids:
        slide_ids.remove(slide_id)
    for offset, slide_id in enumerate(added_ids, 1):
        slide_ids.insert(anchor_index + offset, slide_id)


def build_full(base: Path, output: Path) -> None:
    full = Presentation(base)
    insert_topic(full)
    remove_page_markers(full)
    output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        prefix=f".{output.stem}-",
        suffix=".pptx",
        dir=output.parent,
        delete=False,
    ) as temporary:
        temporary_path = Path(temporary.name)
    try:
        full.save(temporary_path)
        os.replace(temporary_path, output)
        output.chmod(0o644)
    except Exception:
        temporary_path.unlink(missing_ok=True)
        raise


def validate_topic(topic: Path) -> None:
    prs = Presentation(topic)
    if len(prs.slides) != len(DRAWERS):
        raise RuntimeError(f"专题页数错误：{len(prs.slides)}")
    for slide_number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if is_page_marker(shape):
                raise RuntimeError(f"第 {slide_number} 页仍包含右下角数字")
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if run.font.size is None or run.font.size.pt < MIN_FONT_PT:
                        raise RuntimeError(
                            f"第 {slide_number} 页文字小于 {MIN_FONT_PT:g} pt：{run.text!r}"
                        )


def main() -> int:
    parser = argparse.ArgumentParser()
    root = Path(__file__).resolve().parents[2]
    parser.add_argument(
        "--base",
        type=Path,
        default=Path.home() / "Downloads/mygo-defense-full.pptx",
    )
    parser.add_argument(
        "--topic-output",
        type=Path,
        default=root / "output/presentations/mygo-defense-elm-network-6pages.pptx",
    )
    parser.add_argument(
        "--full-output",
        type=Path,
        default=root / "output/presentations/mygo-defense-full.pptx",
    )
    args = parser.parse_args()
    base = args.base.resolve()
    topic = args.topic_output.resolve()
    full = args.full_output.resolve()
    if base == full:
        raise RuntimeError("--base 必须指向独立模板，不能与 --full-output 为同一文件")
    build_topic(base, topic)
    validate_topic(topic)
    build_full(base, full)
    print(topic)
    print(full)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
