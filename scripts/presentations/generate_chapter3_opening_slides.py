#!/usr/bin/env python3
"""生成答辩第三章前四页正文。"""

from __future__ import annotations

from copy import deepcopy
from io import BytesIO
from pathlib import Path
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from generate_engineering_structure_slide import (
    BODY,
    BLUE,
    INK,
    LINE,
    MUTED,
    NAVY,
    PALE_BLUE,
    PALE_GRAY,
    PALE_PURPLE,
    PALE_TEAL,
    PURPLE,
    TEAL,
    WHITE,
    add_arrow_tip,
    add_label,
    add_line,
    add_rect,
    add_text,
    draw_engineering_structure,
    rgb,
    set_title,
)


def clone_slide(prs: Presentation, source):
    """复制正文模板页，不共享可编辑图形。"""
    destination = prs.slides.add_slide(source.slide_layout)
    for shape in list(destination.shapes):
        element = shape._element
        element.getparent().remove(element)

    destination.background.fill.solid()
    destination.background.fill.fore_color.rgb = rgb("FAFBFC")

    for shape in source.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture = destination.shapes.add_picture(
                BytesIO(shape.image.blob),
                shape.left,
                shape.top,
                shape.width,
                shape.height,
            )
            picture.crop_left = shape.crop_left
            picture.crop_right = shape.crop_right
            picture.crop_top = shape.crop_top
            picture.crop_bottom = shape.crop_bottom
            picture.rotation = shape.rotation
            continue
        destination.shapes._spTree.insert_element_before(
            deepcopy(shape._element), "p:extLst"
        )
    return destination


def keep_slide_indices(prs: Presentation, indices: list[int]) -> None:
    keep = set(indices)
    slide_ids = prs.slides._sldIdLst
    for index in reversed(range(len(slide_ids))):
        if index in keep:
            continue
        slide_id = slide_ids[index]
        prs.part.drop_rel(slide_id.rId)
        del slide_ids[index]


def add_stage_box(
    slide,
    x,
    y,
    w,
    h,
    number,
    title,
    detail,
    *,
    fill,
    accent,
    title_size=14.2,
    detail_size=13.0,
):
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.08, h, accent)
    add_text(
        slide,
        number,
        x + 0.20,
        y + 0.13,
        0.34,
        0.24,
        size=12.5,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        title,
        x + 0.62,
        y + 0.10,
        w - 0.80,
        0.30,
        size=title_size,
        color=INK,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        detail,
        x + 0.22,
        y + 0.50,
        w - 0.42,
        h - 0.60,
        size=detail_size,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_flow_arrow(slide, x1, y1, x2, y2, color=BLUE, direction="right"):
    add_line(slide, x1, y1, x2, y2, color, 1.55)
    add_arrow_tip(slide, x2, y2, direction, color, 0.105)


def add_program_stage(
    slide,
    x,
    y,
    w,
    number,
    title,
    detail,
    *,
    fill,
    accent,
    title_size=13.2,
):
    """绘制固定网格的程序装载阶段，避免编号挤压长标题。"""
    add_rect(slide, x, y, w, 1.18, fill)
    add_rect(slide, x, y, 0.08, 1.18, accent)
    add_text(
        slide,
        number,
        x + 0.22,
        y + 0.12,
        0.34,
        0.22,
        size=12.5,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        title,
        x + 0.22,
        y + 0.37,
        w - 0.42,
        0.28,
        size=title_size,
        color=INK,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, x + 0.22, y + 0.72, x + w - 0.22, y + 0.72, LINE, 0.65)
    add_text(
        slide,
        detail,
        x + 0.22,
        y + 0.76,
        w - 0.42,
        0.38,
        size=12.2,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def draw_boot_flow(slide) -> None:
    set_title(slide, "启动控制流")

    # 两条架构入口在 StartContext 处收束。
    add_stage_box(
        slide,
        3.00,
        1.88,
        2.58,
        0.88,
        "A",
        "RISC-V64",
        "OpenSBI · DTB\n复制并固化固件视图",
        fill=PALE_BLUE,
        accent=BLUE,
        title_size=16.2,
        detail_size=12.8,
    )
    add_stage_box(
        slide,
        3.00,
        2.98,
        2.58,
        0.88,
        "B",
        "LoongArch64",
        "EFI · ACPI / DTB\n规范化 EFI 内存图",
        fill=PALE_TEAL,
        accent=TEAL,
        title_size=15.2,
        detail_size=12.8,
    )
    add_flow_arrow(slide, 5.58, 2.32, 5.95, 2.32, BLUE)
    add_flow_arrow(slide, 5.58, 3.42, 5.95, 3.42, TEAL)

    add_rect(slide, 5.95, 1.88, 3.02, 1.98, NAVY)
    add_text(
        slide,
        "StartContext",
        6.25,
        2.04,
        2.20,
        0.34,
        size=20.5,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "稳定、面向数据的启动交接对象",
        6.25,
        2.42,
        2.42,
        0.28,
        size=13.3,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_line(slide, 6.25, 2.83, 8.63, 2.83, "4E769D", 0.8)
    add_text(
        slide,
        "启动协议 · 固件快照 · 内存图",
        6.25,
        2.92,
        2.42,
        0.24,
        size=13.0,
        color="D9E7F1",
        bold=True,
    )
    add_text(
        slide,
        "地址转换 · 分配器与页表回调",
        6.25,
        3.24,
        2.42,
        0.24,
        size=13.0,
        color="D9E7F1",
        bold=True,
    )

    add_flow_arrow(slide, 8.97, 2.87, 9.34, 2.87, PURPLE)
    add_rect(slide, 9.34, 2.22, 3.21, 1.30, PALE_PURPLE)
    add_rect(slide, 9.34, 2.22, 0.08, 1.30, PURPLE)
    add_text(
        slide,
        "__kernel_start_init",
        9.64,
        2.38,
        2.42,
        0.28,
        size=15.5,
        color=PURPLE,
        bold=True,
    )
    add_text(
        slide,
        "校验上下文 · 选择 DTB / ACPI",
        9.64,
        2.80,
        2.50,
        0.30,
        size=13.2,
        color=INK,
        bold=True,
        chinese_font="SimHei",
    )
    add_label(slide, "进入通用主线", 10.64, 3.14, 1.28, PURPLE)

    # 将架构入口明确接入通用启动阶段，避免控制流在视觉上断开。
    add_line(slide, 10.94, 3.52, 10.94, 3.94, PURPLE, 1.55)
    add_line(slide, 10.94, 3.94, 4.40, 3.94, PURPLE, 1.55)
    add_line(slide, 4.40, 3.94, 4.40, 4.18, PURPLE, 1.55)
    add_arrow_tip(slide, 4.40, 4.14, "down", PURPLE, 0.11)

    # 通用启动阶段采用蛇形编号，保持控制流连续。
    xs = (3.00, 6.25, 9.50)
    top_y = 4.18
    bottom_y = 5.57
    stages = [
        (xs[0], top_y, "01", "校验与解析", "固件表 · 命令行 · CPU 拓扑", PALE_BLUE, BLUE),
        (xs[1], top_y, "02", "建立内存", "physical · paging · vmem · heap", PALE_BLUE, BLUE),
        (xs[2], top_y, "03", "设备前置", "filesystem · devtmpfs · devices", PALE_TEAL, TEAL),
        (xs[2], bottom_y, "04", "平台就绪", "PnP · PCI · rootfs · console", PALE_TEAL, TEAL),
        (xs[1], bottom_y, "05", "运行环境", "init task · syscalls · SMP", PALE_PURPLE, PURPLE),
        (xs[0], bottom_y, "06", "运行时与用户入口", "elm-mgr · network · /init", PALE_PURPLE, PURPLE),
    ]
    for x, y, number, title, detail, fill, accent in stages:
        add_stage_box(
            slide,
            x,
            y,
            2.80,
            0.92,
            number,
            title,
            detail,
            fill=fill,
            accent=accent,
            detail_size=13.0,
        )

    add_flow_arrow(slide, 5.80, 4.64, 6.17, 4.64, BLUE)
    add_flow_arrow(slide, 9.05, 4.64, 9.42, 4.64, TEAL)
    add_flow_arrow(slide, 10.90, 5.10, 10.90, 5.49, TEAL, "down")
    add_flow_arrow(slide, 9.50, 6.03, 9.13, 6.03, PURPLE, "left")
    add_flow_arrow(slide, 6.25, 6.03, 5.88, 6.03, PURPLE, "left")


def add_capability_card(slide, x, y, title, lines, *, fill, accent) -> None:
    """绘制基础能力卡片；每项都对应仓库中的实际子系统。"""
    add_rect(slide, x, y, 3.07, 1.72, fill)
    add_rect(slide, x, y, 3.07, 0.46, accent)
    add_text(
        slide,
        title,
        x + 0.22,
        y + 0.09,
        2.63,
        0.28,
        size=16.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    step = 0.28 if len(lines) == 4 else 0.34
    start = y + (0.55 if len(lines) == 4 else 0.61)
    for index, line in enumerate(lines):
        add_text(
            slide,
            line,
            x + 0.22,
            start + index * step,
            2.63,
            0.25,
            size=12.9 if len(lines) == 4 else 13.2,
            color=BODY,
            bold=True,
        )


def draw_capability_overview(slide) -> None:
    set_title(slide, "基础能力全景")

    add_rect(slide, 3.00, 1.88, 9.55, 0.58, NAVY)
    add_text(
        slide,
        "通用宏内核能力",
        3.28,
        2.01,
        2.10,
        0.30,
        size=17.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "进程、内存、I/O、通信与平台机制共同承载完整用户态",
        6.35,
        2.02,
        5.80,
        0.28,
        size=14.0,
        color="D9E7F1",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    cards = [
        (
            3.00,
            2.68,
            "进程与程序",
            [
                "进程 / 线程 / 进程组 / 会话",
                "fork · clone · exec · wait · exit",
                "ELF · shebang · auxv · VDSO",
                "credentials · rlimit · rseq · pidfd",
            ],
            PALE_BLUE,
            BLUE,
        ),
        (
            6.24,
            2.68,
            "调度与同步",
            [
                "EEVDF · per-CPU runqueue · SMP",
                "亲和性 · 迁移 · 时钟与抢占",
                "spinlock · mutex · futex · waitqueue",
            ],
            PALE_TEAL,
            TEAL,
        ),
        (
            9.48,
            2.68,
            "内存管理",
            [
                "物理页 · buddy / slab · vmem",
                "VMA · 缺页 · mmap / brk / mprotect",
                "COW · 文件映射 · ASID / TLB",
            ],
            PALE_PURPLE,
            PURPLE,
        ),
        (
            3.00,
            4.62,
            "VFS 与存储",
            [
                "FdTable · File · Dentry · Inode",
                "路径 · 挂载 · 权限 · 锁 · xattr",
                "extfs / fatfs / tmpfs",
                "procfs / sysfs / devtmpfs",
            ],
            PALE_BLUE,
            BLUE,
        ),
        (
            6.24,
            4.62,
            "IPC 与事件",
            [
                "signal · pipe · socketpair",
                "poll · select · epoll · eventfd",
                "timerfd · signalfd · SysV shm / sem",
            ],
            PALE_TEAL,
            TEAL,
        ),
        (
            9.48,
            4.62,
            "设备、网络与系统服务",
            [
                "ACPI / DTB · PCI · IRQ · DMA",
                "char / block / TTY · devtmpfs",
                "IPv4 / IPv6 · TCP / UDP · Unix socket",
                "console · clock · random · syslog",
            ],
            PALE_PURPLE,
            PURPLE,
        ),
    ]
    for x, y, title, lines, fill, accent in cards:
        add_capability_card(slide, x, y, title, lines, fill=fill, accent=accent)


def add_vfs_layer(slide, y, h, label, summary, detail, *, fill, accent) -> None:
    add_rect(slide, 3.00, y, 9.55, h, fill)
    add_rect(slide, 3.00, y, 1.52, h, accent)
    add_text(
        slide,
        label,
        3.12,
        y,
        1.28,
        h,
        size=15.2,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        summary,
        4.78,
        y + 0.10,
        7.43,
        0.28,
        size=14.2,
        color=INK,
        bold=True,
    )
    add_text(
        slide,
        detail,
        4.78,
        y + 0.43,
        7.43,
        0.26,
        size=13.2,
        color=BODY,
        bold=True,
    )


def draw_vfs_posix(slide) -> None:
    set_title(slide, "VFS 与 POSIX 兼容层")

    add_rect(slide, 3.00, 1.88, 9.55, 0.58, NAVY)
    add_text(
        slide,
        "双重职责",
        3.28,
        2.01,
        1.35,
        0.30,
        size=17.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "向上收口 Linux / POSIX 语义，向下统一文件系统与内核对象",
        5.05,
        2.02,
        7.10,
        0.28,
        size=14.0,
        color="D9E7F1",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    layers = [
        (
            2.70,
            0.80,
            "POSIX 接口",
            "路径与目录 · 文件与描述符 · 元数据与权限 · 事件等待 · socket / 设备文件",
            "openat / read / write / statx · dup / fcntl · poll / epoll · ioctl / mmap",
            PALE_BLUE,
            BLUE,
        ),
        (
            3.70,
            0.80,
            "ABI 投影",
            "Linux flags / errno / stat / dirent / timespec / ioctl",
            "用户布局在边界转换为 typed request，POSIX 常量不下沉到后端",
            PALE_PURPLE,
            PURPLE,
        ),
        (
            4.70,
            0.80,
            "VFS 语义",
            "VfsContext · FdTable / File · Path / Dentry / Inode · Mount / Superblock",
            "root / cwd · mount namespace · credentials / umask · locks / xattr · readiness",
            PALE_TEAL,
            TEAL,
        ),
        (
            5.70,
            0.96,
            "统一后端",
            "extfs · fatfs · tmpfs | procfs · sysfs · devtmpfs | char / block device",
            "pipe · socket · memfd · eventfd / timerfd / signalfd | file-backed mmap",
            PALE_GRAY,
            "607989",
        ),
    ]
    for y, h, label, summary, detail, fill, accent in layers:
        add_vfs_layer(slide, y, h, label, summary, detail, fill=fill, accent=accent)

    for y, color in ((3.60, PURPLE), (4.60, TEAL), (5.60, "607989")):
        add_line(slide, 3.76, y - 0.10, 3.76, y + 0.10, color, 1.45)
        add_arrow_tip(slide, 3.76, y + 0.06, "down", color, 0.10)


def add_titled_panel(slide, x, y, w, h, title, lines, *, fill, accent) -> None:
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, w, 0.50, accent)
    add_text(
        slide,
        title,
        x + 0.22,
        y + 0.10,
        w - 0.44,
        0.29,
        size=16.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    available = h - 0.64
    step = available / max(1, len(lines))
    for index, line in enumerate(lines):
        add_text(
            slide,
            line,
            x + 0.22,
            y + 0.57 + index * step,
            w - 0.44,
            min(0.30, step),
            size=13.0,
            color=BODY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def add_compact_stage(
    slide,
    x,
    y,
    w,
    h,
    title,
    detail,
    *,
    fill,
    accent,
    title_size=14.2,
    detail_size=12.5,
) -> None:
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, w, 0.07, accent)
    add_text(
        slide,
        title,
        x + 0.14,
        y + 0.15,
        w - 0.28,
        0.28,
        size=title_size,
        color=INK,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, x + 0.18, y + 0.50, x + w - 0.18, y + 0.50, LINE, 0.65)
    add_text(
        slide,
        detail,
        x + 0.14,
        y + 0.58,
        w - 0.28,
        h - 0.66,
        size=detail_size,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def draw_filesystem_storage(slide) -> None:
    set_title(slide, "文件系统与存储路径")

    add_rect(slide, 3.00, 1.88, 9.55, 0.58, NAVY)
    add_text(
        slide,
        "统一挂载模型",
        3.28,
        2.01,
        1.85,
        0.30,
        size=17.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "FsRegistry · FsDriver · Superblock · Inode / FileOps",
        6.10,
        2.02,
        6.05,
        0.28,
        size=14.2,
        color="D9E7F1",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    add_titled_panel(
        slide,
        3.00,
        2.70,
        2.92,
        2.08,
        "内存与伪文件系统",
        [
            "tmpfs：内存文件与 /dev/shm",
            "procfs / sysfs：内核状态视图",
            "devtmpfs：设备能力的用户态投影",
        ],
        fill=PALE_BLUE,
        accent=BLUE,
    )
    add_titled_panel(
        slide,
        6.31,
        2.70,
        2.92,
        2.08,
        "磁盘文件系统",
        [
            "ext2 / ext3 / ext4：extent 与写回",
            "JBD2 / fast commit / orphan recovery",
            "FAT：目录、长文件名与读写同步",
        ],
        fill=PALE_TEAL,
        accent=TEAL,
    )
    add_titled_panel(
        slide,
        9.63,
        2.70,
        2.92,
        2.08,
        "块设备后端",
        [
            "BlockFsAdapter / BlockBackend",
            "BIO：范围、方向与完成状态",
            "BlockFunction → virtio-blk",
        ],
        fill=PALE_PURPLE,
        accent=PURPLE,
    )

    add_label(slide, "读写主路径", 3.00, 5.04, 1.08, NAVY)
    stages = [
        ("系统调用", "read / write\nfsync", PALE_BLUE, BLUE),
        ("VFS", "FileOps\nInodeOps", PALE_BLUE, BLUE),
        ("文件系统", "extent / FAT\nwriteback", PALE_TEAL, TEAL),
        ("块 I/O", "BlockBackend\nBIO", PALE_TEAL, TEAL),
        ("设备能力", "BlockFunction\nvirtio-blk", PALE_PURPLE, PURPLE),
    ]
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(stages):
        add_compact_stage(slide, x, 5.48, 1.65, 1.18, title, detail, fill=fill, accent=accent)
        if index != len(stages) - 1:
            add_flow_arrow(slide, x + 1.65, 6.07, x + 1.93, 6.07, accent)
        x += 1.975


def draw_virtual_memory(slide) -> None:
    set_title(slide, "虚拟内存、缺页与 COW")

    add_rect(slide, 3.00, 1.88, 9.55, 0.58, NAVY)
    add_text(
        slide,
        "VmSpace 策略层",
        3.28,
        2.01,
        1.90,
        0.30,
        size=17.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "VmaSet + resident ledger + PgdHandle；页表机械操作由 arch 注入",
        5.55,
        2.02,
        6.60,
        0.28,
        size=14.0,
        color="D9E7F1",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    flow = [
        ("Trap 解码", "FaultDecodeOps\n地址 · 访问类型", PALE_BLUE, BLUE),
        ("通用分派", "dispatch_page_fault\ncurrent VmSpace", PALE_BLUE, BLUE),
        ("VMA 检查", "范围 · 权限\nanon / file / shared", PALE_TEAL, TEAL),
        ("驻留解析", "lazy fault-in\nfile read · COW", PALE_TEAL, TEAL),
        ("页表提交", "UserPgdOps\nPTE · TLB / shootdown", PALE_PURPLE, PURPLE),
    ]
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(flow):
        add_compact_stage(
            slide,
            x,
            2.76,
            1.55,
            1.32,
            title,
            detail,
            fill=fill,
            accent=accent,
            detail_size=12.2,
        )
        if index != len(flow) - 1:
            add_flow_arrow(slide, x + 1.55, 3.42, x + 1.94, 3.42, accent)
        x += 2.00

    add_titled_panel(
        slide,
        3.00,
        4.42,
        4.55,
        2.28,
        "地址空间操作",
        [
            "mmap / brk：建立 VMA 与 backing",
            "munmap / mremap / mprotect：修改映射",
            "mlock / msync：驻留与文件写回语义",
            "copy / pin user：先完成 fault-in 与权限检查",
        ],
        fill=PALE_BLUE,
        accent=BLUE,
    )
    add_titled_panel(
        slide,
        8.00,
        4.42,
        4.55,
        2.28,
        "fork 与 private COW",
        [
            "01 fork：克隆 VMA，驻留页共享并降级写权限",
            "02 Store fault：区分 COW、共享写与权限错误",
            "03 解析：分配 / 复制 / 重映射，保留页所有权",
            "远端 TLB 完成失效后，旧页才允许释放",
        ],
        fill=PALE_PURPLE,
        accent=PURPLE,
    )


def draw_program_lifecycle(slide) -> None:
    set_title(slide, "程序装载与进程生命周期")

    add_label(slide, "execve 装载链", 3.00, 1.88, 1.24, NAVY)
    pipeline = [
        ("01", "execve / VFS", "权限 · ETXTBSY\nshebang", PALE_BLUE, BLUE),
        ("02", "ProcessImageOps", "调度状态机与\nloader 解耦", PALE_BLUE, BLUE),
        ("03", "ELF / interpreter", "PT_LOAD · PIE\n动态解释器", PALE_TEAL, TEAL),
        ("04", "VmSpace / Context", "segments · stack\nauxv · PC / SP", PALE_PURPLE, PURPLE),
    ]
    x = 3.00
    for index, (number, title, detail, fill, accent) in enumerate(pipeline):
        add_program_stage(
            slide,
            x,
            2.30,
            2.22,
            number,
            title,
            detail,
            fill=fill,
            accent=accent,
            title_size=12.2 if index else 13.2,
        )
        if index != len(pipeline) - 1:
            add_flow_arrow(slide, x + 2.22, 2.89, x + 2.40, 2.89, accent)
        x += 2.44

    add_line(slide, 11.44, 3.48, 11.44, 3.72, PURPLE, 1.35)
    add_arrow_tip(slide, 11.44, 3.68, "down", PURPLE, 0.10)
    add_rect(slide, 3.00, 3.72, 9.55, 0.48, PALE_GRAY)
    add_text(
        slide,
        "Task / ThreadGroup 生命周期",
        3.28,
        3.84,
        2.60,
        0.23,
        size=14.0,
        color="465C6C",
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "clone flags 明确资源共享；创建、替换与回收分别提交",
        6.05,
        3.84,
        6.10,
        0.23,
        size=13.2,
        color=BODY,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    add_titled_panel(
        slide,
        3.00,
        4.46,
        2.95,
        2.24,
        "创建",
        [
            "clone / clone3 / fork",
            "复制或共享 VM、FdTable、信号表",
            "TLS · TID · pidfd · thread group",
        ],
        fill=PALE_BLUE,
        accent=BLUE,
    )
    add_titled_panel(
        slide,
        6.30,
        4.46,
        2.95,
        2.24,
        "替换",
        [
            "execve / execveat",
            "替换 VmSpace 与完整用户上下文",
            "close-on-exec · signal / rseq reset",
        ],
        fill=PALE_TEAL,
        accent=TEAL,
    )
    add_titled_panel(
        slide,
        9.60,
        4.46,
        2.95,
        2.24,
        "退出与回收",
        [
            "exit / exit_group → Zombie",
            "SIGCHLD · robust cleanup",
            "wait4 / waitid 完成最终回收",
        ],
        fill=PALE_PURPLE,
        accent=PURPLE,
    )


def add_mechanism_panel(slide, x, title, rows, *, fill, accent) -> None:
    add_rect(slide, x, 2.82, 2.95, 3.88, fill)
    add_rect(slide, x, 2.82, 2.95, 0.56, accent)
    add_text(
        slide,
        title,
        x + 0.24,
        2.95,
        2.47,
        0.30,
        size=16.5,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    row_y = 3.52
    for index, (heading, detail) in enumerate(rows):
        add_text(
            slide,
            heading,
            x + 0.22,
            row_y,
            2.51,
            0.25,
            size=13.8,
            color=INK,
            bold=True,
            chinese_font="SimHei",
        )
        add_text(
            slide,
            detail,
            x + 0.22,
            row_y + 0.25,
            2.51,
            0.25,
            size=12.5,
            color=BODY,
            bold=True,
            align=PP_ALIGN.LEFT,
        )
        if index != len(rows) - 1:
            add_line(slide, x + 0.22, row_y + 0.54, x + 2.73, row_y + 0.54, LINE, 0.65)
        row_y += 0.62


def draw_sched_sync_signal(slide) -> None:
    set_title(slide, "调度、同步与信号")

    add_rect(slide, 3.00, 1.88, 9.55, 0.62, NAVY)
    add_text(
        slide,
        "统一 Task 状态机",
        3.28,
        2.02,
        2.15,
        0.30,
        size=17.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "Runnable · Running · Sleeping · Stopped · Zombie",
        6.15,
        2.03,
        6.00,
        0.28,
        size=14.5,
        color="D9E7F1",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    add_mechanism_panel(
        slide,
        3.00,
        "调度",
        [
            ("调度类", "DL > RT > Fair > Idle"),
            ("公平队列", "EEVDF eligible · deadline"),
            ("SMP 队列", "per-CPU rq · affinity"),
            ("迁移均衡", "placement · migration · steal"),
            ("抢占边界", "tick / wakeup 设置 resched"),
        ],
        fill=PALE_BLUE,
        accent=BLUE,
    )
    add_mechanism_panel(
        slide,
        6.30,
        "同步",
        [
            ("短临界区", "spinlock · acquire / release"),
            ("睡眠等待", "mutex · WaitQueue"),
            ("用户同步", "futex wait / wake / requeue"),
            ("优先级继承", "PI owner-chain propagation"),
            ("内存顺序", "只在跨 CPU 发布边界建立"),
        ],
        fill=PALE_TEAL,
        accent=TEAL,
    )
    add_mechanism_panel(
        slide,
        9.60,
        "信号",
        [
            ("信号状态", "per-task + group pending"),
            ("处理策略", "sigaction · blocked mask"),
            ("返回快路", "user-return work flag"),
            ("投递慢路", "select · signal frame"),
            ("恢复执行", "rt_sigreturn · restart"),
        ],
        fill=PALE_PURPLE,
        accent=PURPLE,
    )


def draw_device_abstraction_overview(slide) -> None:
    set_title(slide, "设备抽象概览")

    # 先给出不依赖代码的三段式心智模型。
    add_rect(slide, 3.00, 1.88, 2.30, 0.68, PALE_BLUE)
    add_rect(slide, 3.00, 1.88, 0.08, 0.68, BLUE)
    add_text(
        slide,
        "PnPDevice",
        3.26,
        1.97,
        1.76,
        0.27,
        size=16.0,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "物理身份 · 总线 · 资源",
        3.26,
        2.25,
        1.76,
        0.22,
        size=12.8,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_flow_arrow(slide, 5.30, 2.22, 5.70, 2.22, BLUE)

    add_rect(slide, 5.70, 1.88, 3.25, 0.68, NAVY)
    add_text(
        slide,
        "DeviceFunction × N",
        5.95,
        1.97,
        2.75,
        0.27,
        size=16.8,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "同一设备可以同时发布多种能力",
        5.95,
        2.25,
        2.75,
        0.22,
        size=13.0,
        color="D9E7F1",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_flow_arrow(slide, 8.95, 2.22, 9.35, 2.22, PURPLE)

    add_rect(slide, 9.35, 1.88, 3.20, 0.68, PALE_PURPLE)
    add_rect(slide, 9.35, 1.88, 0.08, 0.68, PURPLE)
    add_text(
        slide,
        "用户接口 / 内核服务",
        9.61,
        1.97,
        2.66,
        0.27,
        size=15.2,
        color=PURPLE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "/dev · sysfs · kernel / ELM consumer",
        9.61,
        2.25,
        2.66,
        0.22,
        size=12.3,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 左侧剖开一份 Function，解释它不是一个空泛的 trait 名称。
    add_rect(slide, 3.00, 2.84, 5.84, 3.28, PALE_BLUE)
    add_rect(slide, 3.00, 2.84, 5.84, 0.54, BLUE)
    add_text(
        slide,
        "DeviceFunction = 一份可管理的能力契约",
        3.28,
        2.96,
        5.28,
        0.30,
        size=16.5,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    contract_rows = [
        ("类别身份", "class_id + dev_name：稳定、唯一、可查询"),
        ("能力接口", "operation contract：明确它能够完成什么工作"),
        ("调用方式", "类型化对象直接调用，也可通过 opcode 使用通用契约"),
        ("安全退役", "mark_gone → drain_io：先拒绝新访问，再排空 I/O"),
    ]
    row_y = 3.56
    for index, (heading, detail) in enumerate(contract_rows):
        add_rect(slide, 3.28, row_y, 1.10, 0.40, WHITE)
        add_text(
            slide,
            heading,
            3.28,
            row_y,
            1.10,
            0.40,
            size=14.0,
            color=BLUE,
            bold=True,
            chinese_font="SimHei",
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            detail,
            4.62,
            row_y,
            3.82,
            0.40,
            size=13.2,
            color=INK if index < 2 else BODY,
            bold=True,
        )
        row_y += 0.60

    # 右侧直接回答“为什么可拓展”，用新增第 N+1 种能力的过程表达。
    add_rect(slide, 9.14, 2.84, 3.41, 3.28, PALE_PURPLE)
    add_rect(slide, 9.14, 2.84, 3.41, 0.54, PURPLE)
    add_text(
        slide,
        "新增第 N + 1 种能力",
        9.38,
        2.96,
        2.93,
        0.30,
        size=16.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    extension_rows = [
        ("01", "注册新类别", "编号单调分配，不与旧句柄复用"),
        ("02", "实现能力契约", "核心只保存统一 DeviceFunction"),
        ("03", "选择投影方式", "/dev、sysfs 或仅供内核消费"),
    ]
    row_y = 3.55
    for number, heading, detail in extension_rows:
        add_text(
            slide,
            number,
            9.40,
            row_y,
            0.34,
            0.24,
            size=12.5,
            color=PURPLE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            heading,
            9.84,
            row_y - 0.02,
            2.25,
            0.27,
            size=14.2,
            color=INK,
            bold=True,
            chinese_font="SimHei",
        )
        add_text(
            slide,
            detail,
            9.84,
            row_y + 0.29,
            2.25,
            0.25,
            size=12.7,
            color=BODY,
            bold=True,
        )
        row_y += 0.72
    add_rect(slide, 9.40, 5.77, 2.89, 0.22, WHITE)
    add_text(
        slide,
        "PnP core 无需增加类型分支",
        9.40,
        5.76,
        2.89,
        0.25,
        size=13.0,
        color=PURPLE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )

    add_rect(slide, 3.00, 6.34, 9.55, 0.36, PALE_GRAY)
    add_rect(slide, 3.00, 6.34, 0.09, 0.36, "607989")
    add_text(
        slide,
        "字符、块、RTC 只是已有实例；NPU、计算卡等能力可沿同一规则接入，设备文件只是可选投影。",
        3.28,
        6.34,
        8.98,
        0.36,
        size=13.4,
        color="465C6C",
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )


def draw_device_function_capability(slide) -> None:
    set_title(slide, "DeviceFunction 承载真实能力")

    add_rect(slide, 3.00, 1.88, 9.55, 0.58, NAVY)
    add_text(
        slide,
        "以 BlockFunction 为例",
        3.28,
        2.01,
        2.30,
        0.30,
        size=17.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "一次注册同时带入可调用对象、I/O 状态、完成通知与退役语义",
        6.00,
        2.02,
        6.15,
        0.28,
        size=14.0,
        color="D9E7F1",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    stages = [
        ("消费方", "/dev/vd*\nVFS / extfs", PALE_BLUE, BLUE),
        ("BlockFunction", "class + name\nArc<BlockDevice>", PALE_BLUE, BLUE),
        ("BlockDevice", "geometry · limits\nstate · statistics", PALE_TEAL, TEAL),
        ("BIO 队列", "Read / Write\nFlush / Discard", PALE_TEAL, TEAL),
        ("完成语义", "IRQ / poll\ncomplete · wake", PALE_PURPLE, PURPLE),
    ]
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(stages):
        add_compact_stage(
            slide,
            x,
            2.82,
            1.65,
            1.28,
            title,
            detail,
            fill=fill,
            accent=accent,
            title_size=13.6 if index in (1, 2) else 14.2,
            detail_size=12.4,
        )
        if index != len(stages) - 1:
            add_flow_arrow(slide, x + 1.65, 3.46, x + 1.93, 3.46, accent)
        x += 1.975

    add_titled_panel(
        slide,
        3.00,
        4.44,
        4.55,
        2.26,
        "同步与异步共享同一驱动路径",
        [
            "驱动只实现异步 queue_bio，将请求放入硬件队列",
            "同步调用等待 Completion；异步调用返回 Future",
            "IRQ 或主动 drain 处理 used ring 并调用 bio.complete",
            "完成结果唤醒等待任务，同时更新 I/O 统计",
        ],
        fill=PALE_TEAL,
        accent=TEAL,
    )
    add_titled_panel(
        slide,
        8.00,
        4.44,
        4.55,
        2.26,
        "退役不是简单释放引用",
        [
            "Active / Gone 状态首先拒绝新的 open 与提交",
            "mark_gone 让全部旧句柄立即看到设备失效",
            "drain_io 排空已经被硬件接受的请求",
            "完成静默后才撤销注册表与用户态投影",
        ],
        fill=PALE_PURPLE,
        accent=PURPLE,
    )


def add_numbered_process_panel(slide, x, title, rows, *, fill, accent) -> None:
    add_rect(slide, x, 2.98, 4.55, 3.72, fill)
    add_rect(slide, x, 2.98, 4.55, 0.54, accent)
    add_text(
        slide,
        title,
        x + 0.26,
        3.10,
        4.03,
        0.30,
        size=16.2,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    available = 3.03
    step = available / len(rows)
    row_y = 3.66
    for index, (heading, detail) in enumerate(rows, 1):
        add_text(
            slide,
            f"{index:02d}",
            x + 0.24,
            row_y,
            0.38,
            0.24,
            size=12.5,
            color=accent,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            heading,
            x + 0.78,
            row_y - 0.02,
            3.45,
            0.27,
            size=14.0,
            color=INK,
            bold=True,
            chinese_font="SimHei",
        )
        add_text(
            slide,
            detail,
            x + 0.78,
            row_y + 0.27,
            3.45,
            0.24,
            size=12.4,
            color=BODY,
            bold=True,
            align=PP_ALIGN.LEFT,
        )
        if index != len(rows):
            add_line(slide, x + 0.24, row_y + step - 0.04, x + 4.31, row_y + step - 0.04, LINE, 0.65)
        row_y += step


def draw_device_lifecycle(slide) -> None:
    set_title(slide, "生命周期与故障边界")

    states = [
        ("Discovered", PALE_BLUE, BLUE),
        ("Probing", PALE_BLUE, BLUE),
        ("Bound", PALE_TEAL, TEAL),
        ("Removing", PALE_PURPLE, PURPLE),
        ("Gone", PALE_GRAY, "607989"),
    ]
    x = 3.00
    for index, (state, fill, accent) in enumerate(states):
        add_rect(slide, x, 1.88, 1.55, 0.62, fill)
        add_rect(slide, x, 1.88, 0.07, 0.62, accent)
        add_text(
            slide,
            state,
            x + 0.14,
            2.01,
            1.27,
            0.30,
            size=14.2,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if index != len(states) - 1:
            add_flow_arrow(slide, x + 1.55, 2.19, x + 1.93, 2.19, accent)
        x += 2.00
    add_line(slide, 5.76, 2.50, 5.76, 2.72, BLUE, 1.2)
    add_line(slide, 5.76, 2.72, 3.76, 2.72, BLUE, 1.2)
    add_arrow_tip(slide, 3.80, 2.72, "left", BLUE, 0.09)
    add_text(
        slide,
        "probe 失败：逆序撤销副作用，回到 Discovered",
        3.92,
        2.53,
        3.15,
        0.22,
        size=12.5,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_numbered_process_panel(
        slide,
        3.00,
        "probe 是一次提交事务",
        [
            ("识别与进入", "匹配总线身份，进入 Probing"),
            ("建立副作用", "登记子设备、Function 与 IRQ/MSI/DMA 资源"),
            ("原子发布", "设备内附着后再进入全局 registry，随后提交 Bound"),
            ("失败回滚", "叶子优先移除，gone + drain，逆序释放资源"),
        ],
        fill=PALE_BLUE,
        accent=BLUE,
    )
    add_numbered_process_panel(
        slide,
        8.00,
        "remove 先静默再撤销可见性",
        [
            ("隔离设备", "进入 Removing，阻止新的 probe"),
            ("清理拓扑", "子设备按叶子优先递归移除"),
            ("停止数据流", "Function mark_gone，再 drain_io"),
            ("关闭硬件", "driver.remove，释放 PnP-owned resources"),
            ("撤销投影", "发布 Unregistered，解绑 /dev，最终进入 Gone"),
        ],
        fill=PALE_PURPLE,
        accent=PURPLE,
    )


def draw_virtio_block_capability_loop(slide) -> None:
    set_title(slide, "VirtIO Block：能力闭环")

    add_label(slide, "真实执行路径", 3.00, 1.88, 1.18, NAVY)
    xs = (3.00, 6.25, 9.50)
    top_y = 2.24
    bottom_y = 4.02
    stages = [
        (xs[0], top_y, "01", "发现设备", "PCI / MMIO bus\nPnpDevice + resources", PALE_BLUE, BLUE),
        (xs[1], top_y, "02", "驱动 probe", "match IDs · negotiate\nqueue + IRQ ownership", PALE_BLUE, BLUE),
        (xs[2], top_y, "03", "发布能力", "BlockDevice\nBlockFunction(vd*)", PALE_TEAL, TEAL),
        (xs[2], bottom_y, "04", "创建投影", "Registered event\nprojector → /dev/vd*", PALE_TEAL, TEAL),
        (xs[1], bottom_y, "05", "提交请求", "VFS / extfs → BIO\nqueue_bio → virtqueue", PALE_PURPLE, PURPLE),
        (xs[0], bottom_y, "06", "完成请求", "IRQ / poll → drain\nbio.complete → wake", PALE_PURPLE, PURPLE),
    ]
    for x, y, number, title, detail, fill, accent in stages:
        add_stage_box(
            slide,
            x,
            y,
            2.80,
            1.10,
            number,
            title,
            detail,
            fill=fill,
            accent=accent,
            detail_size=12.8,
        )
    add_flow_arrow(slide, 5.80, 2.79, 6.17, 2.79, BLUE)
    add_flow_arrow(slide, 9.05, 2.79, 9.42, 2.79, TEAL)
    add_flow_arrow(slide, 10.90, 3.34, 10.90, 3.94, TEAL, "down")
    add_flow_arrow(slide, 9.50, 4.57, 9.13, 4.57, PURPLE, "left")
    add_flow_arrow(slide, 6.25, 4.57, 5.88, 4.57, PURPLE, "left")

    add_rect(slide, 3.00, 5.42, 5.95, 0.92, PALE_GRAY)
    add_rect(slide, 3.00, 5.42, 0.09, 0.92, "607989")
    add_text(
        slide,
        "反向移除",
        3.28,
        5.51,
        1.18,
        0.28,
        size=15.0,
        color="465C6C",
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "mark_gone → drain → remove IRQ / reset → unbind /dev",
        3.28,
        5.87,
        5.22,
        0.27,
        size=13.2,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_rect(slide, 9.20, 5.42, 3.35, 0.92, PALE_PURPLE)
    add_rect(slide, 9.20, 5.42, 0.09, 0.92, PURPLE)
    add_text(
        slide,
        "实现来源可以变化",
        9.46,
        5.51,
        2.83,
        0.28,
        size=15.0,
        color=PURPLE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "集成组件或受管模块\n内核消费同一稳定能力契约",
        9.46,
        5.82,
        2.83,
        0.40,
        size=12.8,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "从发现到一次读请求完成，每一段都有实际执行者与可观测状态。",
        3.28,
        6.43,
        8.98,
        0.24,
        size=13.3,
        color=BODY,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )


def add_path_row(slide, x, y, w, number, title, detail, accent):
    add_text(
        slide,
        number,
        x,
        y,
        0.34,
        0.28,
        size=11.3,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        title,
        x + 0.50,
        y - 0.02,
        1.48,
        0.30,
        size=14.0,
        color=INK,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        detail,
        x + 1.95,
        y,
        w - 2.00,
        0.28,
        size=11.8,
        color=BODY,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def draw_data_paths(slide) -> None:
    set_title(slide, "用户态兼容与数据承载")

    add_rect(slide, 3.00, 1.88, 9.55, 0.70, NAVY)
    add_text(
        slide,
        "Linux / POSIX 用户态边界",
        3.30,
        2.04,
        2.65,
        0.30,
        size=17.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "系统调用号 · 标志 · errno · 用户结构只在入口转换",
        7.20,
        2.04,
        4.95,
        0.30,
        size=12.7,
        color="D9E7F1",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    add_line(slide, 5.28, 2.58, 5.28, 2.82, BLUE, 1.4)
    add_arrow_tip(slide, 5.28, 2.78, "down", BLUE, 0.09)
    add_line(slide, 10.28, 2.58, 10.28, 2.82, TEAL, 1.4)
    add_arrow_tip(slide, 10.28, 2.78, "down", TEAL, 0.09)

    # 文件 I/O 路径。
    add_rect(slide, 3.00, 2.82, 4.55, 3.10, PALE_BLUE)
    add_rect(slide, 3.00, 2.82, 4.55, 0.54, BLUE)
    add_text(
        slide,
        "文件与 I/O 路径",
        3.28,
        2.94,
        2.20,
        0.28,
        size=16.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    left_rows = [
        ("01", "兼容层", "open · read · write · stat"),
        ("02", "FdTable / VFS", "File · Dentry · Inode · Mount"),
        ("03", "文件系统驱动", "tmpfs · extfs · fatfs · devtmpfs"),
        ("04", "块设备能力", "BlockFunction · block I/O"),
    ]
    y = 3.58
    for index, (number, title, detail) in enumerate(left_rows):
        add_path_row(slide, 3.28, y, 3.98, number, title, detail, BLUE)
        if index != len(left_rows) - 1:
            add_line(slide, 3.28, y + 0.48, 7.25, y + 0.48, "CAD9E5", 0.7)
        y += 0.62

    # 虚拟内存与缺页路径。
    add_rect(slide, 8.00, 2.82, 4.55, 3.10, PALE_TEAL)
    add_rect(slide, 8.00, 2.82, 4.55, 0.54, TEAL)
    add_text(
        slide,
        "地址空间与缺页路径",
        8.28,
        2.94,
        2.48,
        0.28,
        size=16.0,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    right_rows = [
        ("01", "内存接口", "mmap · brk · mprotect · fork"),
        ("02", "VmSpace", "VMA set · resident page map"),
        ("03", "缺页分派", "anonymous · file-backed · shared"),
        ("04", "写时复制", "private COW · PTE / TLB update"),
    ]
    y = 3.58
    for index, (number, title, detail) in enumerate(right_rows):
        add_path_row(slide, 8.28, y, 3.98, number, title, detail, TEAL)
        if index != len(right_rows) - 1:
            add_line(slide, 8.28, y + 0.48, 12.25, y + 0.48, "C8DDDB", 0.7)
        y += 0.62

    # file-backed mmap 是两条数据路径的显式桥梁。
    add_flow_arrow(slide, 7.55, 4.84, 7.95, 4.84, PURPLE)
    add_text(
        slide,
        "mmap",
        7.55,
        4.57,
        0.40,
        0.20,
        size=10.8,
        color=PURPLE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_line(slide, 5.28, 5.92, 5.28, 6.12, BLUE, 1.25)
    add_arrow_tip(slide, 5.28, 6.08, "down", BLUE, 0.09)
    add_line(slide, 10.28, 5.92, 10.28, 6.12, TEAL, 1.25)
    add_arrow_tip(slide, 10.28, 6.08, "down", TEAL, 0.09)
    add_rect(slide, 3.00, 6.12, 9.55, 0.58, PALE_GRAY)
    add_rect(slide, 3.00, 6.12, 0.09, 0.58, "758998")
    add_text(
        slide,
        "统一数据后备",
        3.28,
        6.26,
        1.42,
        0.25,
        size=14.2,
        color="465C6C",
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "physical page · file page cache · block storage",
        5.05,
        6.26,
        4.16,
        0.25,
        size=12.3,
        color=INK,
        bold=True,
    )
    add_text(
        slide,
        "权限与所有权在内核语义中保持一致",
        9.22,
        6.26,
        2.92,
        0.25,
        size=12.2,
        color=BODY,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.RIGHT,
    )


def draw_program_concurrency(slide) -> None:
    set_title(slide, "程序执行与并发")

    add_label(slide, "程序装载链", 3.00, 1.88, 1.08, NAVY)
    pipeline = [
        ("01", "execve / VFS", "打开映像 · 兼容入口"),
        ("02", "ProcessImageOps", "隔离 sched 与 loader"),
        ("03", "ELF / interpreter", "PIE · shebang · PT_LOAD"),
        ("04", "VmSpace / Context", "segments · stack · auxv · PC / SP"),
    ]
    x = 3.00
    box_w = 2.22
    gap = 0.22
    for index, (number, title, detail) in enumerate(pipeline):
        fill = PALE_BLUE if index < 2 else PALE_TEAL if index == 2 else PALE_PURPLE
        accent = BLUE if index < 2 else TEAL if index == 2 else PURPLE
        add_program_stage(
            slide,
            x,
            2.32,
            box_w,
            number,
            title,
            detail,
            fill=fill,
            accent=accent,
            title_size=12.2 if index > 0 else 13.2,
        )
        if index != len(pipeline) - 1:
            add_flow_arrow(
                slide,
                x + box_w,
                2.91,
                x + box_w + gap - 0.04,
                2.91,
                accent,
            )
        x += box_w + gap

    # 装载结果进入统一任务生命周期，再分派到三类运行机制。
    add_line(slide, 11.44, 3.50, 11.44, 3.72, PURPLE, 1.2)
    add_arrow_tip(slide, 11.44, 3.68, "down", PURPLE, 0.085)
    add_line(slide, 7.78, 4.20, 7.78, 4.30, "7D8F9D", 1.0)
    add_line(slide, 4.48, 4.30, 11.08, 4.30, "7D8F9D", 1.0)
    for center, color in ((4.48, BLUE), (7.78, TEAL), (11.08, PURPLE)):
        add_line(slide, center, 4.30, center, 4.46, color, 1.2)
        add_arrow_tip(slide, center, 4.42, "down", color, 0.085)

    add_rect(slide, 3.00, 3.72, 9.55, 0.48, PALE_GRAY)
    add_text(
        slide,
        "进程与线程生命周期",
        3.28,
        3.84,
        1.78,
        0.23,
        size=13.4,
        color="465C6C",
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "clone / fork · thread group · wait / exit · robust cleanup",
        5.24,
        3.84,
        4.80,
        0.23,
        size=12.0,
        color=INK,
        bold=True,
    )
    add_text(
        slide,
        "运行边界保持 Linux ABI",
        10.30,
        3.84,
        1.84,
        0.23,
        size=11.8,
        color=BODY,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.RIGHT,
    )

    panels = [
        (
            3.00,
            "调度",
            "EEVDF fair class",
            "per-CPU runqueue",
            "SMP context switch",
            PALE_BLUE,
            BLUE,
        ),
        (
            6.30,
            "同步",
            "Spinlock：短临界区",
            "Mutex + WaitQueue：睡眠等待",
            "Futex / PI：用户态同步",
            PALE_TEAL,
            TEAL,
        ),
        (
            9.60,
            "信号",
            "task / thread-group pending",
            "return-to-user slow path",
            "signal frame / sigreturn",
            PALE_PURPLE,
            PURPLE,
        ),
    ]
    for x, title, first, second, third, fill, accent in panels:
        add_rect(slide, x, 4.46, 2.95, 2.24, fill)
        add_rect(slide, x, 4.46, 2.95, 0.56, accent)
        add_text(
            slide,
            title,
            x + 0.24,
            4.59,
            0.72,
            0.28,
            size=16.0,
            color=WHITE,
            bold=True,
            chinese_font="SimHei",
        )
        add_text(
            slide,
            first,
            x + 0.28,
            5.21,
            2.40,
            0.28,
            size=12.3,
            color=INK,
            bold=True,
            chinese_font="SimHei",
            align=PP_ALIGN.CENTER,
        )
        add_line(slide, x + 0.28, 5.56, x + 2.67, 5.56, LINE, 0.7)
        add_text(
            slide,
            second,
            x + 0.28,
            5.65,
            2.40,
            0.28,
            size=11.9,
            color=BODY,
            bold=True,
            chinese_font="SimHei",
            align=PP_ALIGN.CENTER,
        )
        add_line(slide, x + 0.28, 6.00, x + 2.67, 6.00, LINE, 0.7)
        add_text(
            slide,
            third,
            x + 0.28,
            6.09,
            2.40,
            0.28,
            size=11.9,
            color=BODY,
            bold=True,
            chinese_font="SimHei",
            align=PP_ALIGN.CENTER,
        )

def build_deck(template: Path, output: Path) -> None:
    prs = Presentation(template)
    source_index = 11
    source = prs.slides[source_index]
    pages = [source]
    pages.extend(clone_slide(prs, source) for _ in range(11))

    draw_engineering_structure(pages[0])
    draw_boot_flow(pages[1])
    draw_capability_overview(pages[2])
    draw_vfs_posix(pages[3])
    draw_filesystem_storage(pages[4])
    draw_virtual_memory(pages[5])
    draw_program_lifecycle(pages[6])
    draw_sched_sync_signal(pages[7])
    draw_device_abstraction_overview(pages[8])
    draw_device_function_capability(pages[9])
    draw_device_lifecycle(pages[10])
    draw_virtio_block_capability_loop(pages[11])

    clone_start = len(prs.slides) - 11
    keep_slide_indices(prs, [source_index, *range(clone_start, clone_start + 11)])
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def build_device_deck(template: Path, output: Path) -> None:
    prs = Presentation(template)
    source_index = 11
    source = prs.slides[source_index]
    pages = [source]
    pages.extend(clone_slide(prs, source) for _ in range(3))

    draw_device_abstraction_overview(pages[0])
    draw_device_function_capability(pages[1])
    draw_device_lifecycle(pages[2])
    draw_virtio_block_capability_loop(pages[3])

    clone_start = len(prs.slides) - 3
    keep_slide_indices(prs, [source_index, *range(clone_start, clone_start + 3)])
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> int:
    root = Path(__file__).resolve().parents[2]
    template = root / "output/presentations/mygo-defense.pptx"
    output = root / "output/presentations/mygo-defense-chapter3-pages1-12.pptx"
    if len(sys.argv) > 1:
        output = Path(sys.argv[1]).resolve()
    build_deck(template, output)
    print(output)
    if len(sys.argv) == 1:
        device_output = root / "output/presentations/mygo-defense-device-abstraction-pages9-12.pptx"
        build_device_deck(template, device_output)
        print(device_output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
