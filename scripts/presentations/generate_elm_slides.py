#!/usr/bin/env python3
"""生成第三章 ELM 专题 11 页答辩稿。

内容依据 docs/ELM.md、libs/elm 和 kernel/src/elm 的当前实现，使用最新全量稿的
第三章正文模板作为视觉底板。该脚本只生成独立文件，不改写全量答辩稿。
"""

from __future__ import annotations

from copy import deepcopy
from io import BytesIO
from pathlib import Path
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from generate_engineering_structure_slide import (
    BG,
    BODY,
    BLUE,
    INK,
    LINE,
    MUTED,
    NAVY,
    PALE_BLUE,
    PALE_GRAY,
    PALE_PURPLE,
    PALE_TEAL,
    PURPLE,
    TEAL,
    WHITE,
    add_arrow_tip,
    add_line,
    add_rect,
    add_text,
    rgb,
    set_run_typefaces,
    style_run,
)


SLIDE_W = 13.333
SLIDE_H = 7.5
CONTENT_X = 2.62
CONTENT_RIGHT = 12.82


def clone_slide(prs: Presentation, source):
    """在同一演示文稿中复制一页，同时复制图片关系。"""
    destination = prs.slides.add_slide(source.slide_layout)
    for shape in list(destination.shapes):
        element = shape._element
        element.getparent().remove(element)

    destination.background.fill.solid()
    destination.background.fill.fore_color.rgb = rgb(BG)
    for shape in source.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture = destination.shapes.add_picture(
                BytesIO(shape.image.blob),
                shape.left,
                shape.top,
                shape.width,
                shape.height,
            )
            picture.crop_left = shape.crop_left
            picture.crop_right = shape.crop_right
            picture.crop_top = shape.crop_top
            picture.crop_bottom = shape.crop_bottom
            picture.rotation = shape.rotation
            continue
        destination.shapes._spTree.insert_element_before(
            deepcopy(shape._element), "p:extLst"
        )
    return destination


def keep_slide(prs: Presentation, index: int) -> None:
    slide_ids = prs.slides._sldIdLst
    for current in reversed(range(len(slide_ids))):
        if current == index:
            continue
        slide_id = slide_ids[current]
        prs.part.drop_rel(slide_id.rId)
        del slide_ids[current]


def set_title(slide, title: str) -> None:
    """替换正文模板中的标题；模板的历史稿标题不是统一占位符。"""
    candidates = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        x = shape.left / 914400
        y = shape.top / 914400
        if 2.8 <= x <= 3.1 and 0.45 <= y <= 0.75:
            candidates.append(shape)
    if not candidates:
        raise RuntimeError("没有找到正文标题文本框")
    box = candidates[0]
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = title
    style_run(
        run,
        size=29,
        color=INK,
        bold=True,
        chinese_font="SimHei",
        latin_font="Times New Roman",
    )


def add_heading(slide, text: str, x: float, y: float, w: float, *, color=INK, size=16.0):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        0.34,
        size=size,
        color=color,
        bold=True,
        chinese_font="SimHei",
        valign=MSO_ANCHOR.TOP,
    )


def add_body(slide, text: str, x: float, y: float, w: float, h: float, *, size=12.8, color=BODY, bold=False, align=PP_ALIGN.LEFT):
    box = add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        chinese_font="SimSun",
        latin_font="Times New Roman",
        align=align,
        valign=MSO_ANCHOR.TOP,
        margin=0.02,
    )
    box.text_frame.word_wrap = True
    return box


def add_card(slide, x: float, y: float, w: float, h: float, title: str, detail: str, *, fill=PALE_BLUE, accent=BLUE, title_size=14.0, detail_size=12.2):
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.08, h, accent)
    add_heading(slide, title, x + 0.22, y + 0.14, w - 0.36, color=INK, size=title_size)
    add_body(slide, detail, x + 0.22, y + 0.52, w - 0.38, h - 0.61, size=detail_size)


def add_caption(slide, text: str, x: float, y: float, w: float, *, color=BODY, size=10.8):
    """将图表因果关系放在图形附近的简短注释中，不使用统一提示语。"""
    return add_body(
        slide,
        text,
        x,
        y,
        w,
        0.28,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_tag(slide, text: str, x: float, y: float, w: float, *, fill=BLUE, color=WHITE, size=11.5):
    add_rect(slide, x, y, w, 0.30, fill)
    add_text(slide, text, x, y + 0.01, w, 0.28, size=size, color=color, bold=True, chinese_font="SimHei", align=PP_ALIGN.CENTER)


def flow(slide, x1, y1, x2, y2, *, color=BLUE, direction="right", width=1.35):
    add_line(slide, x1, y1, x2, y2, color, width)
    add_arrow_tip(slide, x2, y2, direction, color, 0.10)


def node(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
    title_size=13.5,
    detail_size=11.5,
    title_color=None,
    detail_color=None,
):
    # 深色节点显式使用浅色字，避免主题样式把关键信息压成低对比度。
    if title_color is None:
        title_color = WHITE if fill == NAVY else INK
    if detail_color is None:
        detail_color = "D8E7F3" if fill == NAVY else BODY
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.07, h, accent)
    add_text(slide, title, x + 0.16, y + 0.11, w - 0.28, 0.28, size=title_size, color=title_color, bold=True, chinese_font="SimHei", align=PP_ALIGN.CENTER)
    add_body(slide, detail, x + 0.14, y + 0.48, w - 0.24, h - 0.54, size=detail_size, color=detail_color, align=PP_ALIGN.CENTER)


def draw_01(slide):
    set_title(slide, "可拓展内核模块（ELM）的概念定位")
    add_body(slide, "传统内核模块（如 LKM、KLD 等）具有以下结构性缺陷，会导致扩展越多、管理状态越分散，更新与撤销越依赖人工约定；代码虽然能够进入内核，却缺少完整的运行责任边界。", 3.00, 1.78, 9.62, 0.72, size=12.5, bold=True)
    add_card(slide, 3.00, 2.90, 2.16, 1.42, "身份缺口", "模块名、对象地址与版本信息分散；无法稳定回答“当前运行的是谁”。", fill=PALE_GRAY, accent=MUTED, title_size=13.8, detail_size=10.9)
    add_card(slide, 5.48, 2.90, 2.16, 1.42, "关系缺口", "符号依赖不等于能力关系；调用者、提供者和替换影响范围难以还原。", fill=PALE_GRAY, accent=MUTED, title_size=13.8, detail_size=10.9)
    add_card(slide, 7.96, 2.90, 2.16, 1.42, "资源缺口", "回调、队列、内存与设备对象缺少统一所有者和撤销次序。", fill=PALE_GRAY, accent=MUTED, title_size=13.8, detail_size=10.9)
    add_card(slide, 10.44, 2.90, 2.16, 1.42, "证据缺口", "装载、调用、策略和故障没有共同证据链，问题难以复核。", fill=PALE_GRAY, accent=MUTED, title_size=13.8, detail_size=10.9)
    add_heading(slide, "ELM 的责任单元方法", 3.00, 4.63, 3.6, color=NAVY, size=16.0)
    add_body(slide, "ELM 不把扩展只看作一段可装入代码，而把每个扩展建立为一个可识别、可约束、可计量、可撤销、可审计的运行责任单元；内核统一管理其完整进入、运行、换代和退出过程。", 3.00, 5.04, 9.55, 0.80, size=13.0, color=INK, bold=True)
    add_caption(slide, "四类缺口 → 责任单元：身份、关系、资源、证据归于同一管理对象", 3.00, 6.10, 9.58, size=10.6)


def draw_02(slide):
    set_title(slide, "运行责任单元模型")
    add_body(slide, "责任单元方法首先需要一个内核可识别的对象。ELM 将每个受管扩展登记为 Cell；后续的能力、资源、调用、策略和故障都必须归属于这个 Cell。", 3.00, 1.78, 9.56, 0.68, size=12.6, bold=True)
    node(slide, 6.15, 2.96, 3.25, 1.48, "CellRuntime", "内核中的真实运行记录\n不是文件句柄，也不是模块地址", fill=NAVY, accent=BLUE, title_size=18, detail_size=11.9)
    add_card(slide, 3.00, 2.86, 2.50, 1.05, "单元身份", "ElmId：本次启动实例中的唯一身份；不会暴露为裸内核地址。", fill=PALE_PURPLE, accent=PURPLE, title_size=13.8, detail_size=10.2)
    add_card(slide, 10.05, 2.86, 2.55, 1.05, "单元代际", "Generation：同一逻辑单元替换前后的版本代，不允许旧引用混入新实现。", fill=PALE_TEAL, accent=TEAL, title_size=13.8, detail_size=10.2)
    add_card(slide, 3.00, 4.25, 2.50, 1.05, "单元状态", "State：发现、验证、活动、暂停、隔离、退役等调用门禁。", fill=PALE_BLUE, accent=BLUE, title_size=13.8, detail_size=10.2)
    add_card(slide, 10.05, 4.25, 2.55, 1.05, "单元类型", "Kind 与 Parent：说明职责类别以及在管理树中的父级责任。", fill=PALE_GRAY, accent=MUTED, title_size=13.8, detail_size=10.2)
    flow(slide, 5.50, 3.38, 6.12, 3.38, color=PURPLE)
    flow(slide, 9.40, 3.38, 10.02, 3.38, color=TEAL)
    flow(slide, 5.50, 4.76, 6.12, 4.10, color=BLUE)
    flow(slide, 9.40, 4.10, 10.02, 4.76, color=MUTED)
    add_card(slide, 4.12, 5.65, 7.34, 0.72, "CellRuntime 责任账本", "身份 · 父子关系 · 状态 · 代际 · EBI 来源 · ABI · 策略 · 预算 · 活动调用 · 故障 · 所属 Binding 与菜单项", fill=PALE_BLUE, accent=BLUE, title_size=13.8, detail_size=10.9)
    add_caption(slide, "ElmId：责任主体；Generation：新旧隔离；State：调用门禁", 3.00, 6.48, 9.58, size=10.6)


def draw_03(slide):
    set_title(slide, "管理根与受管拓扑")
    add_body(slide, "Cell 有了身份之后，还需要一个不会被普通扩展绕过的约束层，以及一个统一接受外部管理请求的根单元。", 3.00, 1.78, 9.52, 0.55, size=12.6, bold=True)
    add_rect(slide, 3.02, 2.86, 4.28, 2.20, NAVY)
    add_text(slide, "ELM Core", 3.32, 3.08, 1.48, 0.34, size=20, color=WHITE, bold=True)
    add_body(slide, "常驻最小可信根", 5.02, 3.11, 1.90, 0.26, size=12.2, color="D8E7F3", bold=True, align=PP_ALIGN.RIGHT)
    add_body(slide, "状态机与事务提交\n关系图与代际校验\n租约、预算与执行引用\n原生调用门与故障恢复\n事件、审计与健康检查", 3.35, 3.60, 3.55, 1.22, size=11.5, color=WHITE, bold=True)
    add_rect(slide, 8.10, 2.86, 4.50, 2.20, PALE_BLUE)
    add_rect(slide, 8.10, 2.86, 0.09, 2.20, BLUE)
    add_text(slide, "elm-mgr", 8.42, 3.08, 1.50, 0.34, size=20, color=NAVY, bold=True)
    add_body(slide, "根管理 Cell · ElmId = 1", 10.10, 3.11, 2.12, 0.26, size=12.0, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    add_body(slide, "统一接收 elmctl / sys_elm_ctl\n组织父子 Cell 管理树\n承接策略、菜单、装载与绑定请求\n把管理意图转换为 Core 可验证事务\n自身也进入身份、预算与状态拓扑", 8.42, 3.60, 3.75, 1.22, size=11.4, color=INK, bold=True)
    flow(slide, 7.30, 3.94, 8.06, 3.94, color=TEAL)
    node(slide, 4.04, 5.43, 2.46, 0.80, "eki 子单元", "ElmId = 2\n提供 EKI 投影能力", fill=PALE_PURPLE, accent=PURPLE, title_size=13.5, detail_size=10.4)
    node(slide, 7.05, 5.43, 2.46, 0.80, "普通 Cell", "合法父级 + 身份 + 预算\n进入同一拓扑", fill=PALE_TEAL, accent=TEAL, title_size=13.5, detail_size=10.4)
    node(slide, 10.06, 5.43, 2.46, 0.80, "子管理 Cell", "管理机制本身仍受\n相同关系约束", fill=PALE_GRAY, accent=MUTED, title_size=13.5, detail_size=10.4)
    add_caption(slide, "Core：硬约束；elm-mgr：管理入口；Cell：受管对象", 3.00, 6.46, 9.58, size=10.6)


def draw_04(slide):
    set_title(slide, "声明式拓展表面")
    add_body(slide, "传统模块往往在 init 中执行任意注册逻辑，内核只能在副作用发生后理解它。ELM 把关键关系提前声明，使 Core 在运行代码之前构造和校验预期拓扑。", 3.00, 1.78, 9.54, 0.68, size=12.5, bold=True)
    add_rect(slide, 3.02, 2.96, 2.42, 2.54, NAVY)
    add_text(slide, "Manifest", 3.34, 3.18, 1.78, 0.34, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_body(slide, "规范名称\n版本与 Kind\n能力需求\n能力提供\n依赖关系\n开放拓展位置", 3.42, 3.70, 1.62, 1.48, size=12.0, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 6.00, 2.96, 3.05, 1.12, "能力需求声明", "Intent：说明希望消费、提供、观察、控制或拓展哪一种能力。", fill=PALE_PURPLE, accent=PURPLE, title_size=14.0, detail_size=11.0)
    add_card(slide, 9.55, 2.96, 3.05, 1.12, "能力提供声明", "Offer：说明能够提供哪一种契约，以及共享、独占、有序等分发模式。", fill=PALE_TEAL, accent=TEAL, title_size=14.0, detail_size=11.0)
    add_card(slide, 6.00, 4.38, 3.05, 1.12, "显式拓展位置", "Extension Point：由目标单元主动开放，限定名称、契约、组合模式和处理阶段。", fill=PALE_BLUE, accent=BLUE, title_size=14.0, detail_size=10.9)
    add_card(slide, 9.55, 4.38, 3.05, 1.12, "受控拓展挂接", "Extension：经过策略、契约和关系图校验后附着；不存在任意地址补丁入口。", fill=PALE_GRAY, accent=MUTED, title_size=14.0, detail_size=10.9)
    flow(slide, 5.44, 3.54, 5.96, 3.54, color=PURPLE)
    flow(slide, 5.44, 4.94, 5.96, 4.94, color=BLUE)
    add_caption(slide, "Manifest 只声明；策略、预算、契约与证明决定装载和挂接", 3.00, 6.04, 9.58, size=10.5)
    add_body(slide, "可拓展性的来源不是允许任意修改，而是让新的契约和实现通过统一声明进入拓扑，同时保持 Core 与具体设备、网络或 VFS 语义解耦。", 3.00, 6.44, 9.58, 0.35, size=11.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def draw_05(slide):
    set_title(slide, "能力契约与连接句柄")
    add_body(slide, "声明表达“想要什么”和“能够提供什么”，但只有契约、方向、分发模式、访问策略和代际全部兼容，Core 才会提交实际连接。", 3.00, 1.78, 9.54, 0.62, size=12.5, bold=True)
    labels = [("Intent", "能力需求", PALE_PURPLE, PURPLE), ("PortId", "稳定能力端点", PALE_BLUE, BLUE), ("BindingId", "已提交关系", PALE_TEAL, TEAL), ("LeaseId", "可撤销使用权", PALE_GRAY, MUTED), ("Provider", "实际执行后端", NAVY, NAVY)]
    x = 3.00
    widths = [1.56, 1.70, 1.76, 1.68, 2.00]
    for i, ((title, detail, fill, accent), w) in enumerate(zip(labels, widths)):
        node(slide, x, 3.02, w, 1.22, title, detail, fill=fill, accent=accent, title_size=13.3, detail_size=10.4)
        if i < len(labels) - 1:
            flow(slide, x + w, 3.63, x + w + 0.24, 3.63, color=accent)
        x += w + 0.24
    add_card(slide, 3.00, 4.62, 2.90, 1.20, "版本化流契约", "FlowContract 使用 name@version；完整比较契约，而不是按符号前缀或类型名猜测兼容性。", fill=PALE_PURPLE, accent=PURPLE, title_size=13.8, detail_size=10.8)
    add_card(slide, 6.20, 4.62, 2.90, 1.20, "连接关系记录", "Binding 记录 consumer、Port、契约、Generation、Lease 和 active 状态。", fill=PALE_BLUE, accent=BLUE, title_size=13.8, detail_size=10.8)
    add_card(slide, 9.40, 4.62, 3.20, 1.20, "可撤销引用状态", "Lease 记录 owner、权限、Generation、状态和 active_refs；撤销先拒绝新引用，再等待归零。", fill=PALE_TEAL, accent=TEAL, title_size=13.8, detail_size=10.7)
    add_caption(slide, "Intent → PortId → BindingId → LeaseId → Provider：需求、端点、连接、引用、执行", 3.00, 6.16, 9.58, size=10.5)
    add_body(slide, "当前生产启动形成闭环的内建端口：core.log@1 · core.event@1 · mgr.menu.item@1 · mgr.action.invoke@1。", 3.00, 6.54, 9.58, 0.26, size=10.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def draw_06(slide):
    set_title(slide, "分层执行路径")
    add_body(slide, "统一管理不等于所有调用都经过管理器。ELM 将需要动态治理的调用与已经证明稳定的本地调用分开，使安全检查和运行效率同时成立。", 3.00, 1.78, 9.54, 0.64, size=12.5, bold=True)
    add_card(slide, 3.00, 2.98, 4.50, 2.42, "受管能力路径", "适用范围\n运行时发现 · 多提供者选择 · 访问策略 · 审计 · 异步队列 · 背压 · 取消 · 超时\n\n执行链\nCell → PortId → BindingId / LeaseId → Provider → Reply\n\n每次调用可重新验证代际、策略、后端 epoch 和活动引用。", fill=PALE_BLUE, accent=BLUE, title_size=15.2, detail_size=11.0)
    add_card(slide, 8.10, 2.98, 4.50, 2.42, "证明后直连路径", "适用范围\n常驻内核 API · ABI 严格固定的 ELM 间调用\n\n执行链\nkernel-symbol / direct-pinned → 类型化固定槽 → 真实 Rust 函数\n\n完整校验发生在模块代码执行前；调用热路径不再经过 elm-mgr 或 Provider 帧。", fill=PALE_TEAL, accent=TEAL, title_size=15.2, detail_size=11.0)
    add_tag(slide, "治理优先", 3.00, 5.72, 1.10, fill=BLUE, size=11.1)
    add_body(slide, "动态路径保留发现、授权、审计和流控语义。", 4.25, 5.75, 3.18, 0.25, size=11.4, bold=True)
    add_tag(slide, "延迟优先", 8.10, 5.72, 1.10, fill=TEAL, size=11.1)
    add_body(slide, "直连路径把验证前置，不在每次调用中重复付费。", 9.35, 5.75, 3.20, 0.25, size=11.4, bold=True)
    add_caption(slide, "受管路径保留治理；直连路径前置证明", 3.00, 6.27, 9.58, size=10.6)


def draw_07(slide):
    set_title(slide, "装载协议与接口证明")
    add_body(slide, "动态代码进入内核之前，必须同时解决“文件格式如何解释”和“目标接口是否精确兼容”两个问题。ELM Core 不直接理解 EKI、ELF 或未来格式。", 3.00, 1.78, 9.54, 0.62, size=12.5, bold=True)
    stages = [("外部容器", "EKI / 未来格式", PALE_GRAY, MUTED), ("Projection Source", "格式解析与边界校验", PALE_PURPLE, PURPLE), ("EBI", "Core 消费的装载协议对象", PALE_BLUE, BLUE), ("接口证明", "名称 · 契约 · 版本\nProfile · ABI SHA-256 · 能力", PALE_TEAL, TEAL), ("原生镜像", "映射 · 重定位 · 入口", NAVY, NAVY)]
    x = 3.00
    widths = [1.55, 1.92, 1.55, 2.20, 1.80]
    for i, ((title, detail, fill, accent), w) in enumerate(zip(stages, widths)):
        node(slide, x, 3.00, w, 1.34, title, detail, fill=fill, accent=accent, title_size=12.6, detail_size=10.3)
        if i < len(stages) - 1:
            flow(slide, x + w, 3.67, x + w + 0.24, 3.67, color=accent)
        x += w + 0.24
    add_card(slide, 3.00, 4.72, 4.45, 1.20, "EBI 协议边界", "EBI 只表达架构、清单、段、拓扑、生命周期描述符、import/export 和 ABI 摘要；它不是文件格式。", fill=PALE_BLUE, accent=BLUE, title_size=13.8, detail_size=10.7)
    add_card(slide, 8.10, 4.72, 4.50, 1.20, "Kernel API Profile", "由目标内核实际 rmeta / rlib 生成；接口路径、类型布局、rustc、target 或 feature 变化会在执行前拒绝旧镜像。", fill=PALE_TEAL, accent=TEAL, title_size=13.8, detail_size=10.7)
    add_caption(slide, "容器 → Projection Source → EBI → Profile / ABI → 原生入口", 3.00, 6.24, 9.58, size=10.5)


def draw_08(slide):
    set_title(slide, "权限、预算与资源所有权")
    add_body(slide, "传统模块加载成功后往往共享宽泛内核权限，资源成本也难以归属。ELM 将能力策略、预算和长期资源登记同时绑定到 Cell。", 3.00, 1.78, 9.54, 0.60, size=12.5, bold=True)
    add_card(slide, 3.00, 2.92, 3.00, 2.12, "能力策略", "Capability Policy\n\n每次管理动作、Provider 注册/调用、事件订阅和 Mixin 挂接均检查权限位、Cell 状态、Generation 与 policy epoch。\n\n父级不能向子级委派自身没有的能力；策略可进入不可逆 LOCKED 状态。", fill=PALE_PURPLE, accent=PURPLE, title_size=14.5, detail_size=10.7)
    add_card(slide, 6.30, 2.92, 3.00, 2.12, "分层预算", "Resource Budget\n\n动态内存、原生栈、并发调用、Provider 端口与队列、事件订阅、镜像、审计和故障记录进入账本。\n\n父级总预算必须覆盖自身用量与所有存活子级保留预算；超限在操作前拒绝。", fill=PALE_BLUE, accent=BLUE, title_size=14.5, detail_size=10.7)
    add_card(slide, 9.60, 2.92, 3.00, 2.12, "长期资源", "Owned Resource\n\n任务、定时器、工作项、回调、IRQ 回调和异步请求登记所属 Cell 与生命周期操作。\n\n退役按 quiesce → cancel → drain → release 排空，避免镜像释放后留下回调。", fill=PALE_TEAL, accent=TEAL, title_size=14.5, detail_size=10.7)
    add_card(slide, 3.00, 5.42, 4.50, 0.88, "Lease 与 Owned Resource 边界", "Lease 负责绑定与控制面的可撤销使用权；长期异步对象由 Owned Resource 登记完整排空操作。", fill=PALE_GRAY, accent=MUTED, title_size=13.6, detail_size=10.4)
    add_card(slide, 8.10, 5.42, 4.50, 0.88, "当前实现边界", "CPU 预算当前主要用于记账和超限统计，尚未接入调度器节流；部分设备资源会明确阻断 Pause。", fill=PALE_GRAY, accent=MUTED, title_size=13.6, detail_size=10.4)
    add_caption(slide, "策略：允许做什么；预算：最多消耗多少；资源登记：退出顺序", 3.00, 6.46, 9.58, size=10.6)


def draw_09(slide):
    set_title(slide, "生命周期与代际事务")
    add_body(slide, "装载、暂停、替换和卸载不是直接修改状态字段，而是由同一状态机与事务协议控制；任意 blocker 都在提交前暴露。", 3.00, 1.78, 9.54, 0.58, size=12.5, bold=True)
    states = [("Discovered", PALE_GRAY, MUTED), ("Verified", PALE_PURPLE, PURPLE), ("Loaded", PALE_BLUE, BLUE), ("Linked", PALE_TEAL, TEAL), ("Ready", PALE_BLUE, BLUE), ("Active", NAVY, NAVY)]
    x = 3.02
    for i, (label, fill, accent) in enumerate(states):
        w = 1.38 if label != "Discovered" else 1.55
        add_rect(slide, x, 2.88, w, 0.52, fill)
        add_text(slide, label, x, 2.98, w, 0.25, size=11.5, color=WHITE if fill == NAVY else INK, bold=True, align=PP_ALIGN.CENTER)
        if i < len(states)-1:
            flow(slide, x+w, 3.14, x+w+0.22, 3.14, color=accent)
        x += w + 0.22
    add_card(slide, 3.02, 3.80, 3.02, 1.70, "生命周期事务", "Active → Quiescing：关闭新工作入口\nQuiescing：等待调用、Lease、Provider 队列与受托资源排空\nPaused：保留镜像和状态但不服务\nDetached → Retired：摘除拓扑并完成回收", fill=PALE_BLUE, accent=BLUE, title_size=14.2, detail_size=10.6)
    add_card(slide, 6.30, 3.80, 3.02, 1.70, "换代事务", "Generation G → G+1\n影子装载新代 → 初始化 → 静默旧代 → 导出/导入迁移状态 → 提交点切换 Provider、Source、Binding 和 Import → 退役旧镜像", fill=PALE_PURPLE, accent=PURPLE, title_size=14.2, detail_size=10.6)
    add_card(slide, 9.58, 3.80, 3.02, 1.70, "失败回滚", "新代验证、初始化或迁移失败：撤销暂存 Import、Source 与新镜像，并恢复旧代。\n旧代恢复失败：隔离单元并保留诊断状态。\n运行中调用不跨 Generation 延续。", fill=PALE_TEAL, accent=TEAL, title_size=14.2, detail_size=10.6)
    add_caption(slide, "G → G+1：影子装载 → 静默旧代 → 状态迁移 → 提交；失败回滚", 3.00, 5.92, 9.58, size=10.5)
    add_body(slide, "高可维护性的来源：变更失败不会留下半提交拓扑，新旧代际不会共享未经迁移的隐式状态。", 3.00, 6.38, 9.58, 0.30, size=11.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def draw_10(slide):
    set_title(slide, "原生执行故障边界")
    add_body(slide, "宏内核原生扩展不能依赖用户态进程隔离。ELM 为生命周期钩子、原生 Provider、Entry 和迁移入口建立独立调用门与固定恢复出口。", 3.00, 1.78, 9.54, 0.60, size=12.5, bold=True)
    stages = [("执行预检", "Cell / Generation / Policy", PALE_BLUE, BLUE), ("保护调用门", "独立 64 KiB 栈\n双 Guard 页", PALE_PURPLE, PURPLE), ("原生代码", "Hook · Provider\nEntry · Migration", PALE_TEAL, TEAL), ("异常现场", "Fault · Panic\nTimer Timeout", PALE_GRAY, MUTED), ("固定恢复出口", "重写 PC / SP / 返回值", NAVY, NAVY)]
    x = 3.00
    widths = [1.68, 1.78, 1.75, 1.72, 2.00]
    for i, ((title, detail, fill, accent), w) in enumerate(zip(stages, widths)):
        node(slide, x, 2.94, w, 1.20, title, detail, fill=fill, accent=accent, title_size=12.3, detail_size=10.2)
        if i < len(stages)-1:
            flow(slide, x+w, 3.54, x+w+0.22, 3.54, color=accent)
        x += w + 0.22
    add_card(slide, 3.00, 4.56, 4.50, 1.34, "结构化故障现场", "Cell · Phase · CPU · PC · Fault Address · Exception Code · Recovery PC · Recovery SP\n恢复过程不读取故障现场 ra，深层嵌套故障也不会跳回镜像内部。", fill=PALE_PURPLE, accent=PURPLE, title_size=14.2, detail_size=10.7)
    add_card(slide, 8.10, 4.56, 4.50, 1.34, "隔离门禁", "native_faults 与 isolated 进入 CellRuntime；随后阻止新 Binding、Provider 调用和 Import 解析，同时保留 Fault Dump 供诊断与受控 Detach。", fill=PALE_TEAL, accent=TEAL, title_size=14.2, detail_size=10.7)
    add_caption(slide, "预检 → Guard → 原生执行 → fault / panic / timeout → 固定恢复出口", 3.00, 6.20, 9.58, size=10.5)
    add_body(slide, "边界：这是共享地址空间中的软隔离和受控恢复，不是 MMU 沙箱；Kernel Provider 回调不属于原生 ELM Guard。", 3.00, 6.56, 9.58, 0.25, size=10.7, color=BODY, bold=True, align=PP_ALIGN.CENTER)


def draw_11(slide):
    set_title(slide, "运行证据与性质闭环")
    add_body(slide, "只有拓扑、决策、成本、调用与故障能够被共同还原，扩展系统的安全性和可维护性才不是不可验证的口号。", 3.00, 1.78, 9.54, 0.55, size=12.5, bold=True)
    add_card(slide, 3.00, 2.82, 2.86, 1.26, "身份与关系证据", "Cell · Parent · Dependency\nPort · Provider · Binding · Lease\n回答“系统里有哪些对象，它们怎样相连”。", fill=PALE_BLUE, accent=BLUE, title_size=13.8, detail_size=10.6)
    add_card(slide, 6.18, 2.82, 2.86, 1.26, "行为与决策证据", "Event · Audit · Policy · Trust\nLifecycle / Provider / Mixin / Replace Trace\n回答“谁触发、谁批准、为什么允许”。", fill=PALE_PURPLE, accent=PURPLE, title_size=13.8, detail_size=10.6)
    add_card(slide, 9.36, 2.82, 3.24, 1.26, "成本与故障证据", "Execution · Resource Accounting · Owned Resource · Worker · Fault · Diagnostics\n回答“正在做什么、成本归谁、失败在哪里”。", fill=PALE_TEAL, accent=TEAL, title_size=13.8, detail_size=10.6)
    add_card(slide, 3.00, 4.38, 2.18, 1.38, "可拓展性", "Manifest + 版本化契约 + 显式拓展点\n新能力无需把具体子系统类型写进 Core。", fill=PALE_BLUE, accent=BLUE, title_size=14.2, detail_size=10.0)
    add_card(slide, 5.48, 4.38, 2.18, 1.38, "高管理性", "Cell + Binding Graph + 状态机 + elm-mgr\n对象、关系和变更进入统一拓扑。", fill=PALE_PURPLE, accent=PURPLE, title_size=14.2, detail_size=10.0)
    add_card(slide, 7.96, 4.38, 2.18, 1.38, "高效率", "验证前置 + kernel-symbol / direct-pinned\n稳定热路径不重复经过管理器。", fill=PALE_TEAL, accent=TEAL, title_size=14.2, detail_size=10.0)
    add_card(slide, 10.44, 4.38, 2.18, 1.38, "高安全性", "策略、预算、代际、租约、资源排空、恢复与审计\n非法状态在提交前拒绝。", fill=PALE_GRAY, accent=MUTED, title_size=14.2, detail_size=10.0)
    add_caption(slide, "拓扑、审计、执行、资源、故障 → /sys/kernel/elm；Health 17 类检查；Journal 前序哈希链", 3.00, 6.10, 9.58, size=10.3)
    add_body(slide, "当前边界：未登记生产持久化 Journal 后端时运行在易失模式；子系统 Provider 需要由各子系统显式注册，未完成能力不会被包装成已完成。", 3.00, 6.48, 9.58, 0.31, size=10.3, color=BODY, bold=True, align=PP_ALIGN.CENTER)


DRAWERS = [draw_01, draw_02, draw_03, draw_04, draw_05, draw_06, draw_07, draw_08, draw_09, draw_10, draw_11]


def build(template: Path, output: Path) -> None:
    prs = Presentation(template)
    # 第三章设备页的模板包含最新侧栏、Logo、页眉线和页码；去掉其正文图形。
    base = prs.slides[19]
    keep_slide(prs, 19)
    for shape in list(base.shapes):
        # 保留 0..20 的模板元素；第 20 个之后是设备页正文。
        if shape in list(base.shapes)[21:]:
            element = shape._element
            element.getparent().remove(element)
    slides = [base]
    for _ in range(len(DRAWERS) - 1):
        slides.append(clone_slide(prs, base))
    for slide, drawer in zip(slides, DRAWERS):
        # 清除上一页留下的正文，防止模板修改时正文索引发生变化。
        for shape in list(slide.shapes)[21:]:
            element = shape._element
            element.getparent().remove(element)
        drawer(slide)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> int:
    root = Path(__file__).resolve().parents[2]
    template = root / "output/presentations/mygo-defense-full.pptx"
    output = root / "output/presentations/mygo-defense-elm-pages1-11.pptx"
    if len(sys.argv) > 1:
        output = Path(sys.argv[1]).resolve()
    build(template, output)
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
