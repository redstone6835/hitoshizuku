#!/usr/bin/env python3
"""生成答辩第六章“总结展望”的两页正文。"""

from __future__ import annotations

import argparse
import os
from pathlib import Path
import tempfile

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from generate_chapter5_results_slides import (
    BLUE,
    BODY,
    INK,
    LINE,
    MIN_FONT_PT,
    MUTED,
    NAVY,
    PALE_BLUE,
    PALE_GRAY,
    PALE_PURPLE,
    PALE_TEAL,
    PURPLE,
    TEAL,
    WHITE,
    clean_body_template,
    clone_slide,
    enforce_font_floor,
    find_slide,
    heading,
    inches,
    remove_slide,
    set_title,
    slide_texts,
    table_cell,
)
from generate_engineering_structure_slide import style_run


CHAPTER6_TRANSITION = "第六章 · 总结展望"
CHAPTER6_TITLES = ("工程创新", "后续发展方向")


def update_transition(slide) -> None:
    candidates = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and 3.10 <= inches(shape.top) <= 3.45
        and shape.text.strip()
    ]
    if not candidates:
        return
    frame = candidates[0].text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "工程创新与发展方向"
    style_run(
        run,
        size=18,
        color=BODY,
        bold=False,
        chinese_font="SimSun",
        latin_font="Times New Roman",
    )


def draw_table(
    slide,
    headers: tuple[str, str, str],
    rows: tuple[tuple[str, str, str, str, str], ...],
) -> None:
    columns = (3.00, 5.16, 9.04, 12.60)
    for index, text in enumerate(headers):
        table_cell(
            slide,
            text,
            columns[index],
            1.86,
            columns[index + 1] - columns[index],
            0.58,
            size=16,
            color=WHITE,
            bold=True,
            fill=NAVY,
            align=PP_ALIGN.CENTER,
            chinese_font="SimHei",
        )

    y = 2.44
    for title, mechanism, result, fill, accent in rows:
        table_cell(
            slide,
            title,
            columns[0],
            y,
            columns[1] - columns[0],
            0.88,
            size=15,
            color=accent,
            bold=True,
            fill=fill,
            align=PP_ALIGN.CENTER,
            chinese_font="SimHei",
        )
        table_cell(
            slide,
            mechanism,
            columns[1],
            y,
            columns[2] - columns[1],
            0.88,
            size=14,
            color=INK,
            bold=True,
            fill=fill,
        )
        table_cell(
            slide,
            result,
            columns[2],
            y,
            columns[3] - columns[2],
            0.88,
            size=14,
            color=INK,
            fill="FFFFFF",
        )
        y += 0.88


def draw_01(slide) -> None:
    set_title(slide, CHAPTER6_TITLES[0])
    rows = (
        (
            "可管理内核单元运行时",
            "Cell表示运行单元，Generation区分替换前后的实现，Binding与Lease记录调用关系和资源归属",
            "装载、暂停、替换、故障隔离与退役进入同一套生命周期",
            PALE_BLUE,
            BLUE,
        ),
        (
            "审核式Rust接口直连",
            "装载前核验接口名称、版本、能力、ABI摘要与目标Profile",
            "核验通过后直接调用常驻Rust实现，热路径保持直接调用形式",
            PALE_TEAL,
            TEAL,
        ),
        (
            "同源双形态组件构建",
            "同一份Rust源码由y/m/n配置选择静态集成、受管EKI或关闭构建",
            "Modules.toml统一检查组件依赖，双架构分别生成匹配产物",
            PALE_PURPLE,
            PURPLE,
        ),
        (
            "开放式类型化设备模型",
            "PnP统一发现、匹配、资源归属和移除，DeviceFunction承载开放能力与生命周期",
            "字符、块、网络与RTC保留类型化语义，用户视图由投影层生成",
            PALE_BLUE,
            BLUE,
        ),
        (
            "可复核性能证据链",
            "QEMU动态指令、运行地址与同次构建的内核符号快照绑定",
            "固定负载、成对对照和多轮复测把开销归属到函数与责任阶段",
            PALE_GRAY,
            MUTED,
        ),
    )
    draw_table(slide, ("创新设计", "实现方式", "工程落点"), rows)


def draw_02(slide) -> None:
    set_title(slide, CHAPTER6_TITLES[1])
    rows = (
        (
            "设备与ELM生命周期",
            "继续细化设备退出状态机，统一中断、DMA、总线资源、协议栈句柄与用户态节点的回收顺序",
            "形成可审计的PnP资源图与稳定的设备属性ABI",
            PALE_BLUE,
            BLUE,
        ),
        (
            "内存压力与回收",
            "统筹伙伴分配器、Slab、页缓存、文件回写和匿名页，完善分级回收与缓存收缩器",
            "形成统一压力模型与可配置的内存审计能力",
            PALE_TEAL,
            TEAL,
        ),
        (
            "VFS与POSIX语义",
            "以系统调用语义矩阵持续验证标志组合、错误码和副作用，统一权限检查与目录项状态",
            "形成稳定的兼容语义与文件系统驱动契约",
            PALE_PURPLE,
            PURPLE,
        ),
        (
            "网络并发与协议能力",
            "分层组织设备中断、收发批处理、协议定时器和套接字唤醒，扩展IPv6、错误队列与缓冲所有权",
            "形成面向多连接负载的并发推进与协议接口",
            PALE_BLUE,
            BLUE,
        ),
        (
            "用户程序与多核运行",
            "继续完善ELF校验、解释器、TLS、辅助向量和栈随机化，优化VMA管理、运行队列迁移与唤醒选核",
            "持续提升复杂用户程序和多核长期负载的运行能力",
            PALE_GRAY,
            MUTED,
        ),
    )
    draw_table(slide, ("发展方向", "下一阶段工作", "目标形态"), rows)


DRAWERS = (draw_01, draw_02)


def insert_chapter(prs: Presentation) -> list:
    transition = find_slide(prs, CHAPTER6_TRANSITION)
    transition_index = list(prs.slides).index(transition)
    existing = list(prs.slides)[transition_index + 1 : -1]
    if not existing:
        raise RuntimeError("第六章缺少可复用正文模板")
    template = existing[0]
    update_transition(transition)

    added = []
    for drawer in DRAWERS:
        slide = clone_slide(prs, template)
        clean_body_template(slide)
        drawer(slide)
        enforce_font_floor(slide)
        added.append(slide)

    for slide in existing:
        remove_slide(prs, slide)

    transition_index = list(prs.slides).index(transition)
    slide_ids = prs.slides._sldIdLst
    added_ids = list(slide_ids)[-len(added) :]
    for slide_id in added_ids:
        slide_ids.remove(slide_id)
    for offset, slide_id in enumerate(added_ids, 1):
        slide_ids.insert(transition_index + offset, slide_id)
    return added


def atomic_save(prs: Presentation, output: Path) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        prefix=f".{output.stem}-",
        suffix=".pptx",
        dir=output.parent,
        delete=False,
    ) as temporary:
        temporary_path = Path(temporary.name)
    try:
        prs.save(temporary_path)
        os.replace(temporary_path, output)
        output.chmod(0o644)
    except Exception:
        temporary_path.unlink(missing_ok=True)
        raise


def build_topic(full_output: Path, topic_output: Path) -> None:
    prs = Presentation(full_output)
    keep_titles = set(CHAPTER6_TITLES)
    for slide in list(prs.slides):
        if not (slide_texts(slide) & keep_titles):
            remove_slide(prs, slide)
    if len(prs.slides) != len(CHAPTER6_TITLES):
        raise RuntimeError(f"第六章专题页数错误：{len(prs.slides)}")
    atomic_save(prs, topic_output)


def validate(path: Path, *, expected_slides: int | None = None) -> None:
    prs = Presentation(path)
    if expected_slides is not None and len(prs.slides) != expected_slides:
        raise RuntimeError(f"{path} 页数错误：{len(prs.slides)}")
    titles_found = []
    forbidden = ("不足", "缺少", "尚未", "做得不好", "不完善", "不能")
    for slide_number, slide in enumerate(prs.slides, 1):
        texts = slide_texts(slide)
        titles_found.extend(title for title in CHAPTER6_TITLES if title in texts)
        if not (texts & set(CHAPTER6_TITLES)):
            continue
        joined = "\n".join(texts)
        for phrase in forbidden:
            if phrase in joined:
                raise RuntimeError(f"第 {slide_number} 页包含负向自评：{phrase}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip() or run.font.size is None:
                        continue
                    if run.font.size.pt < MIN_FONT_PT:
                        raise RuntimeError(
                            f"第 {slide_number} 页文字小于 {MIN_FONT_PT:g} pt：{run.text!r}"
                        )
    if sorted(titles_found) != sorted(CHAPTER6_TITLES):
        raise RuntimeError("第六章标题集合不完整或重复")


def main() -> int:
    parser = argparse.ArgumentParser()
    root = Path(__file__).resolve().parents[2]
    parser.add_argument(
        "--base",
        type=Path,
        default=root / "output/presentations/mygo-defense-full.pptx",
    )
    parser.add_argument(
        "--full-output",
        type=Path,
        default=root / "output/presentations/mygo-defense-full.pptx",
    )
    parser.add_argument(
        "--topic-output",
        type=Path,
        default=root / "output/presentations/mygo-defense-chapter6-outlook-2pages.pptx",
    )
    args = parser.parse_args()

    base = args.base.resolve()
    full_output = args.full_output.resolve()
    topic_output = args.topic_output.resolve()
    prs = Presentation(base)
    insert_chapter(prs)
    atomic_save(prs, full_output)
    validate(full_output)
    build_topic(full_output, topic_output)
    validate(topic_output, expected_slides=len(CHAPTER6_TITLES))
    print(full_output)
    print(topic_output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
