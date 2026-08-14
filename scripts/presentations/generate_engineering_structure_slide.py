#!/usr/bin/env python3
"""生成答辩第三章第一页“工程结构”。"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400

BG = "FAFBFC"
NAVY = "082D5B"
BLUE = "155A9D"
TEAL = "167C80"
PURPLE = "55427E"
INK = "17212B"
BODY = "435766"
MUTED = "748491"
LINE = "D9E3EA"
PALE_BLUE = "EAF2F8"
PALE_TEAL = "E8F3F2"
PALE_PURPLE = "F2EFF7"
PALE_GRAY = "EEF2F5"
WHITE = "FFFFFF"

CJK_RE = re.compile(r"([\u3400-\u9fff\uf900-\ufaff]+)")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_no_line(shape) -> None:
    shape.line.fill.background()


def flatten_shape(shape) -> None:
    """移除主题效果引用，保证新增图形保持纯平面样式。"""
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)
    properties = shape._element.spPr
    for tag in ("a:effectLst", "a:effectDag"):
        effect = properties.find(qn(tag))
        if effect is not None:
            properties.remove(effect)
    properties.append(OxmlElement("a:effectLst"))


def add_rect(slide, x, y, w, h, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        set_no_line(shape)
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    flatten_shape(shape)
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.25, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    flatten_shape(line)
    return line


def add_arrow_tip(slide, x, y, direction, color=BLUE, size=0.10):
    rotations = {"up": 0, "right": 90, "down": 180, "left": 270}
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x - size / 2),
        Inches(y - size / 2),
        Inches(size),
        Inches(size),
    )
    shape.rotation = rotations[direction]
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    set_no_line(shape)
    flatten_shape(shape)
    return shape


def set_run_typefaces(run, chinese_font: str, latin_font: str) -> None:
    """同时写入 latin/ea/cs，避免 Office 按主题字体回退。"""
    properties = run._r.get_or_add_rPr()
    for tag, typeface in (
        ("a:latin", latin_font),
        ("a:ea", chinese_font),
        ("a:cs", latin_font),
    ):
        element = properties.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            properties.append(element)
        element.set("typeface", typeface)
        if tag == "a:ea" and typeface == "SimHei":
            element.set("panose", "02010609060101010101")
            element.set("charset", "-122")


def style_run(run, *, size, color, bold, chinese_font, latin_font):
    run.font.name = latin_font
    set_run_typefaces(run, chinese_font, latin_font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color=BODY,
    bold=False,
    chinese_font="SimSun",
    latin_font="Times New Roman",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(0)
    frame.margin_bottom = Inches(0)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    for chunk in filter(None, CJK_RE.split(text)):
        run = paragraph.add_run()
        run.text = chunk
        style_run(
            run,
            size=size,
            color=color,
            bold=bold,
            chinese_font=chinese_font,
            latin_font=latin_font,
        )
    return box


def add_label(slide, text, x, y, w, fill, color=WHITE):
    add_rect(slide, x, y, w, 0.30, fill)
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        0.30,
        size=12.5,
        color=color,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )


def set_title(slide, text: str) -> None:
    for shape in slide.shapes:
        if shape.name == "Rectangle 20":
            flatten_shape(shape)
        if getattr(shape, "text", "") == "[本页标题]":
            frame = shape.text_frame
            frame.clear()
            frame.margin_left = 0
            frame.margin_right = 0
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = text
            run.font.name = "Times New Roman"
            set_run_typefaces(run, "SimHei", "Times New Roman")
            run.font.size = Pt(29)
            run.font.bold = True
            run.font.color.rgb = rgb(INK)
            return
    raise RuntimeError("没有找到正文页标题占位符")


def keep_only_slide(prs: Presentation, index: int) -> None:
    slide_ids = prs.slides._sldIdLst  # python-pptx 暂无公开的删除幻灯片 API。
    for current in reversed(range(len(slide_ids))):
        if current == index:
            continue
        slide_id = slide_ids[current]
        prs.part.drop_rel(slide_id.rId)
        del slide_ids[current]


def draw_engineering_structure(slide) -> None:
    """在已经具有第三章正文模板的幻灯片上绘制工程结构。"""
    set_title(slide, "工程结构")

    # 连接线先绘制，保证节点色块始终处于最上层。
    add_line(slide, 4.575, 2.84, 4.575, 3.18, BLUE, 1.5)
    add_arrow_tip(slide, 4.575, 3.14, "down", BLUE)
    add_line(slide, 8.11, 2.84, 8.11, 3.18, BLUE, 1.5)
    add_arrow_tip(slide, 8.11, 3.14, "down", BLUE)
    add_line(slide, 8.11, 4.10, 8.11, 4.45, TEAL, 1.5)
    add_arrow_tip(slide, 8.11, 4.41, "down", TEAL)
    add_line(slide, 4.55, 5.23, 4.55, 5.88, BLUE, 1.35)
    add_arrow_tip(slide, 4.55, 5.84, "down", BLUE)
    # kernel 直接整合 libs；线路沿主图左边界绕行，避免穿过中部结构。
    add_line(slide, 3.00, 2.37, 2.82, 2.37, "718494", 1.2)
    add_line(slide, 2.82, 2.37, 2.82, 6.29, "718494", 1.2)
    add_line(slide, 2.82, 6.29, 3.00, 6.29, "718494", 1.2)
    add_arrow_tip(slide, 2.96, 6.29, "right", "718494", 0.09)

    # drivers 沿 arch 与 libs 之间的留白登记到 general。
    add_line(slide, 10.15, 5.76, 6.30, 5.76, PURPLE, 1.25)
    add_line(slide, 6.30, 5.76, 6.30, 4.28, PURPLE, 1.25)
    add_line(slide, 6.30, 4.28, 6.15, 4.28, PURPLE, 1.25)
    add_arrow_tip(slide, 6.19, 4.28, "left", PURPLE, 0.09)
    # 主分层区域。
    add_rect(slide, 3.00, 1.92, 6.75, 0.92, NAVY)
    add_rect(slide, 3.00, 1.92, 0.10, 0.92, BLUE)
    add_text(
        slide,
        "kernel",
        3.28,
        2.04,
        1.35,
        0.34,
        size=23,
        color=WHITE,
        bold=True,
        latin_font="Times New Roman",
    )
    add_text(
        slide,
        "最终镜像的整合与运行时编排",
        4.70,
        2.01,
        2.55,
        0.30,
        size=14,
        color=WHITE,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "启动 · 系统调用 · 进程入口 · ELM 运行时",
        4.70,
        2.34,
        4.55,
        0.25,
        size=13.5,
        color="D8E7F3",
        bold=True,
    )
    add_label(slide, "编排层", 8.78, 2.12, 0.70, BLUE)

    add_rect(slide, 3.00, 3.18, 3.15, 2.05, PALE_BLUE)
    add_rect(slide, 3.00, 3.18, 0.09, 2.05, BLUE)
    add_text(
        slide,
        "general",
        3.28,
        3.34,
        1.55,
        0.34,
        size=22,
        color=NAVY,
        bold=True,
    )
    add_label(slide, "稳定抽象", 5.02, 3.34, 0.88, BLUE)
    add_text(
        slide,
        "架构无关的抽象与管理机制",
        3.28,
        3.79,
        2.55,
        0.30,
        size=14,
        color=INK,
        bold=True,
        chinese_font="SimHei",
    )
    add_line(slide, 3.28, 4.18, 5.86, 4.18, "CAD9E5", 0.8)
    add_text(
        slide,
        "设备身份 · PnP · 资源归属",
        3.28,
        4.29,
        2.55,
        0.28,
        size=13.2,
        color=BODY,
        bold=True,
    )
    add_text(
        slide,
        "开放能力 · 生命周期 · 用户态投影",
        3.28,
        4.64,
        2.62,
        0.28,
        size=13.2,
        color=BODY,
        bold=True,
    )

    add_rect(slide, 6.48, 3.18, 3.27, 0.92, PALE_TEAL)
    add_rect(slide, 6.48, 3.18, 0.09, 0.92, TEAL)
    add_text(
        slide,
        "hal",
        6.76,
        3.30,
        0.72,
        0.30,
        size=21,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        "常用架构能力统一入口",
        7.58,
        3.30,
        1.86,
        0.28,
        size=13.2,
        color=INK,
        bold=True,
        chinese_font="SimHei",
    )
    add_text(
        slide,
        "中断 · 时间 · 内存 · 用户上下文",
        6.76,
        3.65,
        2.68,
        0.23,
        size=13.0,
        color=BODY,
        bold=True,
    )

    add_rect(slide, 6.48, 4.45, 3.27, 1.20, "E5EBF2")
    add_rect(slide, 6.48, 4.45, 0.09, 1.20, "526E91")
    add_text(
        slide,
        "arch",
        6.76,
        4.57,
        0.92,
        0.30,
        size=21,
        color="365675",
        bold=True,
    )
    add_label(slide, "独立实现", 8.70, 4.54, 0.82, "526E91")
    add_text(
        slide,
        "RISC-V64     LoongArch64",
        6.76,
        4.93,
        2.60,
        0.24,
        size=13.2,
        color=INK,
        bold=True,
    )
    add_text(
        slide,
        "boot · trap · paging · SMP · VDSO",
        6.76,
        5.18,
        2.66,
        0.21,
        size=12.8,
        color=BODY,
        bold=True,
    )
    add_rect(slide, 3.00, 5.88, 6.75, 0.82, PALE_GRAY)
    add_rect(slide, 3.00, 5.88, 0.10, 0.82, "7D8F9D")
    add_text(
        slide,
        "libs",
        3.28,
        6.00,
        0.84,
        0.28,
        size=20,
        color="485D6C",
        bold=True,
    )
    add_text(
        slide,
        "sched · mm · vfs · fs · net · elm",
        4.18,
        5.98,
        3.13,
        0.27,
        size=13.5,
        color=INK,
        bold=True,
    )
    add_text(
        slide,
        "可复用子系统底座",
        7.72,
        5.99,
        1.62,
        0.26,
        size=13.0,
        color=BODY,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.RIGHT,
    )

    # drivers 是侧向插接的可配置实现区，不画成核心的第五层。
    add_rect(slide, 10.15, 1.92, 2.40, 4.78, PALE_PURPLE)
    add_rect(slide, 10.15, 1.92, 2.40, 0.92, PURPLE)
    add_text(
        slide,
        "drivers",
        10.42,
        2.03,
        1.25,
        0.34,
        size=22,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "集成构建 · ELM 装载",
        10.43,
        2.39,
        1.80,
        0.22,
        size=13.2,
        color="E8E1F3",
        bold=True,
        chinese_font="SimHei",
    )
    add_label(slide, "可配置实现区", 10.40, 2.96, 1.90, PURPLE)

    driver_rows = [
        ("01", "识别硬件", "match · probe · remove"),
        ("02", "实现设备能力", "MMIO · IRQ · DMA · I/O"),
        ("03", "注册抽象能力", "DeviceFunction · resource"),
    ]
    row_y = 3.40
    for number, heading, detail in driver_rows:
        add_text(
            slide,
            number,
            10.40,
            row_y,
            0.32,
            0.24,
            size=12.5,
            color=PURPLE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            heading,
            10.84,
            row_y - 0.02,
            1.32,
            0.26,
            size=14.0,
            color=INK,
            bold=True,
            chinese_font="SimHei",
        )
        add_text(
            slide,
            detail,
            10.84,
            row_y + 0.29,
            1.40,
            0.22,
            size=12.8,
            color=BODY,
            bold=True,
        )
        if number != "03":
            add_line(slide, 10.40, row_y + 0.66, 12.28, row_y + 0.66, "DDD7E8", 0.7)
        row_y += 0.88

    add_rect(slide, 10.40, 6.12, 1.90, 0.38, WHITE)
    add_text(
        slide,
        "硬件细节保持在实现边界内",
        10.43,
        6.12,
        1.84,
        0.38,
        size=13.0,
        color=PURPLE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )

    # kernel 对 drivers 的构建或装载关系，具体方式由配置决定。
    add_line(slide, 9.75, 2.37, 10.15, 2.37, PURPLE, 1.3)
    add_arrow_tip(slide, 10.11, 2.37, "right", PURPLE, 0.09)

    # 标签最后绘制，用白底截断线路，避免被 arch 或 libs 节点遮挡。
    add_rect(slide, 8.42, 5.64, 1.08, 0.24, WHITE)
    add_text(
        slide,
        "登记能力",
        8.42,
        5.64,
        1.08,
        0.24,
        size=12.5,
        color=PURPLE,
        bold=True,
        chinese_font="SimHei",
        align=PP_ALIGN.CENTER,
    )


def build_slide(template: Path, output: Path) -> None:
    prs = Presentation(template)
    slide_index = 11
    draw_engineering_structure(prs.slides[slide_index])
    keep_only_slide(prs, slide_index)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> int:
    root = Path(__file__).resolve().parents[2]
    template = root / "output/presentations/mygo-defense.pptx"
    output = root / "output/presentations/mygo-defense-chapter3-page1.pptx"
    if len(sys.argv) > 1:
        output = Path(sys.argv[1]).resolve()
    build_slide(template, output)
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
