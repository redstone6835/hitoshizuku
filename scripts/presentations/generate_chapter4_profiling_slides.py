#!/usr/bin/env python3
"""生成答辩第四章“调试方法与指令成本建模”并接入完整答辩稿。"""

from __future__ import annotations

import argparse
from copy import deepcopy
from io import BytesIO
import os
from pathlib import Path
import tempfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from generate_engineering_structure_slide import (
    BG,
    BODY,
    BLUE,
    INK,
    LINE,
    MUTED,
    NAVY,
    PALE_BLUE,
    PALE_GRAY,
    PALE_PURPLE,
    PALE_TEAL,
    PURPLE,
    TEAL,
    WHITE,
    add_arrow_tip,
    add_line,
    add_rect,
    add_text,
    rgb,
    set_run_typefaces,
    style_run,
)


EMU_PER_INCH = 914400
MIN_FONT_PT = 14.0
CONTENT_X = 3.00
CONTENT_RIGHT = 12.60
CONTENT_W = CONTENT_RIGHT - CONTENT_X

CHAPTER4_TRANSITION = "第四章 · 调试方法"
CHAPTER5_TRANSITION = "第五章 · 结果结论"
CHAPTER4_TITLES = (
    "内核性能证据体系",
    "受控指令截窗模型",
    "单次系统调用的指令对照",
    "动态指令数量与执行成本",
    "探针与基线的成对差分",
    "指令语义与执行上下文",
    "混杂因素与稳健回归",
    "Super-run 与分层 Bootstrap",
    "统计区间与质量分级",
    "机器学习辅助的结构发现",
    "独立留出与时间稳定性",
    "统计建模与机器学习协同闭环",
)


def inches(value: int) -> float:
    return value / EMU_PER_INCH


def clone_slide(prs: Presentation, source):
    destination = prs.slides.add_slide(source.slide_layout)
    for shape in list(destination.shapes):
        element = shape._element
        element.getparent().remove(element)
    destination.background.fill.solid()
    destination.background.fill.fore_color.rgb = rgb(BG)
    for shape in source.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture = destination.shapes.add_picture(
                BytesIO(shape.image.blob),
                shape.left,
                shape.top,
                shape.width,
                shape.height,
            )
            picture.crop_left = shape.crop_left
            picture.crop_right = shape.crop_right
            picture.crop_top = shape.crop_top
            picture.crop_bottom = shape.crop_bottom
            picture.rotation = shape.rotation
            continue
        destination.shapes._spTree.insert_element_before(
            deepcopy(shape._element), "p:extLst"
        )
    return destination


def slide_texts(slide) -> set[str]:
    return {
        shape.text.strip()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    }


def find_slide(prs: Presentation, exact_text: str):
    for slide in prs.slides:
        if exact_text in slide_texts(slide):
            return slide
    raise RuntimeError(f"没有找到幻灯片：{exact_text}")


def remove_slide(prs: Presentation, slide) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        if prs.part.related_part(slide_id.rId) is slide.part:
            prs.part.drop_rel(slide_id.rId)
            prs.slides._sldIdLst.remove(slide_id)
            return
    raise RuntimeError("没有找到待删除幻灯片关系")


def is_template_chrome(shape) -> bool:
    x, y = inches(shape.left), inches(shape.top)
    w, h = inches(shape.width), inches(shape.height)
    if x + w <= 2.20:
        return True
    if 2.55 <= x <= 2.75 and 0.52 <= y <= 0.75 and h <= 0.90:
        return True
    if 2.80 <= x <= 3.15 and 0.45 <= y <= 0.75 and w >= 8.0:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.LINE and 1.50 <= y <= 1.66:
        return True
    return False


def clean_body_template(slide) -> None:
    for shape in list(slide.shapes):
        if not is_template_chrome(shape):
            element = shape._element
            element.getparent().remove(element)


def set_title(slide, title: str) -> None:
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        x, y = inches(shape.left), inches(shape.top)
        if 2.80 <= x <= 3.15 and 0.45 <= y <= 0.75:
            candidates.append(shape)
    if not candidates:
        raise RuntimeError("没有找到正文标题文本框")
    box = candidates[0]
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = title
    style_run(
        run,
        size=30,
        color=INK,
        bold=True,
        chinese_font="SimHei",
        latin_font="Times New Roman",
    )


def body(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 15.0,
    color=BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    if size < MIN_FONT_PT:
        raise ValueError(f"正文文字不得小于 {MIN_FONT_PT:g} pt：{text}")
    box = add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        chinese_font="SimSun",
        latin_font="Times New Roman",
        align=align,
        valign=valign,
        margin=0.05,
    )
    box.text_frame.word_wrap = True
    return box


def heading(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float = 0.34,
    *,
    size: float = 19.0,
    color=INK,
    align=PP_ALIGN.LEFT,
):
    box = add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=max(size, MIN_FONT_PT),
        color=color,
        bold=True,
        chinese_font="SimHei",
        latin_font="Times New Roman",
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.03,
    )
    box.text_frame.word_wrap = True
    return box


def lead(slide, text: str, *, size: float = 16.0) -> None:
    body(
        slide,
        text,
        CONTENT_X,
        1.78,
        CONTENT_W,
        0.62,
        size=size,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
    title_size: float = 18.0,
    detail_size: float = 15.0,
    center: bool = False,
) -> None:
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.08, h, accent)
    heading(
        slide,
        title,
        x + 0.22,
        y + 0.13,
        w - 0.40,
        0.34,
        size=title_size,
        color=NAVY,
        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
    )
    body(
        slide,
        detail,
        x + 0.22,
        y + 0.57,
        w - 0.40,
        h - 0.67,
        size=detail_size,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
    )


def compact_row(
    slide,
    y: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
    h: float = 0.62,
    title_w: float = 2.08,
    detail_size: float = 15.0,
) -> None:
    add_rect(slide, CONTENT_X, y, CONTENT_W, h, fill)
    add_rect(slide, CONTENT_X, y, 0.08, h, accent)
    heading(
        slide,
        title,
        CONTENT_X + 0.23,
        y + 0.12,
        title_w,
        h - 0.22,
        size=17,
        color=NAVY,
    )
    body(
        slide,
        detail,
        CONTENT_X + 0.33 + title_w,
        y + 0.08,
        CONTENT_W - title_w - 0.55,
        h - 0.14,
        size=detail_size,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def boundary(
    slide,
    text: str,
    *,
    label: str,
    fill=PALE_GRAY,
    accent=MUTED,
    label_w: float = 1.68,
) -> None:
    add_rect(slide, CONTENT_X, 6.14, CONTENT_W, 0.64, fill)
    add_rect(slide, CONTENT_X, 6.14, 0.08, 0.64, accent)
    heading(
        slide,
        label,
        CONTENT_X + 0.23,
        6.25,
        label_w,
        0.31,
        size=17,
        color=NAVY,
    )
    body(
        slide,
        text,
        CONTENT_X + 0.34 + label_w,
        6.18,
        CONTENT_W - label_w - 0.55,
        0.56,
        size=14,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def arrow(slide, x1: float, y: float, x2: float, *, color=BLUE) -> None:
    add_line(slide, x1, y, x2, y, color, 1.5)
    add_arrow_tip(slide, x2, y, "right", color, 0.11)


def stage_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    detail: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
    title_size: float = 17,
    detail_size: float = 14,
) -> None:
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.07, h, accent)
    heading(
        slide,
        title,
        x + 0.16,
        y + 0.12,
        w - 0.28,
        0.34,
        size=title_size,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    if detail:
        body(
            slide,
            detail,
            x + 0.14,
            y + 0.52,
            w - 0.24,
            h - 0.60,
            size=detail_size,
            color=BODY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def metric(
    slide,
    x: float,
    y: float,
    w: float,
    value: str,
    label: str,
    *,
    fill=PALE_BLUE,
    accent=BLUE,
) -> None:
    add_rect(slide, x, y, w, 0.92, fill)
    add_rect(slide, x, y, 0.07, 0.92, accent)
    heading(
        slide,
        value,
        x + 0.16,
        y + 0.10,
        w - 0.26,
        0.38,
        size=20,
        color=accent,
        align=PP_ALIGN.CENTER,
    )
    body(
        slide,
        label,
        x + 0.12,
        y + 0.52,
        w - 0.20,
        0.28,
        size=14,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def enforce_font_floor(slide) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                if run.font.size is None or run.font.size.pt < MIN_FONT_PT:
                    run.font.size = Pt(MIN_FONT_PT)
                if run.font.name is None:
                    set_run_typefaces(run, "SimSun", "Times New Roman")


def draw_01(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[0])
    lead(slide, "性能结论沿五层证据收敛：固定负载、还原路径、统计工作量、估计成本，并把结果归因到责任阶段。")

    labels = (
        ("固定负载", "保持比较条件一致", PALE_BLUE, BLUE),
        ("动态路径", "还原真实执行顺序", PALE_TEAL, TEAL),
        ("执行计数", "统计实际工作量", PALE_PURPLE, PURPLE),
        ("成本估计", "给出中心值与区间", PALE_BLUE, BLUE),
        ("热点归因", "对应函数与责任阶段", PALE_GRAY, MUTED),
    )
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(labels):
        stage_box(slide, x, 2.55, 1.64, 1.50, title, detail, fill=fill, accent=accent)
        if index != len(labels) - 1:
            arrow(slide, x + 1.67, 3.30, x + 1.88, color=accent)
        x += 1.96

    compact_row(slide, 4.36, "宏观比较", "整体时间与进度回答“是否变快”，指令工作量和热点回答“变化来自哪里”。", fill=PALE_BLUE, accent=BLUE, h=0.68, title_w=1.76)
    compact_row(slide, 5.18, "微观比较", "固定同一类操作，对照两条动态路径并定位结构差异。", fill=PALE_TEAL, accent=TEAL, h=0.68, title_w=1.76)
    boundary(slide, "统计模型量化成本与不确定性，机器学习检查剩余结构并反馈实验设计。", label="建模分工")


def draw_02(slide) -> None:
    set_title(slide, METHOD_TITLES[1])
    lead(slide, "制作 C/ASM 探针并在 MyGO 与 Linux 中运行；QEMU-TCG 识别探针头、限定窗口并记录两条动态内核路径。")

    steps = (
        ("编写探针", "C / ASM 固定操作", PALE_BLUE, BLUE),
        ("双系统运行", "MyGO / Linux", PALE_PURPLE, PURPLE),
        ("识别探针头", "限定采集窗口", PALE_TEAL, TEAL),
        ("记录指令", "指令 / PC", PALE_PURPLE, PURPLE),
        ("过滤并查表", "PC → 函数", PALE_GRAY, MUTED),
    )
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(steps):
        stage_box(slide, x, 2.40, 1.66, 1.38, title, detail, fill=fill, accent=accent)
        if index != len(steps) - 1:
            arrow(slide, x + 1.68, 3.09, x + 1.90, color=accent)
        x += 1.96

    panel(slide, 3.00, 4.08, 4.55, 1.62, "动态路径分析", "按责任阶段和函数语义对齐两条路径，定位结构差异来自哪个内核责任。", fill=PALE_BLUE, accent=BLUE, title_size=18, detail_size=15)
    panel(slide, 7.86, 4.08, 4.74, 1.62, "具体优化内容", "结合函数内指令数量和汇编差异，定位额外搬运、检查、分支或慢路径。", fill=PALE_TEAL, accent=TEAL, title_size=18, detail_size=15)
    boundary(slide, "探针法观察一次固定操作，是同一指令归因思路的微观形态。", label="微观输出")


def draw_03(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[2])
    lead(slide, "单次系统调用按责任阶段对齐：比较入口、分发、业务处理和返回，而不是机械比较同名函数。")

    stages = (
        ("ecall", "用户边界"),
        ("trap 入口", "进入内核"),
        ("syscall 分发", "选择责任阶段"),
        ("业务处理", "完成目标操作"),
        ("返回检查", "准备返回用户态"),
        ("返回用户态", "恢复执行"),
    )
    x = 3.00
    for index, (title, detail) in enumerate(stages):
        fill, accent = (PALE_TEAL, TEAL) if index in (0, 5) else (PALE_BLUE, BLUE)
        if index == 3:
            fill, accent = PALE_PURPLE, PURPLE
        stage_box(slide, x, 2.35, 1.35, 1.18, title, detail, fill=fill, accent=accent, title_size=16, detail_size=14)
        if index != len(stages) - 1:
            arrow(slide, x + 1.37, 2.94, x + 1.55, color=accent)
        x += 1.60

    compact_row(slide, 3.84, "阶段指令数", "分别统计入口、分发、业务和返回，保留责任结构。", fill=PALE_BLUE, accent=BLUE, h=0.58, title_w=1.82)
    compact_row(slide, 4.54, "首个动态分歧", "从两条路径的第一个不同位置开始定位结构差异。", fill=PALE_PURPLE, accent=PURPLE, h=0.68, title_w=1.82)
    compact_row(slide, 5.34, "责任对齐", "把差异归入对应阶段，避免只按函数名称下结论。", fill=PALE_TEAL, accent=TEAL, h=0.58, title_w=1.82)
    boundary(slide, "单次轨迹解释路径结构，长负载计数补足其出现频率。", label="两类尺度")


def draw_04(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[3])
    lead(slide, "动态指令数描述工作量；将其与语义成本结合，才能得到可解释的相对性能估计。")

    add_rect(slide, 3.00, 2.48, 9.60, 1.06, NAVY)
    body(slide, "T_kernel  ≈  Σₖ  Nₖ · θₖ", 3.20, 2.64, 9.20, 0.68, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    panel(slide, 3.00, 3.82, 2.92, 1.74, "Nₖ · 动态次数", "在固定窗口中统计每类指令的真实执行次数。", fill=PALE_BLUE, accent=BLUE, title_size=19, detail_size=15)
    panel(slide, 6.24, 3.82, 2.92, 1.74, "θₖ · 上下文成本", "按执行语义估计中心成本与不确定性区间。", fill=PALE_PURPLE, accent=PURPLE, title_size=19, detail_size=15)
    panel(slide, 9.48, 3.82, 3.12, 1.74, "函数与阶段归因", "将成本汇总到函数和系统责任阶段。", fill=PALE_TEAL, accent=TEAL, title_size=19, detail_size=15)
    boundary(slide, "模型输出相对成本构成，并把中心值和区间同时用于归因。", label="归因输出")


def draw_05(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[4])
    lead(slide, "成对差分让目标路径与同形态基线共享公共成本，差值用于估计目标操作的增量。")

    panel(slide, 3.00, 2.36, 4.16, 1.78, "目标窗口", "公共准备 → 目标操作 → 公共收束", fill=PALE_BLUE, accent=BLUE, title_size=20, detail_size=16, center=True)
    panel(slide, 8.44, 2.36, 4.16, 1.78, "基线窗口", "相同准备 → 对照操作 → 相同收束", fill=PALE_TEAL, accent=TEAL, title_size=20, detail_size=16, center=True)
    body(slide, "−", 7.36, 2.86, 0.88, 0.48, size=28, color=PURPLE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_rect(slide, 3.00, 4.39, 9.60, 0.76, PALE_PURPLE)
    add_rect(slide, 3.00, 4.39, 0.08, 0.76, PURPLE)
    body(slide, "dᵢ = (T_probe,i − T_baseline,i) / (N_probe,i − N_baseline,i)", 3.24, 4.49, 9.12, 0.52, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    metric(slide, 3.00, 5.34, 2.90, "成对窗口", "共享公共成本", fill=PALE_BLUE, accent=BLUE)
    metric(slide, 6.23, 5.34, 2.92, "交错顺序", "降低顺序影响", fill=PALE_TEAL, accent=TEAL)
    metric(slide, 9.48, 5.34, 3.12, "固定规模", "保证可比性", fill=PALE_GRAY, accent=MUTED)


def draw_06(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[5])
    lead(slide, "成本按指令身份和执行语义分层；不同数据流、控制流和访存情境分别估计，避免平均掉关键差异。")

    rows = (
        ("整数算术", "依赖与独立数据流", PALE_BLUE, BLUE),
        ("条件分支", "方向与历史", PALE_TEAL, TEAL),
        ("访存与栈", "缓存与地址形态", PALE_PURPLE, PURPLE),
        ("跳转与调用", "控制流类别", PALE_BLUE, BLUE),
        ("原子操作", "成功与失败语义", PALE_TEAL, TEAL),
        ("浮点与系统指令", "操作语义与上下文", PALE_GRAY, MUTED),
    )
    y = 2.42
    for title, detail, fill, accent in rows:
        compact_row(slide, y, title, detail, fill=fill, accent=accent, h=0.52, title_w=1.76, detail_size=15)
        y += 0.61

    boundary(slide, "数据流和控制流差异在成本估计中保持可见，模型不会把不同语义静默合并。", label="数据流结论")


def draw_07(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[6])
    lead(slide, "成对差分之后，仍需用稳健模型吸收运行差异、顺序影响、时间漂移和规模效应。")

    add_rect(slide, 3.00, 2.39, 9.60, 0.76, NAVY)
    body(slide, "dᵢ = θ + α_run(i) + βₒOᵢ + β_dDᵢ + β_bBᵢ + εᵢ", 3.16, 2.50, 9.28, 0.52, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    panels = (
        (3.00, "目标量 θ", "目标指令相对 baseline 的中心成本。", PALE_BLUE, BLUE),
        (5.45, "运行效应 α", "吸收不同运行的整体差异。", PALE_TEAL, TEAL),
        (7.90, "顺序与漂移", "表示实验顺序和时间位置。", PALE_PURPLE, PURPLE),
        (10.35, "规模效应", "检查成本随工作量的变化。", PALE_GRAY, MUTED),
    )
    for x, title, detail, fill, accent in panels:
        panel(slide, x, 3.43, 2.25, 1.56, title, detail, fill=fill, accent=accent, title_size=17, detail_size=14)

    compact_row(slide, 5.24, "稳健拟合", "主体样本保持高权重，异常和长尾样本降低影响，输出中心成本及区间。", fill=PALE_BLUE, accent=BLUE, h=0.58, title_w=1.74)
    boundary(slide, "回归结果同时保留运行效应和目标成本，便于比较并归因。", label="模型输出")


def draw_08(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[7])
    lead(slide, "Bootstrap 以完整运行组合为独立样本，并在运行内部保持时间块结构，得到更稳健的不确定性区间。")

    body(slide, "一个独立运行组合", 3.00, 2.36, 2.16, 0.34, size=16, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    sequence = (("A", "测量", BLUE), ("B", "对照", TEAL), ("B", "对照", TEAL), ("A", "测量", BLUE))
    x = 5.05
    for index, (letter, detail, accent) in enumerate(sequence):
        add_rect(slide, x, 2.28, 1.46, 0.78, PALE_BLUE if letter == "A" else PALE_TEAL)
        heading(slide, letter, x + 0.10, 2.36, 0.38, 0.32, size=20, color=accent, align=PP_ALIGN.CENTER)
        body(slide, detail, x + 0.46, 2.34, 0.88, 0.34, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if index != len(sequence) - 1:
            arrow(slide, x + 1.47, 2.67, x + 1.67, color=accent)
        x += 1.82

    compact_row(slide, 3.40, "顶层重采样", "以完整运行组合为单位抽样，保留独立性。", fill=PALE_PURPLE, accent=PURPLE, h=0.68, title_w=1.86)
    compact_row(slide, 4.22, "时间块重采样", "保留相邻窗口的相关结构和慢漂移。", fill=PALE_TEAL, accent=TEAL, h=0.68, title_w=1.86)
    compact_row(slide, 5.04, "全量重新拟合", "每次重采样都重新估计中心值和区间。", fill=PALE_BLUE, accent=BLUE, h=0.68, title_w=1.86)
    boundary(slide, "分层 Bootstrap 输出稳定的成本区间，并把运行结构纳入不确定性。", label="重采样结论")


def draw_09(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[8])
    lead(slide, "质量评价同时关注点估计、跨运行稳定性和采集完整性，最终输出中心值、区间和质量等级。")

    add_rect(slide, 3.00, 2.40, 4.44, 1.18, NAVY)
    body(slide, "M_b = maxₖ |(θ̂ₖ⁽ᵇ⁾ − θ̂ₖ) / sₖ|", 3.16, 2.55, 4.12, 0.40, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    body(slide, "同时区间控制多上下文比较的整体误报", 3.18, 2.99, 4.08, 0.30, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    panel(slide, 7.76, 2.40, 4.84, 1.18, "正成本锚点", "用稳定的正成本参照校准相对指令成本。", fill=PALE_TEAL, accent=TEAL, title_size=18, detail_size=14)

    left = (
        ("样本闭合", "样本完整、区间可复现"),
        ("模型稳定", "估计对运行条件不敏感"),
    )
    right = (
        ("跨运行一致性", "不同运行得到相近结论"),
        ("环境记录", "保留必要的主机与时间上下文"),
    )
    for index, (title, detail) in enumerate(left):
        panel(slide, 3.00, 3.86 + index * 1.02, 4.44, 0.88, title, detail, fill=PALE_BLUE, accent=BLUE, title_size=17, detail_size=14)
    for index, (title, detail) in enumerate(right):
        panel(slide, 7.76, 3.86 + index * 1.02, 4.84, 0.88, title, detail, fill=PALE_PURPLE if index == 0 else PALE_GRAY, accent=PURPLE if index == 0 else MUTED, title_size=17, detail_size=14)

    boundary(slide, "输出保留中心值、同时区间和质量标签；下游归因据此区分稳定结论与探索估计。", label="分级产物")


def draw_10(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[9])
    lead(slide, "机器学习作为结构发现工具，读取差分样本并寻找稳健统计模型尚未表达的非线性关系。")

    stages = (
        ("差分样本", "目标成本与上下文", PALE_BLUE, BLUE),
        ("特征展开", "语义、执行模式和位置", PALE_TEAL, TEAL),
        ("结构模型", "学习非线性关系", PALE_PURPLE, PURPLE),
        ("残差结构", "发现剩余交互", PALE_BLUE, BLUE),
        ("实验反馈", "更新探针与采样", PALE_GRAY, MUTED),
    )
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(stages):
        stage_box(slide, x, 2.38, 1.66, 1.44, title, detail, fill=fill, accent=accent)
        if index != len(stages) - 1:
            arrow(slide, x + 1.68, 3.10, x + 1.90, color=accent)
        x += 1.96

    panel(slide, 3.00, 4.12, 4.55, 1.60, "结构发现", "识别上下文之间的非线性关联，找出需要进一步解释的成本差异。", fill=PALE_BLUE, accent=BLUE, title_size=19, detail_size=15)
    panel(slide, 7.86, 4.12, 4.74, 1.60, "采样反馈", "依据残差、区间和稳定性调整下一轮实验重点。", fill=PALE_TEAL, accent=TEAL, title_size=19, detail_size=15)
    boundary(slide, "机器学习带来小幅但稳定的结构增量，价值在于发现方向并改善实验设计。", label="结构结论")


def draw_11(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[10])
    lead(slide, "留出评估用独立运行组合划分训练、校准和测试，并同时检验随机重复与时间前向稳定性。")

    panel(slide, 3.00, 2.36, 4.55, 1.62, "随机完整组留出", "以独立运行组合划分训练、校准和测试，评价重复性。", fill=PALE_BLUE, accent=BLUE, title_size=19, detail_size=15, center=True)
    panel(slide, 7.86, 2.36, 4.74, 1.62, "时间前向留出", "按采集顺序前向测试，评价跨时间的稳定性。", fill=PALE_TEAL, accent=TEAL, title_size=19, detail_size=15, center=True)

    metric(slide, 3.00, 4.26, 2.90, "0.12258 ns", "HGB OOF MAE", fill=PALE_PURPLE, accent=PURPLE)
    metric(slide, 6.23, 4.26, 2.92, "0.12392 ns", "context+batch 基线", fill=PALE_BLUE, accent=BLUE)
    metric(slide, 9.48, 4.26, 3.12, "0.00134 ns", "增量改善", fill=PALE_TEAL, accent=TEAL)

    compact_row(slide, 5.40, "时间结构发现", "前向测试显示结论具有时间稳定性，同时暴露出可继续优化的慢漂移。", fill=PALE_GRAY, accent=MUTED, h=0.58, title_w=1.92, detail_size=14)
    boundary(slide, "随机与时间前向留出共同支持模型的泛化结论，并反馈后续采集协议。", label="协议反馈")


def draw_12(slide) -> None:
    set_title(slide, CHAPTER4_TITLES[11])
    lead(slide, "统计估计给出可解释的成本和区间，机器学习发现剩余结构，二者共同反馈下一轮受控实验。")

    stages = (
        ("受控微基准", "固定上下文", PALE_BLUE, BLUE),
        ("成对差分", "消除公共成本", PALE_TEAL, TEAL),
        ("稳健统计", "中心与同时区间", PALE_PURPLE, PURPLE),
        ("ML 结构检查", "残差与前向稳定性", PALE_BLUE, BLUE),
        ("实验协议更新", "探针、分层与门禁", PALE_GRAY, MUTED),
    )
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(stages):
        stage_box(slide, x, 2.36, 1.66, 1.42, title, detail, fill=fill, accent=accent)
        if index != len(stages) - 1:
            arrow(slide, x + 1.68, 3.07, x + 1.90, color=accent)
        x += 1.96

    compact_row(slide, 4.10, "指令层产物", "动态次数、中心成本、区间和稳定性标签。", fill=PALE_BLUE, accent=BLUE, h=0.64, title_w=1.82)
    compact_row(slide, 4.88, "函数层产物", "把成本区间汇总到函数和系统责任阶段。", fill=PALE_TEAL, accent=TEAL, h=0.64, title_w=1.82)
    compact_row(slide, 5.66, "现行实验设计", "以独立运行组合进行随机与时间前向留出。", fill=PALE_PURPLE, accent=PURPLE, h=0.40, title_w=1.82, detail_size=14)
    boundary(slide, "动态指令流从调试日志转化为可计数、可定价、可归因的证据，并持续反馈新的实验方向。", label="方法收束")


METHOD_TITLES = (
    "三种调试方法与演进关系",
    "调试探针法：微观调用路径",
    "整体测量法：1200 秒宏观窗口",
    "统计分析法：构建可信指令权重",
    "加权归因与优化优先级",
)


def draw_method_01(slide) -> None:
    set_title(slide, METHOD_TITLES[0])
    lead(slide, "三种方法共享“动态指令 → 内核地址 → 函数归属”的证据链，并依次完成微观定位、宏观筛选与加权定序。")

    panel(slide, 3.00, 2.38, 2.92, 2.02, "调试探针法 · 微观", "C/ASM 固定操作；对齐 MyGO 与 Linux 的单次内核路径，定位动态函数序列与函数内差异。", fill=PALE_BLUE, accent=BLUE, title_size=18, detail_size=15)
    panel(slide, 6.24, 2.38, 2.92, 2.02, "整体测量法 · 宏观", "同一负载运行 1200 秒；按函数汇总内核动态指令，筛选长期累计热点。", fill=PALE_TEAL, accent=TEAL, title_size=18, detail_size=15)
    panel(slide, 9.48, 2.38, 3.12, 2.02, "统计分析法 · 加权", "测量指令在特定上下文中的相对时间，隔离环境噪声并建立可信权重表。", fill=PALE_PURPLE, accent=PURPLE, title_size=18, detail_size=15)

    compact_row(slide, 4.72, "统一思路", "TCG 记录指令及其 PC，过滤内核地址，再由当次构建的符号表映射到函数。", fill=PALE_BLUE, accent=BLUE, h=0.62, title_w=1.76)
    compact_row(slide, 5.48, "演进原因", "指令数只能产生候选热点，不能直接代表时间；最终必须引入指令权重。", fill=PALE_PURPLE, accent=PURPLE, h=0.62, title_w=1.76)
    boundary(slide, "探针与整体测量分别观察微观路径和宏观累计量，统计分析为二者提供统一成本尺度。", label="方法关系")


def draw_method_03(slide) -> None:
    set_title(slide, METHOD_TITLES[2])
    lead(slide, "在固定资源与同一负载定义下，各运行 1200 秒并记录动态指令，从全局视角筛选内核函数热点。")

    steps = (
        ("固定条件", "资源与负载一致", PALE_BLUE, BLUE),
        ("1200 秒", "覆盖长期行为", PALE_PURPLE, PURPLE),
        ("记录指令", "TCG 指令与 PC", PALE_TEAL, TEAL),
        ("过滤内核", "保留内核 PC", PALE_PURPLE, PURPLE),
        ("函数聚合", "符号表 → 指令构成", PALE_GRAY, MUTED),
    )
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(steps):
        stage_box(slide, x, 2.38, 1.66, 1.42, title, detail, fill=fill, accent=accent)
        if index != len(steps) - 1:
            arrow(slide, x + 1.68, 3.09, x + 1.90, color=accent)
        x += 1.96

    panel(slide, 3.00, 4.08, 4.55, 1.62, "宏观输出", "按函数汇总动态指令总量与构成，识别长窗口中累计工作量高的候选热点。", fill=PALE_BLUE, accent=BLUE, title_size=18, detail_size=15)
    panel(slide, 7.86, 4.08, 4.74, 1.62, "与探针法的关系", "符号归因链相同；探针法解释单次路径，整体法衡量长负载中的累计工作量。", fill=PALE_TEAL, accent=TEAL, title_size=18, detail_size=15)
    boundary(slide, "跨系统比较须按完成操作数或有效进度归一化；原始数量只用于筛选候选热点。", label="宏观边界")


def draw_method_04(slide) -> None:
    set_title(slide, METHOD_TITLES[3])
    lead(slide, "统计分析为指令及其执行上下文建立可信权重：先测量初始时间增量，再隔离环境噪声与时间漂移。")

    steps = (
        ("改造探针", "加入固定指令段", PALE_BLUE, BLUE),
        ("前后对比", "时间差 / 指令增量", PALE_TEAL, TEAL),
        ("噪声隔离", "差分、稳健拟合", PALE_PURPLE, PURPLE),
        ("Bootstrap", "复核区间与稳定性", PALE_BLUE, BLUE),
        ("ML 检查", "比较残差与方案", PALE_GRAY, MUTED),
    )
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(steps):
        stage_box(slide, x, 2.30, 1.66, 1.48, title, detail, fill=fill, accent=accent, title_size=16, detail_size=14)
        if index != len(steps) - 1:
            arrow(slide, x + 1.68, 3.04, x + 1.90, color=accent)
        x += 1.96

    add_rect(slide, 3.00, 4.10, 4.55, 0.78, NAVY)
    body(slide, "wₖ⁽⁰⁾  ≈  ΔT / ΔNₖ", 3.20, 4.15, 4.15, 0.34, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    body(slide, "单一指令类 · 固定执行上下文", 3.20, 4.51, 4.15, 0.24, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_rect(slide, 7.86, 4.10, 4.74, 0.78, PALE_TEAL)
    add_rect(slide, 7.86, 4.10, 0.08, 0.78, TEAL)
    heading(slide, "可信权重表", 8.08, 4.28, 1.48, 0.34, size=17, color=NAVY)
    body(slide, "中心权重 · 区间 · 稳定性标签", 9.66, 4.19, 2.70, 0.48, size=14, color=BODY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    compact_row(slide, 5.18, "统计方法", "成对差分吸收公共成本，稳健回归处理漂移，分层 Bootstrap 保留运行结构。", fill=PALE_PURPLE, accent=PURPLE, h=0.62, title_w=1.72)
    boundary(slide, "Bootstrap 评估不确定性；独立留出与机器学习检查残差结构，辅助筛选可信方案。", label="筛选原则")


def draw_method_05(slide) -> None:
    set_title(slide, METHOD_TITLES[4])
    lead(slide, "将微观动态路径、宏观累计指令和可信上下文权重合并，才能得到可解释的性能开销与优化优先级。")

    stages = (
        ("探针路径", "具体调用与指令差异", PALE_BLUE, BLUE),
        ("累计指令", "函数内总量与构成", PALE_TEAL, TEAL),
        ("上下文权重", "相对时间与区间", PALE_PURPLE, PURPLE),
        ("加权成本", "数量 × 权重", PALE_BLUE, BLUE),
        ("优化定位", "优先级与具体内容", PALE_GRAY, MUTED),
    )
    x = 3.00
    for index, (title, detail, fill, accent) in enumerate(stages):
        stage_box(slide, x, 2.34, 1.66, 1.42, title, detail, fill=fill, accent=accent, title_size=16, detail_size=14)
        if index != len(stages) - 1:
            arrow(slide, x + 1.68, 3.05, x + 1.90, color=accent)
        x += 1.96

    add_rect(slide, 3.00, 4.08, 9.60, 0.78, NAVY)
    body(slide, "C(f)  ≈  ∑  N(f, k, c) · w(k, c)", 3.18, 4.19, 9.24, 0.52, size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    compact_row(slide, 5.14, "工程结果", "优先优化“累计加权成本高”的函数，再回到探针动态路径确定具体修改内容。", fill=PALE_TEAL, accent=TEAL, h=0.64, title_w=1.72)
    boundary(slide, "权重和成本只用于同一 QEMU 环境中的相对比较；详细统计参数与验证结果保存在补充文档。", label="结论边界")


DRAWERS = (
    draw_method_01,
    draw_02,
    draw_method_03,
    draw_method_04,
    draw_method_05,
)

ACTIVE_CHAPTER4_TITLES = METHOD_TITLES


def update_transition(slide) -> None:
    replacement = "从微观探针到长窗口测量，再以可信指令权重确定优化优先级。"
    candidates = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and 3.10 <= inches(shape.top) <= 3.45
        and shape.text.strip()
    ]
    if not candidates:
        return
    box = candidates[0]
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = replacement
    style_run(
        run,
        size=18,
        color=BODY,
        bold=False,
        chinese_font="SimSun",
        latin_font="Times New Roman",
    )


def insert_chapter(prs: Presentation) -> list:
    transition = find_slide(prs, CHAPTER4_TRANSITION)
    next_transition = find_slide(prs, CHAPTER5_TRANSITION)
    transition_index = list(prs.slides).index(transition)
    next_index = list(prs.slides).index(next_transition)
    if next_index <= transition_index + 1:
        raise RuntimeError("第四章缺少可复用正文模板")

    existing = list(prs.slides)[transition_index + 1 : next_index]
    template = existing[0]
    update_transition(transition)

    added = []
    for drawer in DRAWERS:
        slide = clone_slide(prs, template)
        clean_body_template(slide)
        drawer(slide)
        enforce_font_floor(slide)
        added.append(slide)

    for slide in existing:
        remove_slide(prs, slide)

    transition_index = list(prs.slides).index(transition)
    slide_ids = prs.slides._sldIdLst
    added_ids = list(slide_ids)[-len(added) :]
    for slide_id in added_ids:
        slide_ids.remove(slide_id)
    for offset, slide_id in enumerate(added_ids, 1):
        slide_ids.insert(transition_index + offset, slide_id)
    return added


def atomic_save(prs: Presentation, output: Path) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        prefix=f".{output.stem}-",
        suffix=".pptx",
        dir=output.parent,
        delete=False,
    ) as temporary:
        temporary_path = Path(temporary.name)
    try:
        prs.save(temporary_path)
        os.replace(temporary_path, output)
        output.chmod(0o644)
    except Exception:
        temporary_path.unlink(missing_ok=True)
        raise


def build_topic(full_output: Path, topic_output: Path) -> None:
    prs = Presentation(full_output)
    keep_titles = set(ACTIVE_CHAPTER4_TITLES)
    for slide in list(prs.slides):
        if not (slide_texts(slide) & keep_titles):
            remove_slide(prs, slide)
    if len(prs.slides) != len(ACTIVE_CHAPTER4_TITLES):
        raise RuntimeError(f"第四章专题页数错误：{len(prs.slides)}")
    atomic_save(prs, topic_output)


def validate(path: Path, *, expected_slides: int | None = None) -> None:
    prs = Presentation(path)
    if expected_slides is not None and len(prs.slides) != expected_slides:
        raise RuntimeError(f"{path} 页数错误：{len(prs.slides)}")
    titles_found = []
    for slide_number, slide in enumerate(prs.slides, 1):
        texts = slide_texts(slide)
        titles_found.extend(title for title in ACTIVE_CHAPTER4_TITLES if title in texts)
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip() or run.font.size is None:
                        continue
                    if any(title in texts for title in ACTIVE_CHAPTER4_TITLES) and run.font.size.pt < MIN_FONT_PT:
                        raise RuntimeError(
                            f"第 {slide_number} 页文字小于 {MIN_FONT_PT:g} pt：{run.text!r}"
                        )
    if sorted(titles_found) != sorted(ACTIVE_CHAPTER4_TITLES):
        raise RuntimeError("第四章标题集合不完整或重复")


def main() -> int:
    parser = argparse.ArgumentParser()
    root = Path(__file__).resolve().parents[2]
    parser.add_argument(
        "--base",
        type=Path,
        default=root / "output/presentations/mygo-defense-full.pptx",
    )
    parser.add_argument(
        "--full-output",
        type=Path,
        default=root / "output/presentations/mygo-defense-full.pptx",
    )
    parser.add_argument(
        "--topic-output",
        type=Path,
        default=root / "output/presentations/mygo-defense-chapter4-profiling-5pages.pptx",
    )
    args = parser.parse_args()

    base = args.base.resolve()
    full_output = args.full_output.resolve()
    topic_output = args.topic_output.resolve()
    prs = Presentation(base)
    insert_chapter(prs)
    atomic_save(prs, full_output)
    validate(full_output)
    build_topic(full_output, topic_output)
    validate(topic_output, expected_slides=len(ACTIVE_CHAPTER4_TITLES))
    print(full_output)
    print(topic_output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
